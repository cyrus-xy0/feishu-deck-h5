"""Focused contract tests for PPTX Template Design System extraction."""

from __future__ import annotations

import json
import subprocess
import sys
from base64 import b64decode
from pathlib import Path

import pytest


ASSETS = Path(__file__).resolve().parent.parent / "assets"
if str(ASSETS) not in sys.path:
    sys.path.insert(0, str(ASSETS))

et = pytest.importorskip("extract_template", reason="python-pptx/lxml not importable")
pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


def _template(path: Path, *, width: float = 16, height: float = 3, slides: int = 2) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    for index in range(slides):
        layout = prs.slide_layouts[0] if index == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = "Template Cover" if index == 0 else f"Content {index}"
        for placeholder in slide.placeholders:
            if placeholder is not slide.shapes.title and getattr(placeholder, "has_text_frame", False):
                placeholder.text = "Replaceable content"
    prs.save(path)
    return path


def test_extract_preserves_16_by_3_canvas_and_emits_review(tmp_path: Path):
    source = _template(tmp_path / "wide-template.pptx")
    out = tmp_path / "out"
    dossier, pack, preview = et.extract_template(
        source,
        out,
        role_mappings={"cover": "slide:1", "raw": "slide:2"},
        aliases={"agenda": "raw", "end": "cover"},
    )

    assert dossier["canvas"]["source_width_emu"] == int(Inches(16))
    assert dossier["canvas"]["source_height_emu"] == int(Inches(3))
    assert dossier["canvas"]["aspect_ratio"]["label"] == "16:3"
    assert dossier["canvas"]["recommended_design_canvas"] == {
        "width": 1920,
        "height": 360,
        "policy": "preserve-source-aspect-ratio",
    }
    assert pack["canvas"] == dossier["canvas"]
    assert preview.is_file()
    preview_text = preview.read_text(encoding="utf-8")
    assert "16:3" in preview_text
    assert "六类 Layout 覆盖" in preview_text
    assert "<img" not in preview_text.lower(), "review must not smuggle a whole-page screenshot"


def test_partial_coverage_and_aliases_are_explicit(tmp_path: Path):
    source = _template(tmp_path / "partial.pptx")
    dossier = et.build_dossier(source, tmp_path / "extract")
    pack = et.build_template_pack(
        dossier,
        template_id="partial",
        role_mappings={"cover": "slide:1", "raw": "slide:2"},
        aliases={"agenda": "raw", "end": "cover"},
    )

    statuses = {role: item["status"] for role, item in pack["layout_coverage"].items()}
    assert statuses == {
        "cover": "native",
        "raw": "native",
        "section": "unsupported",
        "quote": "unsupported",
        "agenda": "alias",
        "end": "alias",
    }
    assert pack["layout_coverage"]["agenda"]["alias_to"] == "raw"
    assert pack["layout_coverage"]["end"]["alias_to"] == "cover"
    assert set(pack["layouts"]) == {"role-cover", "role-raw"}
    assert pack["extraction_report"]["missing_roles"] == ["section", "quote"]


def test_alias_cycle_and_alias_to_missing_role_are_refused(tmp_path: Path):
    source = _template(tmp_path / "bad-alias.pptx", slides=1)
    dossier = et.build_dossier(source, tmp_path / "extract")

    with pytest.raises(et.ExtractionError, match="cycle"):
        et.build_template_pack(
            dossier,
            template_id="cycle",
            aliases={"cover": "raw", "raw": "cover"},
        )

    with pytest.raises(et.ExtractionError, match="unsupported"):
        et.build_template_pack(
            dossier,
            template_id="missing-target",
            role_mappings={"cover": "slide:1"},
            aliases={"agenda": "raw"},
        )


def test_missing_roles_remain_unsupported_without_invention(tmp_path: Path):
    source = _template(tmp_path / "missing.pptx", slides=1)
    dossier = et.build_dossier(source, tmp_path / "extract")
    pack = et.build_template_pack(
        dossier,
        template_id="missing",
        role_mappings={"cover": "slide:1"},
    )

    assert pack["layout_coverage"]["cover"]["status"] == "native"
    for role in ("raw", "section", "quote", "agenda", "end"):
        coverage = pack["layout_coverage"][role]
        assert coverage["status"] == "unsupported"
        assert "layout_id" not in coverage
    assert set(pack["layouts"]) == {"role-cover"}
    assert pack["status"] == "draft"
    assert pack["policies"]["missing_layout_behavior"] == "block"


def test_explicit_derived_role_is_low_confidence_and_approval_required(tmp_path: Path):
    source = _template(tmp_path / "derived.pptx", slides=2)
    dossier = et.build_dossier(source, tmp_path / "extract")
    pack = et.build_template_pack(
        dossier,
        template_id="derived",
        role_mappings={"raw": "slide:2"},
        derives={"quote": "raw"},
    )

    coverage = pack["layout_coverage"]["quote"]
    assert coverage["status"] == "derived"
    assert coverage["confidence"] <= 0.55
    layout = pack["layouts"][coverage["layout_id"]]
    assert layout["status"] == "derived"
    assert layout["source"]["approval_required"] is True
    assert all(slot["approval_required"] for slot in layout["slots"])


def test_inherited_master_typography_is_resolved_without_rewriting(tmp_path: Path):
    source = _template(tmp_path / "typography.pptx", slides=1)
    dossier, pack, _ = et.extract_template(
        source,
        tmp_path / "typography-out",
        role_mappings={"cover": "slide:1"},
    )
    title = next(
        slot for slot in pack["layouts"]["role-cover"]["slots"]
        if slot["placeholder_type"] in {"title", "center_title"}
    )
    # Default Office master title is 44 pt. The slide/layout carries only
    # partial overrides, so this proves the extractor resolves the inheritance
    # chain rather than dropping the template's font size or inventing one.
    assert title["style"]["font_size_pt"] == 44.0
    assert title["style"]["font_size"] > 0
    assert "color_source" in title["style"]
    assert title["resolved_text_styles"][0]["size_pt"] == 44.0
    assert title["text_styles"] != title["resolved_text_styles"]
    assert dossier["masters"][0]["text_styles"]["title"][0]["size_pt"] == 44.0


def test_master_theme_clrmap_resolves_tx1_to_concrete_color(tmp_path: Path):
    source = _template(tmp_path / "theme-context.pptx", slides=1)
    dossier, pack, _ = et.extract_template(
        source,
        tmp_path / "theme-context-out",
        role_mappings={"cover": "slide:1"},
    )
    master = dossier["masters"][0]
    assert master["theme_part"].startswith("ppt/theme/")
    assert master["color_map"]["tx1"] == "dk1"
    title = next(
        slot for slot in pack["layouts"]["role-cover"]["slots"]
        if slot["semantic_name"] == "title"
    )
    color = title["style"]["color_source"]["color"]
    assert color["kind"] in {"rgb", "system"}
    assert color["source_theme_token"] == "tx1"
    assert color["resolved_theme_token"] == "dk1"
    assert title["style"]["color"].startswith("#")
    assert any(token["token"] == "tx1" for token in pack["tokens"]["colors"])


def test_theme_helpers_are_contextual_not_first_theme_global():
    theme_a = {
        "part": "ppt/theme/theme1.xml",
        "colors": {"accent1": {"kind": "rgb", "value": "111111", "transforms": []}},
        "fonts": {"major": {"latin": "Theme A"}, "minor": {"latin": "Theme A Body"}},
    }
    theme_b = {
        "part": "ppt/theme/theme2.xml",
        "colors": {"accent1": {"kind": "rgb", "value": "FF0000", "transforms": []}},
        "fonts": {"major": {"latin": "Theme B"}, "minor": {"latin": "Theme B Body"}},
    }
    family, token = et._theme_font_family(theme_b, "title")
    assert (family, token) == ("Theme B", "major.latin")
    resolved = et._resolve_color_fact(
        {"kind": "theme", "value": "tx1", "transforms": []},
        theme_b,
        {"tx1": "accent1"},
    )
    assert resolved["value"] == "FF0000"
    assert resolved["theme_part"] == "ppt/theme/theme2.xml"


def test_cover_freeform_textboxes_get_explicit_reviewable_semantics(tmp_path: Path):
    source = tmp_path / "freeform-cover.pptx"
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(3)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(0.7), Inches(12), Inches(0.7))
    title.text_frame.paragraphs[0].text = "Template title"
    title.text_frame.paragraphs[0].runs[0].font.size = Pt(34)
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(1.7), Inches(10), Inches(0.4))
    subtitle.text_frame.paragraphs[0].text = "Template subtitle"
    subtitle.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    prs.save(source)

    _, pack, _ = et.extract_template(
        source,
        tmp_path / "freeform-cover-out",
        role_mappings={"cover": "slide:1"},
    )
    layout = pack["layouts"]["role-cover"]
    mapped = {
        slot["semantic_name"]: slot
        for slot in layout["slots"]
        if slot.get("semantic_name") in {"title", "subtitle"}
    }
    assert set(mapped) == {"title", "subtitle"}
    assert mapped["title"]["semantic_mapping_method"] == "cover-freeform-largest-text"
    assert mapped["subtitle"]["semantic_mapping_method"] == "cover-freeform-secondary-text"
    assert mapped["title"]["approval_required"] is True
    assert layout["needs_confirmation"]
    assert any("role-cover:" in item for item in pack["extraction_report"]["needs_confirmation"])


def test_group_children_are_composed_to_absolute_slide_geometry():
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(3)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    group.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(0.5), Inches(2), Inches(0.5),
    )
    group.left = Inches(4)
    group.top = Inches(1)
    group.width = Inches(6)
    group.height = Inches(1)
    fact = et._shape_fact(
        group,
        layer="layout",
        owner_id="layout-test",
        z_index=0,
        canvas=et._canvas_facts(prs),
        asset_by_sha={},
        asset_path_by_id={},
    )
    child = fact["children"][0]
    assert child["source_geometry_emu"] != child["geometry_emu"]
    assert child["geometry_emu"]["x"] == int(Inches(4))
    assert child["geometry_emu"]["y"] == int(Inches(1))
    assert child["geometry"]["x"] == 480.0
    assert child["geometry"]["y"] == 120.0


def test_source_stack_plane_uses_real_placeholder_shape_tree_order():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1), Inches(1), Inches(1),
    )
    canvas = et._canvas_facts(prs)

    # Added after both title/subtitle placeholders -> source foreground.
    foreground = et._container_fixed_elements(
        slide,
        layer="layout",
        owner_id="foreground",
        canvas=canvas,
        asset_by_sha={},
        asset_path_by_id={},
    )[0]
    assert foreground["source_stack_plane"] == "foreground"
    assert foreground["source_stack_evidence"]["placeholder_z_indices"] == [0, 1]

    # Move the same OOXML node between the two placeholders -> interleaved.
    tree = slide.shapes._spTree
    tree.remove(rect._element)
    title_index = tree.index(slide.shapes.title._element)
    tree.insert(title_index + 1, rect._element)
    placeholders = et._container_placeholders(
        slide,
        layer="layout",
        owner_id="interleaved",
        canvas=canvas,
    )
    assert [slot["z_index"] for slot in placeholders] == [0, 2]
    interleaved = et._container_fixed_elements(
        slide,
        layer="layout",
        owner_id="interleaved",
        canvas=canvas,
        asset_by_sha={},
        asset_path_by_id={},
    )[0]
    assert interleaved["z_index"] == 1
    assert interleaved["source_stack_plane"] == "interleaved"

    # Move before every placeholder -> source background.
    tree.remove(rect._element)
    title_index = tree.index(slide.shapes.title._element)
    tree.insert(title_index, rect._element)
    background = et._container_fixed_elements(
        slide,
        layer="layout",
        owner_id="background",
        canvas=canvas,
        asset_by_sha={},
        asset_path_by_id={},
    )[0]
    assert background["source_stack_plane"] == "background"


def test_no_placeholder_defaults_background_and_group_children_inherit_plane():
    prs = Presentation()
    blank = prs.slides.add_slide(prs.slide_layouts[6])
    group = blank.shapes.add_group_shape()
    group.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1), Inches(1), Inches(1), Inches(1),
    )
    fact = et._container_fixed_elements(
        blank,
        layer="master",
        owner_id="group-stack",
        canvas=et._canvas_facts(prs),
        asset_by_sha={},
        asset_path_by_id={},
    )[0]
    assert fact["source_stack_plane"] == "background"
    child = fact["children"][0]
    assert child["source_stack_plane"] == "background"
    assert child["source_stack_evidence"]["inherited_from"] == fact["id"]


def test_interleaved_fixed_vi_enters_layout_review_gate(tmp_path: Path):
    source = _template(tmp_path / "interleaved-review.pptx", slides=1)
    dossier = et.build_dossier(source, tmp_path / "interleaved-review-extract")
    dossier["masters"][0]["fixed_elements"].append({
        "id": "master-interleaved-vi",
        "source_stack_plane": "interleaved",
    })
    pack = et.build_template_pack(
        dossier,
        template_id="interleaved-review",
        role_mappings={"cover": "slide:1"},
    )
    needs = pack["layouts"]["role-cover"]["needs_confirmation"]
    assert any("master-interleaved-vi" in item for item in needs)
    assert any(
        "master-interleaved-vi" in item
        for item in pack["extraction_report"]["needs_confirmation"]
    )


def test_section_title_placeholder_remains_the_section_title(tmp_path: Path):
    source = _template(tmp_path / "section-title.pptx", slides=2)
    dossier = et.build_dossier(source, tmp_path / "section-extract")
    pack = et.build_template_pack(
        dossier,
        template_id="section-title",
        role_mappings={"section": "slide:2"},
    )
    title = next(
        slot for slot in pack["layouts"]["role-section"]["slots"]
        if slot.get("slot_kind") == "title"
    )
    assert title["semantic_name"] == "title"
    assert title["approval_required"] is True


def test_embedded_media_is_pack_relative_and_not_a_slide_screenshot(tmp_path: Path):
    source = _template(tmp_path / "media.pptx", slides=1)
    image = tmp_path / "pixel.png"
    image.write_bytes(b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8"
        "z8BQDwAFgwJ/lEZz3gAAAABJRU5ErkJggg=="
    ))
    prs = Presentation(source)
    prs.slides[0].shapes.add_picture(str(image), Inches(1), Inches(1), Inches(2), Inches(1))
    prs.save(source)

    out = tmp_path / "media-out"
    dossier, pack, _ = et.extract_template(
        source,
        out,
        role_mappings={"cover": "slide:1"},
    )
    media_slot = next(
        slot for slot in pack["layouts"]["role-cover"]["slots"]
        if slot["slot_kind"] == "freeform_image"
    )
    assert media_slot["src"].startswith("assets/")
    assert (out / media_slot["src"]).is_file()
    assert media_slot["asset_id"] in {asset["id"] for asset in dossier["assets"]}
    assert pack["policies"]["whole_page_screenshots"] == "forbidden"


def test_cli_refuses_bad_alias_without_writing_pack(tmp_path: Path):
    source = _template(tmp_path / "cli.pptx", slides=1)
    out = tmp_path / "cli-out"
    result = subprocess.run(
        [
            sys.executable,
            str(ASSETS / "extract_template.py"),
            str(source),
            str(out),
            "--role",
            "cover=slide:1",
            "--alias",
            "agenda=raw",
        ],
        capture_output=True,
        text=True,
    )
    assert result.returncode == 2
    assert "unsupported" in result.stderr
    assert not (out / "template-pack.json").exists()


def test_emitted_json_is_parseable_and_has_six_coverage_keys(tmp_path: Path):
    source = _template(tmp_path / "json.pptx", slides=1)
    out = tmp_path / "json-out"
    et.extract_template(source, out, role_mappings={"cover": "layout:1"})
    dossier = json.loads((out / "template-dossier.json").read_text(encoding="utf-8"))
    pack = json.loads((out / "template-pack.json").read_text(encoding="utf-8"))
    assert dossier["untrusted"] is True
    assert tuple(pack["layout_coverage"].keys()) == et.ROLES
