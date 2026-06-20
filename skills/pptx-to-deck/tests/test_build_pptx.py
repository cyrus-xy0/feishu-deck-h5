"""Tests for build_pptx.py (PPTX → deck.json canvas reconstruction backend).

Run with the skill's interpreter that has python-pptx + lxml installed, e.g.:
    skills/pptx-to-deck/.venv/bin/python -m pytest skills/pptx-to-deck/tests/ -q

Skips gracefully when python-pptx is not importable so a bare `python3` does not
hard-fail the suite.

Regression coverage:
  · H2 — build_font_scheme() must populate _FONT_SCHEME. It calls
    etree.fromstring(theme_part.blob); `etree` used to be imported only LOCALLY
    inside build_theme_map(), so the call raised NameError on every invocation
    and the broad try/except left _FONT_SCHEME permanently empty (theme-font
    runs silently lost their typeface). This test fails on that regression
    because the default python-pptx template ships a fontScheme.
"""
import sys
from pathlib import Path

import pytest

# build_pptx lives in ../assets relative to this tests/ dir.
ASSETS = Path(__file__).resolve().parent.parent / "assets"
if str(ASSETS) not in sys.path:
    sys.path.insert(0, str(ASSETS))

# python-pptx is the hard dependency of build_pptx; skip the whole module if it
# (or build_pptx's other imports) cannot be imported in this interpreter.
bp = pytest.importorskip("build_pptx", reason="python-pptx not importable")
Presentation = pytest.importorskip("pptx").Presentation


def test_build_font_scheme_populates_font_scheme():
    """H2 regression: default template has a fontScheme, so build_font_scheme()
    must leave _FONT_SCHEME non-empty (it was permanently empty when `etree` was
    only imported function-locally and the call NameError'd silently)."""
    prs = Presentation()  # default python-pptx template ships a theme fontScheme
    bp.build_font_scheme(prs)
    assert bp._FONT_SCHEME, (
        "build_font_scheme left _FONT_SCHEME empty — the etree import / theme "
        "fontScheme parse is broken (H2 regression)."
    )
    # default Office theme major/minor latin fonts are present
    assert "mj-lt" in bp._FONT_SCHEME
    assert "mn-lt" in bp._FONT_SCHEME


def test_resolve_theme_font_after_build():
    """A '+mn-lt' theme-font reference resolves to a real typeface once
    build_font_scheme has populated the scheme (and passes plain names through)."""
    prs = Presentation()
    bp.build_font_scheme(prs)
    resolved = bp._resolve_theme_font("+mn-lt")
    assert resolved and resolved == bp._FONT_SCHEME["mn-lt"]
    # a non-reference (no leading '+') is returned verbatim
    assert bp._resolve_theme_font("Arial") == "Arial"


def test_pic_crop_parses_src_rect():
    """_pic_crop reads an <a:srcRect> off shape._element and returns
    [l, r, t, b] crop fractions (vals are 1/1000 of a percent)."""
    from lxml import etree
    from pptx.oxml.ns import qn

    # minimal element tree carrying a srcRect: l=10%, r=5%, t=0, b=20%
    pic = etree.SubElement(etree.Element(qn("p:pic")), qn("p:blipFill"))
    src = etree.SubElement(pic, qn("a:srcRect"))
    src.set("l", "10000")   # 10.000%
    src.set("r", "5000")    # 5.000%
    src.set("b", "20000")   # 20.000%

    class _FakeShape:
        _element = pic.getparent()

    crop = bp._pic_crop(_FakeShape())
    assert crop == [0.1, 0.05, 0.0, 0.2]


def test_pic_crop_none_when_absent():
    """No <a:srcRect> → None (full image, object-fit handled elsewhere)."""
    from lxml import etree
    from pptx.oxml.ns import qn

    class _FakeShape:
        _element = etree.Element(qn("p:pic"))

    assert bp._pic_crop(_FakeShape()) is None


def _shape_with_line():
    """A real autoshape on a real slide so shape.line round-trips through
    python-pptx (needed to exercise the noFill XML path + the color fallback)."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    return slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(914400), Emu(914400))


def test_border_obj_nofill_returns_none():
    """Regression: <a:ln w=...><a:noFill/> means NO border. Before the fix
    _border_obj only checked width>0, failed to resolve a color, and fabricated
    a phantom #888888 box around every such shape (every text box got a gray
    border on import)."""
    shp = _shape_with_line()
    shp.line.width = bp.EMU_PER_PT * 1  # width>0 …
    shp.line.fill.background()          # …but explicit <a:noFill/>

    assert bp._line_is_nofill(shp) is True
    assert bp._border_obj(shp, None) is None, (
        "noFill line must yield no border, not a phantom #888888 one"
    )


def test_border_obj_real_border_preserved():
    """A genuinely-colored line still produces a border (the noFill fix must not
    strip real borders)."""
    from pptx.dml.color import RGBColor

    shp = _shape_with_line()
    shp.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    shp.line.width = bp.EMU_PER_PT * 2

    assert bp._line_is_nofill(shp) is False
    border = bp._border_obj(shp, None)
    assert border is not None and border["color"].lower() == "#ff0000"


def test_line_is_nofill_false_when_no_ln():
    """No <a:ln> at all → not a noFill line (don't accidentally suppress the
    width-based border path)."""
    from lxml import etree
    from pptx.oxml.ns import qn

    class _FakeShape:
        _element = etree.SubElement(etree.Element(qn("p:sp")), qn("p:spPr"))

    # spPr present but no a:ln child
    fake = _FakeShape()
    fake._element = fake._element.getparent()
    assert bp._line_is_nofill(fake) is False


# ── letterbox-seam ("黑边") prevention: _fullbleed_frame_css ─────────────────────
# A full-bleed PPTX page paints its bg only inside the 16:9 .slide; in present
# mode the .slide-frame fills the viewport and shows the GENERIC content-bg in the
# letterbox → a seam on non-16:9 screens. compose_slide bakes a per-slide
# custom_css that mirrors the page's OWN backing onto the frame. The framework
# runtime heal (markBleedPanels) only sees CSS backgrounds, not these <img>/rect
# elements, so this static rule is what actually closes the seam for imports.

_FULL_IMG = {"id": "i1", "type": "image", "src": "input/bg-001.jpg",
             "x": 0, "y": 0, "w": 1920, "h": 1080}


def test_fullbleed_frame_css_image_emits_frame_rule():
    css = bp._fullbleed_frame_css("slide-007", [_FULL_IMG])
    assert css is not None
    assert bp._FRAME_BG_MARKER in css
    # targets the FRAME (parent of .slide) for THIS slide key, present mode only
    assert '.deck[data-mode="present"] .slide-frame' in css
    assert ':has(> .slide[data-slide-key="slide-007"])' in css
    assert 'url("input/bg-001.jpg")' in css
    assert "background-size: cover" in css


def test_fullbleed_frame_css_layers_top_first():
    """Element order is bottom→top; CSS background-image lists TOP first, so the
    emitted layer order must be REVERSED (top image wins the paint)."""
    bottom = dict(_FULL_IMG, id="b", src="input/bottom.jpg")
    top = dict(_FULL_IMG, id="t", src="input/top.png")
    css = bp._fullbleed_frame_css("slide-001", [bottom, top])  # bottom listed first
    i_top = css.index('url("input/top.png")')
    i_bot = css.index('url("input/bottom.jpg")')
    assert i_top < i_bot, "top layer must come first in CSS background-image"


def test_fullbleed_frame_css_solid_only_flat_fills():
    """A page whose only full-bleed backing is a solid rect → flat-fill the frame
    that colour (prevents a coloured slide seaming against the dark letterbox)."""
    rect = {"id": "r", "type": "shape", "kind": "rect", "fill": "#0B0F18",
            "x": 0, "y": 0, "w": 1920, "h": 1080}
    css = bp._fullbleed_frame_css("slide-003", [rect])
    assert css is not None
    assert "background-color: #0B0F18" in css
    assert "background-image: none" in css


def test_fullbleed_frame_css_solid_backs_image_stack():
    """When both a full-bleed solid and image exist, the solid is the base colour
    under the image stack (not #000)."""
    rect = {"id": "r", "type": "shape", "kind": "rect", "fill": "#123456",
            "x": 0, "y": 0, "w": 1920, "h": 1080}
    css = bp._fullbleed_frame_css("slide-004", [rect, _FULL_IMG])
    assert "background-color: #123456" in css
    assert 'url("input/bg-001.jpg")' in css


def test_fullbleed_frame_css_none_when_not_fullbleed():
    """A page with no full-bleed backing (small/inset elements only) gets NO
    custom_css — it keeps the framework's generic letterbox."""
    small = {"id": "s", "type": "image", "src": "input/icon.png",
             "x": 800, "y": 400, "w": 320, "h": 240}
    assert bp._fullbleed_frame_css("slide-009", [small]) is None
    assert bp._fullbleed_frame_css("slide-010", []) is None
