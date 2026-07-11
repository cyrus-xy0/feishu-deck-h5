#!/usr/bin/env python3
"""Extract a PowerPoint template into a reviewable design-system draft.

The extractor is intentionally different from ``build_pptx.py``:

* it reads presentation, master, layout, placeholder, theme and media facts;
* it preserves source geometry in EMU and proposes a ratio-matched H5 canvas;
* it emits a factual ``template-dossier.json`` and a *draft*
  ``template-pack.json``;
* it never renders or stores a whole-slide screenshot;
* it never invents missing semantic layouts.  The six supported roles are
  explicit and missing roles remain ``unsupported`` until the caller maps,
  derives, or aliases them.

Run with the sibling skill interpreter (python-pptx + lxml):

    skills/pptx-to-deck/.venv/bin/python3 \
      skills/pptx-to-deck/assets/extract_template.py template.pptx out/ \
      --role cover=slide:1 --role raw=slide:2 \
      --alias agenda=raw --alias end=cover

Selectors are 1-based: ``slide:N``, ``layout:N`` or
``layout-name:Exact PowerPoint layout name``.  ``--derive quote=raw`` is an
explicit request to seed a new, approval-required layout from another role.
"""

from __future__ import annotations

import argparse
import copy
import hashlib
import html
import json
import os
import posixpath
import re
import sys
import zipfile
from datetime import datetime, timezone
from fractions import Fraction
from pathlib import Path, PurePosixPath
from typing import Any, Iterable

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


SCHEMA_VERSION = "1.0"
PACK_VERSION = "0.1.0"
ROLES = ("cover", "raw", "section", "quote", "agenda", "end")
ROLE_LABELS = {
    "cover": "封面",
    "raw": "内容",
    "section": "章节",
    "quote": "金句",
    "agenda": "目录",
    "end": "封底",
}

EMU_PER_INCH = 914400
EMU_PER_POINT = 12700
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


class ExtractionError(ValueError):
    """A caller-visible contract error; never silently downgraded."""


def _now() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def _sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def _slug(value: str) -> str:
    value = re.sub(r"[^A-Za-z0-9\u4e00-\u9fff]+", "-", value.strip()).strip("-")
    return value.lower() or "pptx-template"


def _enum_name(value: Any) -> str:
    name = getattr(value, "name", None)
    if name:
        return str(name).lower()
    text = str(value or "unknown")
    text = re.sub(r"\s*\([^)]*\)\s*$", "", text)
    return re.sub(r"[^a-z0-9]+", "_", text.lower()).strip("_") or "unknown"


def _part_name(obj: Any) -> str | None:
    try:
        return str(obj.part.partname).lstrip("/")
    except Exception:
        return None


def _element(obj: Any) -> Any:
    el = getattr(obj, "element", None)
    if el is None:
        el = getattr(obj, "_element", None)
    return el


class _Xf:
    """Compose group-local OOXML coordinates into absolute slide EMU."""

    def __init__(self, ox: float = 0, oy: float = 0, sx: float = 1, sy: float = 1):
        self.ox, self.oy, self.sx, self.sy = ox, oy, sx, sy

    def x(self, value: float) -> float:
        return self.ox + value * self.sx

    def y(self, value: float) -> float:
        return self.oy + value * self.sy

    def w(self, value: float) -> float:
        return value * self.sx

    def h(self, value: float) -> float:
        return value * self.sy

    def enter_group(self, group: Any) -> "_Xf":
        gx, gy = self.x(float(group.left or 0)), self.y(float(group.top or 0))
        gw, gh = self.w(float(group.width or 0)), self.h(float(group.height or 0))
        child_x = child_y = 0.0
        child_w = float(group.width or 1)
        child_h = float(group.height or 1)
        el = _element(group)
        grp_sp_pr = el.find(qn("p:grpSpPr")) if el is not None else None
        xfrm = grp_sp_pr.find(qn("a:xfrm")) if grp_sp_pr is not None else None
        if xfrm is not None:
            child_off = xfrm.find(qn("a:chOff"))
            child_ext = xfrm.find(qn("a:chExt"))
            if child_off is not None:
                child_x = float(child_off.get("x", 0))
                child_y = float(child_off.get("y", 0))
            if child_ext is not None:
                child_w = float(child_ext.get("cx", 1)) or 1
                child_h = float(child_ext.get("cy", 1)) or 1
        sx, sy = gw / child_w, gh / child_h
        return _Xf(gx - child_x * sx, gy - child_y * sy, sx, sy)


def _source_geometry(shape: Any) -> dict[str, int]:
    def value(name: str) -> int:
        try:
            return int(getattr(shape, name) or 0)
        except Exception:
            return 0

    return {
        "x": value("left"),
        "y": value("top"),
        "w": value("width"),
        "h": value("height"),
    }


def _geometry(shape: Any, canvas: dict[str, Any], xf: _Xf | None = None) -> dict[str, Any]:
    local = _source_geometry(shape)
    xf = xf or _Xf()
    x = int(round(xf.x(local["x"])))
    y = int(round(xf.y(local["y"])))
    w = int(round(xf.w(local["w"])))
    h = int(round(xf.h(local["h"])))
    sw, sh = canvas["source_width_emu"], canvas["source_height_emu"]
    return {
        "x": x,
        "y": y,
        "w": w,
        "h": h,
        "normalized": {
            "x": round(x / sw, 8) if sw else 0,
            "y": round(y / sh, 8) if sh else 0,
            "w": round(w / sw, 8) if sw else 0,
            "h": round(h / sh, 8) if sh else 0,
        },
    }


def _design_geometry(geometry_emu: dict[str, Any], canvas: dict[str, Any]) -> dict[str, float]:
    """Project exact EMU geometry onto the ratio-matched H5 design canvas."""
    dw = canvas["recommended_design_canvas"]["width"]
    dh = canvas["recommended_design_canvas"]["height"]
    sw, sh = canvas["source_width_emu"], canvas["source_height_emu"]
    return {
        "x": round(geometry_emu["x"] * dw / sw, 3),
        "y": round(geometry_emu["y"] * dh / sh, 3),
        "w": round(geometry_emu["w"] * dw / sw, 3),
        "h": round(geometry_emu["h"] * dh / sh, 3),
    }


def _color_fact(parent: Any) -> dict[str, Any] | None:
    if parent is None:
        return None
    for tag, kind, attr in (
        ("srgbClr", "rgb", "val"),
        ("schemeClr", "theme", "val"),
        ("sysClr", "system", "lastClr"),
        ("prstClr", "preset", "val"),
    ):
        node = parent if etree.QName(parent).localname == tag else parent.find(qn(f"a:{tag}"))
        if node is None:
            continue
        value = node.get(attr) or node.get("val")
        transforms = []
        for child in node:
            transforms.append({
                "name": etree.QName(child).localname,
                "value": child.get("val"),
            })
        return {"kind": kind, "value": value, "transforms": transforms}
    return None


def _fill_fact(node: Any) -> dict[str, Any] | None:
    if node is None:
        return None
    for tag in ("noFill", "solidFill", "gradFill", "blipFill", "pattFill"):
        fill = node.find(qn(f"a:{tag}"))
        if fill is None:
            continue
        if tag == "noFill":
            return {"type": "none"}
        if tag == "solidFill":
            return {"type": "solid", "color": _color_fact(fill)}
        if tag == "gradFill":
            stops = []
            for stop in fill.findall(f".//{qn('a:gs')}"):
                stops.append({
                    "position": int(stop.get("pos", "0")),
                    "color": _color_fact(stop),
                })
            lin = fill.find(qn("a:lin"))
            return {
                "type": "gradient",
                "angle": int(lin.get("ang")) if lin is not None and lin.get("ang") else None,
                "stops": stops,
            }
        if tag == "blipFill":
            return {"type": "picture"}
        return {"type": "pattern"}
    return None


def _line_fact(node: Any) -> dict[str, Any] | None:
    if node is None:
        return None
    line = node.find(qn("a:ln"))
    if line is None:
        return None
    return {
        "width_emu": int(line.get("w", "0")),
        "fill": _fill_fact(line),
        "dash": (line.find(qn("a:prstDash")).get("val")
                 if line.find(qn("a:prstDash")) is not None else None),
    }


def _rpr_fact(node: Any) -> dict[str, Any]:
    fact: dict[str, Any] = {
        "size_pt": round(int(node.get("sz")) / 100, 2) if node.get("sz") else None,
        "bold": node.get("b"),
        "italic": node.get("i"),
        "language": node.get("lang"),
        "fill": _fill_fact(node),
    }
    for key, tag in (("latin", "latin"), ("east_asian", "ea"), ("complex", "cs")):
        font = node.find(qn(f"a:{tag}"))
        if font is not None and font.get("typeface"):
            fact[key] = font.get("typeface")
    return fact


def _dedupe_styles(styles: Iterable[dict[str, Any]]) -> list[dict[str, Any]]:
    output = []
    seen: set[str] = set()
    for fact in styles:
        encoded = json.dumps(fact, ensure_ascii=False, sort_keys=True)
        if encoded not in seen:
            seen.add(encoded)
            output.append(fact)
    return output


def _text_styles(shape: Any) -> list[dict[str, Any]]:
    """Return exact style facts present in the OOXML; do not synthesize defaults."""
    el = _element(shape)
    if el is None:
        return []
    nodes = (
        el.findall(f".//{qn('a:rPr')}")
        + el.findall(f".//{qn('a:defRPr')}")
        + el.findall(f".//{qn('a:endParaRPr')}")
    )
    return _dedupe_styles(_rpr_fact(node) for node in nodes)


def _master_text_styles(master: Any) -> dict[str, list[dict[str, Any]]]:
    """Exact master txStyles by title/body/other, including paragraph levels."""
    el = _element(master)
    output: dict[str, list[dict[str, Any]]] = {"title": [], "body": [], "other": []}
    if el is None:
        return output
    tx_styles = el.find(qn("p:txStyles"))
    if tx_styles is None:
        return output
    for bucket, tag in (("title", "p:titleStyle"), ("body", "p:bodyStyle"), ("other", "p:otherStyle")):
        branch = tx_styles.find(qn(tag))
        if branch is None:
            continue
        level_styles = []
        for level in range(1, 10):
            paragraph = branch.find(qn(f"a:lvl{level}pPr"))
            if paragraph is None:
                continue
            rpr = paragraph.find(qn("a:defRPr"))
            if rpr is None:
                continue
            fact = _rpr_fact(rpr)
            fact["paragraph_level"] = level - 1
            level_styles.append(fact)
        output[bucket] = _dedupe_styles(level_styles)
    return output


def _runtime_color(fill: dict[str, Any] | None) -> str | None:
    if not fill or fill.get("type") != "solid":
        return None
    color = fill.get("color") or {}
    if color.get("kind") in {"rgb", "system"} and color.get("value"):
        return f"#{color['value']}"
    if color.get("kind") == "theme" and color.get("value"):
        return f"theme:{color['value']}"
    return color.get("value")


def _runtime_style_from_facts(
    styles: list[dict[str, Any]],
    canvas: dict[str, Any],
    *,
    alignment: str | None = None,
    line_height: dict[str, Any] | None = None,
) -> dict[str, Any]:
    first = styles[0] if styles else {}
    family = first.get("east_asian") or first.get("latin") or first.get("complex")
    size_pt = first.get("size_pt")
    size_px = None
    if size_pt is not None:
        size_px = round(
            size_pt * EMU_PER_POINT
            * canvas["recommended_design_canvas"]["width"]
            / canvas["source_width_emu"],
            3,
        )
    bold = first.get("bold")
    weight = 700 if str(bold).lower() in {"1", "true"} else 400 if bold is not None else None
    return {
        "font_family": family,
        "font_size": size_px,
        "font_size_pt": size_pt,
        "font_weight": weight,
        "color": _runtime_color(first.get("fill")),
        "color_source": copy.deepcopy(first.get("fill")),
        "alignment": alignment,
        "line_height": line_height,
    }


def _runtime_text_style(shape: Any, canvas: dict[str, Any]) -> dict[str, Any]:
    """Renderer-friendly style view; exact source facts remain in text_styles."""
    styles = _text_styles(shape)
    alignment = None
    line_height = None
    try:
        paragraph = shape.text_frame.paragraphs[0]
        alignment = _enum_name(paragraph.alignment) if paragraph.alignment is not None else None
        ppr = paragraph._p.find(qn("a:pPr"))
        if ppr is not None:
            pct = ppr.find(f".//{qn('a:spcPct')}")
            pts = ppr.find(f".//{qn('a:spcPts')}")
            if pct is not None and pct.get("val"):
                line_height = {"unit": "percent", "value": int(pct.get("val")) / 1000}
            elif pts is not None and pts.get("val"):
                line_height = {"unit": "pt", "value": int(pts.get("val")) / 100}
    except Exception:
        pass
    return _runtime_style_from_facts(
        styles,
        canvas,
        alignment=alignment,
        line_height=line_height,
    )


def _related_asset_id(shape: Any, asset_by_sha: dict[str, str]) -> str | None:
    el = _element(shape)
    if el is None:
        return None
    for blip in el.findall(f".//{qn('a:blip')}"):
        rel_id = blip.get(qn("r:embed"))
        if not rel_id:
            continue
        try:
            blob = shape.part.related_part(rel_id).blob
            return asset_by_sha.get(_sha256(blob))
        except Exception:
            continue
    return None


def _shape_fact(
    shape: Any,
    *,
    layer: str,
    owner_id: str,
    z_index: int,
    canvas: dict[str, Any],
    asset_by_sha: dict[str, str],
    asset_path_by_id: dict[str, str],
    xf: _Xf | None = None,
) -> dict[str, Any]:
    shape_id = getattr(shape, "shape_id", z_index + 1)
    xf = xf or _Xf()
    geometry_emu = _geometry(shape, canvas, xf)
    fact: dict[str, Any] = {
        "id": f"{owner_id}-shape-{shape_id}",
        "source_layer": layer,
        "source_part": _part_name(shape),
        "shape_id": int(shape_id),
        "name": getattr(shape, "name", None),
        "shape_type": _enum_name(getattr(shape, "shape_type", None)),
        "z_index": z_index,
        "source_geometry_emu": _source_geometry(shape),
        "geometry_emu": geometry_emu,
        "geometry": _design_geometry(geometry_emu, canvas),
        "rotation_degrees": float(getattr(shape, "rotation", 0) or 0),
        "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
    }
    el = _element(shape)
    sppr = None
    if el is not None:
        for tag in ("p:spPr", "p:grpSpPr"):
            sppr = el.find(qn(tag))
            if sppr is not None:
                break
    style = {
        "fill": _fill_fact(sppr),
        "line": _line_fact(sppr),
        "text_styles": _text_styles(shape),
        "text": _runtime_text_style(shape, canvas),
    }
    if any(value not in (None, [], {}) for value in style.values()):
        fact["style"] = style
    try:
        if shape.has_text_frame:
            fact["text"] = shape.text_frame.text
    except Exception:
        pass
    asset_id = _related_asset_id(shape, asset_by_sha)
    if asset_id:
        fact["asset_id"] = asset_id
        fact["src"] = asset_path_by_id.get(asset_id)
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        child_xf = xf.enter_group(shape)
        fact["children"] = [
            _shape_fact(
                child,
                layer=layer,
                owner_id=fact["id"],
                z_index=i,
                canvas=canvas,
                asset_by_sha=asset_by_sha,
                asset_path_by_id=asset_path_by_id,
                xf=child_xf,
            )
            for i, child in enumerate(shape.shapes)
        ]
    return fact


def _slot_kind(placeholder_type: str) -> tuple[str, bool]:
    if placeholder_type in {"title", "center_title", "vertical_title"}:
        return "title", False
    if placeholder_type in {"subtitle"}:
        return "subtitle", False
    if placeholder_type in {"body", "object", "vertical_body"}:
        return "body", False
    if placeholder_type in {"picture", "media_clip"}:
        return "image", False
    if placeholder_type in {"chart", "table", "smart_art"}:
        return "visualization", False
    if placeholder_type == "date":
        return "date", True
    if placeholder_type == "footer":
        return "footer", True
    if placeholder_type == "slide_number":
        return "slide_number", True
    if placeholder_type == "header":
        return "header", True
    return "content", False


def _placeholder_fact(
    shape: Any,
    *,
    layer: str,
    owner_id: str,
    canvas: dict[str, Any],
    z_index: int,
) -> dict[str, Any]:
    try:
        idx = int(shape.placeholder_format.idx)
        placeholder_type = _enum_name(shape.placeholder_format.type)
    except Exception:
        idx = -1
        placeholder_type = "unknown"
    kind, system = _slot_kind(placeholder_type)
    prompt = ""
    try:
        prompt = shape.text_frame.text
    except Exception:
        pass
    geometry_emu = _geometry(shape, canvas)
    return {
        "id": f"{owner_id}-slot-{idx if idx >= 0 else getattr(shape, 'shape_id', 0)}",
        "placeholder_index": idx,
        "placeholder_type": placeholder_type,
        "z_index": z_index,
        "slot_kind": kind,
        "semantic_name": kind,
        "system_field": system,
        "name": getattr(shape, "name", None),
        "prompt_text": prompt,
        "geometry_emu": geometry_emu,
        "geometry": _design_geometry(geometry_emu, canvas),
        "text_styles": _text_styles(shape),
        "style": _runtime_text_style(shape, canvas),
        "source": {
            "layer": layer,
            "part": _part_name(shape),
            "shape_id": int(getattr(shape, "shape_id", 0)),
        },
        "confidence": 1.0,
        "approval_required": False,
    }


def _background_fact(
    container: Any,
    asset_by_sha: dict[str, str],
    asset_path_by_id: dict[str, str],
) -> dict[str, Any] | None:
    el = _element(container)
    if el is None:
        return None
    bg = el.find(f".//{qn('p:bg')}")
    if bg is None:
        return None
    fact: dict[str, Any] = {"source_part": _part_name(container)}
    bg_pr = bg.find(qn("p:bgPr"))
    if bg_pr is not None:
        fact["fill"] = _fill_fact(bg_pr)
    bg_ref = bg.find(qn("p:bgRef"))
    if bg_ref is not None:
        fact["theme_reference"] = {
            "index": int(bg_ref.get("idx", "0")),
            "color": _color_fact(bg_ref),
        }
    for blip in bg.findall(f".//{qn('a:blip')}"):
        rel_id = blip.get(qn("r:embed"))
        if rel_id:
            try:
                blob = container.part.related_part(rel_id).blob
                fact["asset_id"] = asset_by_sha.get(_sha256(blob))
                fact["src"] = asset_path_by_id.get(fact["asset_id"])
            except Exception:
                pass
    return fact


def _background_element(
    background: dict[str, Any] | None,
    *,
    owner_id: str,
    layer: str,
    canvas: dict[str, Any],
) -> dict[str, Any] | None:
    if not background:
        return None
    geometry_emu = {
        "x": 0,
        "y": 0,
        "w": canvas["source_width_emu"],
        "h": canvas["source_height_emu"],
        "normalized": {"x": 0, "y": 0, "w": 1, "h": 1},
    }
    element = {
        "id": f"{owner_id}-background",
        "source_layer": layer,
        "source_part": background.get("source_part"),
        "shape_id": 0,
        "name": "Background",
        "shape_type": "background",
        "z_index": -1,
        "geometry_emu": geometry_emu,
        "geometry": _design_geometry(geometry_emu, canvas),
        "rotation_degrees": 0,
        "is_placeholder": False,
        "style": {"fill": background.get("fill"), "theme_reference": background.get("theme_reference")},
        "fixed_by_source": True,
        "source_stack_plane": "background",
        "source_stack_evidence": {
            "basis": "ppt-background-node",
            "container_part": background.get("source_part"),
            "fixed_z_index": -1,
            "placeholder_z_indices": [],
        },
    }
    if background.get("asset_id"):
        element["asset_id"] = background["asset_id"]
        element["src"] = background.get("src")
    return element


def _source_part_from_rels(rels_path: str) -> str | None:
    path = PurePosixPath(rels_path)
    if path.parent.name != "_rels" or not path.name.endswith(".rels"):
        return None
    return str(path.parent.parent / path.name[:-5])


def _media_usage(zf: zipfile.ZipFile) -> dict[str, list[str]]:
    usage: dict[str, list[str]] = {}
    for rels_path in zf.namelist():
        if not rels_path.endswith(".rels"):
            continue
        source_part = _source_part_from_rels(rels_path)
        if not source_part:
            continue
        try:
            root = etree.fromstring(zf.read(rels_path))
        except Exception:
            continue
        for rel in root.findall(f"{{{REL_NS}}}Relationship"):
            if rel.get("TargetMode") == "External":
                continue
            target = rel.get("Target") or ""
            resolved = posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))
            resolved = resolved.lstrip("/")
            if resolved.startswith("ppt/media/"):
                usage.setdefault(resolved, []).append(source_part)
    return usage


def _extract_assets(pptx_path: Path, output_dir: Path) -> tuple[list[dict[str, Any]], dict[str, str]]:
    assets_dir = output_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    records: list[dict[str, Any]] = []
    by_sha: dict[str, str] = {}
    with zipfile.ZipFile(pptx_path) as zf:
        usage = _media_usage(zf)
        for member in sorted(name for name in zf.namelist() if name.startswith("ppt/media/") and not name.endswith("/")):
            data = zf.read(member)
            digest = _sha256(data)
            if digest in by_sha:
                existing = next(record for record in records if record["id"] == by_sha[digest])
                existing["source_parts"].append(member)
                existing["used_by_parts"] = sorted(set(existing["used_by_parts"] + usage.get(member, [])))
                continue
            asset_id = f"asset-{digest[:12]}"
            filename = Path(member).name
            destination = assets_dir / filename
            if destination.exists() and _sha256(destination.read_bytes()) != digest:
                destination = assets_dir / f"{Path(filename).stem}-{digest[:8]}{Path(filename).suffix}"
            destination.write_bytes(data)
            record = {
                "id": asset_id,
                "path": destination.relative_to(output_dir).as_posix(),
                "source_parts": [member],
                "used_by_parts": sorted(set(usage.get(member, []))),
                "sha256": digest,
                "bytes": len(data),
                "extension": destination.suffix.lower().lstrip("."),
            }
            records.append(record)
            by_sha[digest] = asset_id
    return records, by_sha


def _theme_facts(pptx_path: Path) -> list[dict[str, Any]]:
    themes: list[dict[str, Any]] = []
    with zipfile.ZipFile(pptx_path) as zf:
        for part in sorted(name for name in zf.namelist() if name.startswith("ppt/theme/") and name.endswith(".xml")):
            try:
                root = etree.fromstring(zf.read(part))
            except Exception:
                continue
            colors: dict[str, Any] = {}
            scheme = root.find(f".//{{{DRAWING_NS}}}clrScheme")
            if scheme is not None:
                for child in scheme:
                    colors[etree.QName(child).localname] = _color_fact(child)

            def font_branch(tag: str) -> dict[str, Any]:
                branch = root.find(f".//{{{DRAWING_NS}}}{tag}")
                if branch is None:
                    return {}
                fact: dict[str, Any] = {}
                for key in ("latin", "ea", "cs"):
                    node = branch.find(f"{{{DRAWING_NS}}}{key}")
                    if node is not None:
                        fact[key] = node.get("typeface") or ""
                scripts = {}
                for node in branch.findall(f"{{{DRAWING_NS}}}font"):
                    if node.get("script"):
                        scripts[node.get("script")] = node.get("typeface") or ""
                if scripts:
                    fact["scripts"] = scripts
                return fact

            themes.append({
                "part": part,
                "name": root.get("name"),
                "color_scheme_name": scheme.get("name") if scheme is not None else None,
                "colors": colors,
                "fonts": {
                    "major": font_branch("majorFont"),
                    "minor": font_branch("minorFont"),
                },
            })
    return themes


_DEFAULT_CLR_MAP = {
    "bg1": "lt1",
    "tx1": "dk1",
    "bg2": "lt2",
    "tx2": "dk2",
    "accent1": "accent1",
    "accent2": "accent2",
    "accent3": "accent3",
    "accent4": "accent4",
    "accent5": "accent5",
    "accent6": "accent6",
    "hlink": "hlink",
    "folHlink": "folHlink",
}


def _master_theme_part(master: Any) -> str | None:
    try:
        for rel in master.part.rels.values():
            if str(rel.reltype).endswith("/theme"):
                return str(rel.target_part.partname).lstrip("/")
    except Exception:
        pass
    return None


def _color_map(container: Any, base: dict[str, str] | None = None) -> dict[str, str]:
    """Resolve master clrMap plus optional layout/slide override."""
    result = dict(base or _DEFAULT_CLR_MAP)
    el = _element(container)
    if el is None:
        return result
    direct = el.find(qn("p:clrMap"))
    if direct is not None:
        result.update({str(key): str(value) for key, value in direct.attrib.items()})
    override = el.find(f".//{qn('p:clrMapOvr')}")
    if override is not None:
        mapping = override.find(qn("a:overrideClrMapping"))
        if mapping is not None:
            result.update({str(key): str(value) for key, value in mapping.attrib.items()})
    return result


def _theme_record(themes: list[dict[str, Any]], theme_part: str | None) -> dict[str, Any] | None:
    if theme_part:
        for theme in themes:
            if theme.get("part") == theme_part:
                return theme
    return None


def _theme_context(
    context_id: str,
    theme: dict[str, Any] | None,
    color_map: dict[str, str],
) -> dict[str, Any]:
    return {
        "context_id": context_id,
        "theme_part": theme.get("part") if theme else None,
        "color_map": copy.deepcopy(color_map),
    }


def _resolve_color_fact(
    color: dict[str, Any] | None,
    theme: dict[str, Any] | None,
    color_map: dict[str, str],
) -> dict[str, Any] | None:
    if not color:
        return None
    result = copy.deepcopy(color)
    if result.get("kind") != "theme" or not result.get("value"):
        return result
    source_token = str(result["value"])
    resolved_token = color_map.get(source_token, source_token)
    base = copy.deepcopy((theme or {}).get("colors", {}).get(resolved_token))
    if not isinstance(base, dict):
        result["resolved_theme_token"] = resolved_token
        return result
    base["transforms"] = list(base.get("transforms") or []) + list(result.get("transforms") or [])
    base["source_theme_token"] = source_token
    base["resolved_theme_token"] = resolved_token
    base["theme_part"] = (theme or {}).get("part")
    return base


def _resolve_fill_fact(
    fill: dict[str, Any] | None,
    theme: dict[str, Any] | None,
    color_map: dict[str, str],
) -> dict[str, Any] | None:
    if not fill:
        return None
    result = copy.deepcopy(fill)
    if result.get("type") == "solid":
        result["color"] = _resolve_color_fact(result.get("color"), theme, color_map)
    elif result.get("type") == "gradient":
        for stop in result.get("stops") or []:
            if isinstance(stop, dict):
                stop["color"] = _resolve_color_fact(stop.get("color"), theme, color_map)
    return result


def _apply_theme_to_element(
    element: dict[str, Any],
    *,
    context: dict[str, Any],
    theme: dict[str, Any] | None,
    color_map: dict[str, str],
) -> None:
    element["theme_context"] = copy.deepcopy(context)
    style = element.get("style")
    if isinstance(style, dict):
        if style.get("fill") is not None:
            style["source_fill"] = copy.deepcopy(style["fill"])
            style["fill"] = _resolve_fill_fact(style["fill"], theme, color_map)
        reference = style.get("theme_reference")
        if isinstance(reference, dict) and reference.get("color"):
            style["source_theme_reference"] = copy.deepcopy(reference)
            resolved = _resolve_color_fact(reference.get("color"), theme, color_map)
            style["theme_reference"] = {**reference, "color": resolved}
            if style.get("fill") is None and resolved:
                style["fill"] = {"type": "solid", "color": resolved}
        line = style.get("line")
        if isinstance(line, dict) and line.get("fill") is not None:
            line["source_fill"] = copy.deepcopy(line["fill"])
            line["fill"] = _resolve_fill_fact(line["fill"], theme, color_map)
        text = style.get("text")
        if isinstance(text, dict) and text.get("color_source") is not None:
            text["color_source"] = _resolve_fill_fact(text["color_source"], theme, color_map)
            text["color"] = _runtime_color(text["color_source"])
            text["theme_context"] = copy.deepcopy(context)
    for child in element.get("children") or []:
        if isinstance(child, dict):
            _apply_theme_to_element(
                child,
                context=context,
                theme=theme,
                color_map=color_map,
            )


def _canvas_facts(prs: Any) -> dict[str, Any]:
    width, height = int(prs.slide_width), int(prs.slide_height)
    fraction = Fraction(width, height)
    if fraction.numerator > 100 or fraction.denominator > 100:
        fraction = Fraction(width / height).limit_denominator(100)
    ratio_label = f"{fraction.numerator}:{fraction.denominator}"
    design_width = 1920
    design_height = max(1, round(design_width * height / width))
    orientation = "landscape" if width > height else "portrait" if height > width else "square"
    return {
        "source_width_emu": width,
        "source_height_emu": height,
        "width_inches": round(width / EMU_PER_INCH, 6),
        "height_inches": round(height / EMU_PER_INCH, 6),
        "width_points": round(width / EMU_PER_POINT, 3),
        "height_points": round(height / EMU_PER_POINT, 3),
        "aspect_ratio": {
            "label": ratio_label,
            "value": round(width / height, 8),
        },
        "orientation": orientation,
        "recommended_design_canvas": {
            "width": design_width,
            "height": design_height,
            "policy": "preserve-source-aspect-ratio",
        },
        "design_width": design_width,
        "design_height": design_height,
    }


def _container_fixed_elements(
    container: Any,
    *,
    layer: str,
    owner_id: str,
    canvas: dict[str, Any],
    asset_by_sha: dict[str, str],
    asset_path_by_id: dict[str, str],
) -> list[dict[str, Any]]:
    records = []
    shapes = list(container.shapes)
    placeholder_z_indices = []
    for z_index, shape in enumerate(shapes):
        try:
            if shape.is_placeholder:
                placeholder_z_indices.append(z_index)
        except Exception:
            continue

    def plane_for(z_index: int) -> str:
        if not placeholder_z_indices or z_index < min(placeholder_z_indices):
            return "background"
        if z_index > max(placeholder_z_indices):
            return "foreground"
        return "interleaved"

    def stamp(
        element: dict[str, Any],
        plane: str,
        evidence: dict[str, Any],
        inherited_from: str | None = None,
    ) -> None:
        element["source_stack_plane"] = plane
        child_evidence = copy.deepcopy(evidence)
        if inherited_from:
            child_evidence["inherited_from"] = inherited_from
            child_evidence["child_local_z_index"] = element.get("z_index")
        element["source_stack_evidence"] = child_evidence
        for child in element.get("children") or []:
            if isinstance(child, dict):
                stamp(child, plane, evidence, element["id"])

    for z_index, shape in enumerate(shapes):
        try:
            if shape.is_placeholder:
                continue
        except Exception:
            pass
        fact = _shape_fact(
            shape,
            layer=layer,
            owner_id=owner_id,
            z_index=z_index,
            canvas=canvas,
            asset_by_sha=asset_by_sha,
            asset_path_by_id=asset_path_by_id,
        )
        fact["fixed_by_source"] = True
        plane = plane_for(z_index)
        evidence = {
            "basis": "ooxml-shape-tree-order",
            "container_part": _part_name(container),
            "fixed_z_index": z_index,
            "placeholder_z_indices": list(placeholder_z_indices),
        }
        stamp(fact, plane, evidence)
        records.append(fact)
    return records


def _container_placeholders(
    container: Any,
    *,
    layer: str,
    owner_id: str,
    canvas: dict[str, Any],
) -> list[dict[str, Any]]:
    records = []
    for z_index, shape in enumerate(container.shapes):
        try:
            if not shape.is_placeholder:
                continue
        except Exception:
            continue
        records.append(_placeholder_fact(
            shape,
            layer=layer,
            owner_id=owner_id,
            canvas=canvas,
            z_index=z_index,
        ))
    return records


def _slot_style_bucket(slot: dict[str, Any]) -> str:
    placeholder_type = slot.get("placeholder_type")
    if placeholder_type in {"title", "center_title", "vertical_title"}:
        return "title"
    if placeholder_type in {"body", "object", "vertical_body", "subtitle"}:
        return "body"
    return "other"


def _theme_font_family(theme: dict[str, Any] | None, bucket: str) -> tuple[str | None, str | None]:
    if not theme:
        return None, None
    branch_name = "major" if bucket == "title" else "minor"
    branch = theme.get("fonts", {}).get(branch_name, {})
    scripts = branch.get("scripts", {}) if isinstance(branch.get("scripts"), dict) else {}
    for token, value in (
        (f"{branch_name}.script.Hans", scripts.get("Hans")),
        (f"{branch_name}.ea", branch.get("ea")),
        (f"{branch_name}.latin", branch.get("latin")),
    ):
        if value:
            return value, token
    return None, None


def _enrich_slot_style(
    slot: dict[str, Any],
    *,
    canvas: dict[str, Any],
    theme: dict[str, Any] | None,
    color_map: dict[str, str],
    theme_context: dict[str, Any],
    fallback_styles: list[dict[str, Any]] | None = None,
    fallback_source: str | None = None,
) -> None:
    exact_styles = slot.get("text_styles", [])
    base = copy.deepcopy((fallback_styles or [{}])[0]) if fallback_styles else {}
    for exact in exact_styles:
        for key, value in exact.items():
            if value is not None:
                base[key] = copy.deepcopy(value)
    resolved = [base] if base else []
    slot["resolved_text_styles"] = resolved
    slot["style_source"] = (
        f"direct-over-{fallback_source}" if exact_styles and fallback_styles
        else "direct" if exact_styles
        else fallback_source or "unresolved"
    )
    prior = slot.get("style", {})
    runtime = _runtime_style_from_facts(
        resolved,
        canvas,
        alignment=prior.get("alignment"),
        line_height=prior.get("line_height"),
    )
    if not runtime.get("font_family") or str(runtime.get("font_family")).startswith("+"):
        family, token = _theme_font_family(theme, _slot_style_bucket(slot))
        runtime["font_family"] = family
        runtime["font_family_token"] = token
    if runtime.get("color_source") is not None:
        runtime["color_source"] = _resolve_fill_fact(runtime["color_source"], theme, color_map)
        runtime["color"] = _runtime_color(runtime["color_source"])
    runtime["theme_context"] = copy.deepcopy(theme_context)
    slot["theme_context"] = copy.deepcopy(theme_context)
    slot["style"] = runtime


def _enrich_layout_slots(
    slots: list[dict[str, Any]],
    *,
    master_text_styles: dict[str, list[dict[str, Any]]],
    canvas: dict[str, Any],
    theme: dict[str, Any] | None,
    color_map: dict[str, str],
    theme_context: dict[str, Any],
) -> None:
    for slot in slots:
        bucket = _slot_style_bucket(slot)
        _enrich_slot_style(
            slot,
            canvas=canvas,
            theme=theme,
            color_map=color_map,
            theme_context=theme_context,
            fallback_styles=master_text_styles.get(bucket, []),
            fallback_source=f"master.txStyles.{bucket}",
        )


def _enrich_slide_slots(
    slots: list[dict[str, Any]],
    *,
    layout_slots: list[dict[str, Any]],
    canvas: dict[str, Any],
    theme: dict[str, Any] | None,
    color_map: dict[str, str],
    theme_context: dict[str, Any],
) -> None:
    layout_by_idx = {slot.get("placeholder_index"): slot for slot in layout_slots}
    for slot in slots:
        fallback = layout_by_idx.get(slot.get("placeholder_index"))
        _enrich_slot_style(
            slot,
            canvas=canvas,
            theme=theme,
            color_map=color_map,
            theme_context=theme_context,
            fallback_styles=(fallback or {}).get("resolved_text_styles", []),
            fallback_source=(f"layout.placeholder.{slot.get('placeholder_index')}"
                             if fallback else "unresolved"),
        )


def _sample_elements(
    slide: Any,
    *,
    owner_id: str,
    canvas: dict[str, Any],
    asset_by_sha: dict[str, str],
    asset_path_by_id: dict[str, str],
) -> list[dict[str, Any]]:
    records = []
    for z_index, shape in enumerate(slide.shapes):
        try:
            if shape.is_placeholder:
                continue
        except Exception:
            pass
        fact = _shape_fact(
            shape,
            layer="slide",
            owner_id=owner_id,
            z_index=z_index,
            canvas=canvas,
            asset_by_sha=asset_by_sha,
            asset_path_by_id=asset_path_by_id,
        )
        shape_type = fact["shape_type"]
        has_text = bool((fact.get("text") or "").strip())
        replaceable = has_text or shape_type in {"picture", "chart", "table", "media"}
        fact["replaceable_candidate"] = replaceable
        fact["candidate_confidence"] = 0.65 if replaceable else 0.25
        fact["approval_required"] = True
        records.append(fact)
    return records


def _role_scores(name: str, placeholder_types: set[str], text: str = "") -> list[tuple[str, float, str]]:
    haystack = f"{name} {text}".lower()
    scores: dict[str, tuple[float, str]] = {}

    def add(role: str, score: float, reason: str) -> None:
        if role not in scores or score > scores[role][0]:
            scores[role] = (score, reason)

    keyword_sets = {
        "cover": ("title slide", "cover", "封面", "首页"),
        "raw": ("title and content", "content", "正文", "内容", "body"),
        "section": ("section header", "section", "章节", "过渡"),
        "quote": ("quote", "quotation", "金句", "引语"),
        "agenda": ("agenda", "table of contents", "contents", "目录"),
        "end": ("closing", "thank you", "thanks", "end", "封底", "谢谢", "感谢"),
    }
    for role, keywords in keyword_sets.items():
        if any(keyword in haystack for keyword in keywords):
            add(role, 0.94, f"name/text matched {role} keyword")
    if {"title", "subtitle"}.issubset(placeholder_types) or {"center_title", "subtitle"}.issubset(placeholder_types):
        add("cover", 0.82, "title and subtitle placeholders")
    if placeholder_types.intersection({"body", "object", "picture", "table", "chart"}):
        add("raw", 0.80, "content-bearing placeholder")
    if "title" in placeholder_types and not placeholder_types.intersection({"body", "object"}):
        add("section", 0.58, "title-led sparse layout")
    return [(role, score, reason) for role, (score, reason) in scores.items()]


def _role_candidates(layouts: list[dict[str, Any]], slides: list[dict[str, Any]]) -> dict[str, list[dict[str, Any]]]:
    candidates: dict[str, list[dict[str, Any]]] = {role: [] for role in ROLES}
    for layout in layouts:
        placeholder_types = {item["placeholder_type"] for item in layout["placeholders"]}
        for role, confidence, reason in _role_scores(layout.get("name") or "", placeholder_types):
            candidates[role].append({
                "source_kind": "layout",
                "selector": f"layout:{layout['index']}",
                "layout_id": layout["id"],
                "name": layout.get("name"),
                "confidence": confidence,
                "reason": reason,
            })
    for slide in slides:
        text = " ".join(filter(None, [slide.get("title"), slide.get("text_sample")]))
        placeholder_types = {item["placeholder_type"] for item in slide["placeholders"]}
        for role, confidence, reason in _role_scores("", placeholder_types, text):
            if confidence < 0.9:
                continue
            candidates[role].append({
                "source_kind": "slide",
                "selector": f"slide:{slide['index']}",
                "slide_id": slide["id"],
                "layout_id": slide["layout_id"],
                "name": slide.get("title"),
                "confidence": confidence,
                "reason": reason,
            })
    for role in ROLES:
        candidates[role].sort(key=lambda item: (-item["confidence"], item["selector"]))
        candidates[role] = candidates[role][:8]
    return candidates


def build_dossier(pptx_path: Path, output_dir: Path) -> dict[str, Any]:
    """Read source facts and extract embedded media; never render a page image."""
    if not pptx_path.is_file() or pptx_path.suffix.lower() != ".pptx":
        raise ExtractionError(f"PPTX not found or unsupported: {pptx_path}")
    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        raise ExtractionError(f"cannot open PPTX: {exc}") from exc

    output_dir.mkdir(parents=True, exist_ok=True)
    assets, asset_by_sha = _extract_assets(pptx_path, output_dir)
    asset_path_by_id = {asset["id"]: asset["path"] for asset in assets}
    canvas = _canvas_facts(prs)
    themes = _theme_facts(pptx_path)

    masters: list[dict[str, Any]] = []
    layouts: list[dict[str, Any]] = []
    layout_by_part: dict[str, str] = {}
    layout_counter = 0

    for master_index, master in enumerate(prs.slide_masters, 1):
        master_id = f"master-{master_index}"
        master_part = _part_name(master)
        master_theme_part = _master_theme_part(master)
        master_theme = _theme_record(themes, master_theme_part)
        master_color_map = _color_map(master)
        master_theme_context = _theme_context(master_id, master_theme, master_color_map)
        master_background = _background_fact(master, asset_by_sha, asset_path_by_id)
        master_fixed = _container_fixed_elements(
            master,
            layer="master",
            owner_id=master_id,
            canvas=canvas,
            asset_by_sha=asset_by_sha,
            asset_path_by_id=asset_path_by_id,
        )
        master_background_element = _background_element(
            master_background,
            owner_id=master_id,
            layer="master",
            canvas=canvas,
        )
        if master_background_element:
            master_fixed.insert(0, master_background_element)
        for element in master_fixed:
            _apply_theme_to_element(
                element,
                context=master_theme_context,
                theme=master_theme,
                color_map=master_color_map,
            )
        master_text_styles = _master_text_styles(master)
        master_placeholders = _container_placeholders(
            master,
            layer="master",
            owner_id=master_id,
            canvas=canvas,
        )
        _enrich_layout_slots(
            master_placeholders,
            master_text_styles=master_text_styles,
            canvas=canvas,
            theme=master_theme,
            color_map=master_color_map,
            theme_context=master_theme_context,
        )
        master_record = {
            "id": master_id,
            "index": master_index,
            "name": getattr(master, "name", None),
            "part": master_part,
            "theme_part": master_theme_part,
            "color_map": master_color_map,
            "theme_context": master_theme_context,
            "background": master_background,
            "fixed_elements": master_fixed,
            "placeholders": master_placeholders,
            "text_styles": master_text_styles,
            "layout_ids": [],
        }
        masters.append(master_record)

        for layout in master.slide_layouts:
            layout_counter += 1
            layout_id = f"layout-{layout_counter}"
            layout_part = _part_name(layout)
            if layout_part:
                layout_by_part[layout_part] = layout_id
            layout_color_map = _color_map(layout, master_color_map)
            layout_theme_context = _theme_context(layout_id, master_theme, layout_color_map)
            layout_background = _background_fact(layout, asset_by_sha, asset_path_by_id)
            layout_fixed = _container_fixed_elements(
                layout,
                layer="layout",
                owner_id=layout_id,
                canvas=canvas,
                asset_by_sha=asset_by_sha,
                asset_path_by_id=asset_path_by_id,
            )
            layout_background_element = _background_element(
                layout_background,
                owner_id=layout_id,
                layer="layout",
                canvas=canvas,
            )
            if layout_background_element:
                layout_fixed.insert(0, layout_background_element)
            for element in layout_fixed:
                _apply_theme_to_element(
                    element,
                    context=layout_theme_context,
                    theme=master_theme,
                    color_map=layout_color_map,
                )
            layout_placeholders = _container_placeholders(
                layout,
                layer="layout",
                owner_id=layout_id,
                canvas=canvas,
            )
            _enrich_layout_slots(
                layout_placeholders,
                master_text_styles=master_text_styles,
                canvas=canvas,
                theme=master_theme,
                color_map=layout_color_map,
                theme_context=layout_theme_context,
            )
            record = {
                "id": layout_id,
                "index": layout_counter,
                "name": getattr(layout, "name", None),
                "part": layout_part,
                "master_id": master_id,
                "theme_part": master_theme_part,
                "color_map": layout_color_map,
                "theme_context": layout_theme_context,
                "background": layout_background,
                "fixed_elements": layout_fixed,
                "placeholders": layout_placeholders,
                "used_by_slides": [],
            }
            layouts.append(record)
            master_record["layout_ids"].append(layout_id)

    layouts_by_id = {record["id"]: record for record in layouts}
    slides: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(prs.slides, 1):
        part = _part_name(slide)
        layout_part = _part_name(slide.slide_layout)
        layout_id = layout_by_part.get(layout_part or "")
        if layout_id and layout_id in layouts_by_id:
            layouts_by_id[layout_id]["used_by_slides"].append(slide_index)
        title = ""
        try:
            if slide.shapes.title is not None:
                title = slide.shapes.title.text or ""
        except Exception:
            pass
        text_items = []
        for shape in slide.shapes:
            try:
                if shape.has_text_frame and shape.text_frame.text:
                    text_items.append(shape.text_frame.text)
            except Exception:
                continue
        slide_id = f"slide-{slide_index}"
        source_layout = layouts_by_id.get(layout_id or "", {})
        source_master = next(
            (item for item in masters if item["id"] == source_layout.get("master_id")),
            None,
        )
        slide_color_map = _color_map(slide, (source_layout or {}).get("color_map"))
        slide_theme = _theme_record(themes, (source_master or {}).get("theme_part"))
        slide_theme_context = _theme_context(slide_id, slide_theme, slide_color_map)
        slide_placeholders = _container_placeholders(
            slide,
            layer="slide",
            owner_id=slide_id,
            canvas=canvas,
        )
        _enrich_slide_slots(
            slide_placeholders,
            layout_slots=layouts_by_id.get(layout_id or "", {}).get("placeholders", []),
            canvas=canvas,
            theme=slide_theme,
            color_map=slide_color_map,
            theme_context=slide_theme_context,
        )
        slide_sample_elements = _sample_elements(
            slide,
            owner_id=slide_id,
            canvas=canvas,
            asset_by_sha=asset_by_sha,
            asset_path_by_id=asset_path_by_id,
        )
        for element in slide_sample_elements:
            _apply_theme_to_element(
                element,
                context=slide_theme_context,
                theme=slide_theme,
                color_map=slide_color_map,
            )
        slides.append({
            "id": slide_id,
            "index": slide_index,
            "slide_id": int(getattr(slide, "slide_id", slide_index)),
            "part": part,
            "layout_id": layout_id,
            "master_id": layouts_by_id.get(layout_id or "", {}).get("master_id"),
            "theme_part": (source_master or {}).get("theme_part"),
            "color_map": slide_color_map,
            "theme_context": slide_theme_context,
            "title": title,
            "text_sample": "\n".join(text_items)[:2000],
            "placeholders": slide_placeholders,
            "sample_elements": slide_sample_elements,
        })

    candidates = _role_candidates(layouts, slides)
    warnings = [
        "Semantic role candidates are recommendations only; no role is activated without an explicit --role, --derive, or --alias declaration.",
        "Master/layout non-placeholder elements are fixed-by-source candidates; confirm which are protected corporate VI before approving the pack.",
        "Font facts are preserved from OOXML, but local font availability is not asserted by this extractor.",
    ]
    if not any(layout["placeholders"] for layout in layouts):
        warnings.append("No layout placeholders were found; safe areas and replaceable slots require manual confirmation.")
    if not themes:
        warnings.append("No readable PowerPoint theme part was found.")

    return {
        "schema_version": SCHEMA_VERSION,
        "source": {
            "path": str(pptx_path.resolve()),
            "filename": pptx_path.name,
            "sha256": _sha256_file(pptx_path),
            "slide_count": len(slides),
            "master_count": len(masters),
            "layout_count": len(layouts),
            "extracted_at": _now(),
        },
        "canvas": canvas,
        "theme": {"themes": themes},
        "masters": masters,
        "layouts": layouts,
        "slides": slides,
        "assets": assets,
        "role_candidates": candidates,
        "confidence": {
            "overall": 0.82 if layouts and themes else 0.65,
            "needs_confirmation": [
                "semantic role mappings",
                "locked corporate VI elements",
                "font availability and embedding rights",
                "derived layout safe areas",
            ],
        },
        "warnings": warnings,
        "untrusted": True,
    }


def _bindings(values: dict[str, str] | Iterable[str] | None, option: str) -> dict[str, str]:
    if values is None:
        return {}
    if isinstance(values, dict):
        pairs = list(values.items())
    else:
        pairs = []
        for value in values:
            if "=" not in value:
                raise ExtractionError(f"{option} expects ROLE=VALUE, got: {value}")
            pairs.append(tuple(value.split("=", 1)))
    result: dict[str, str] = {}
    for raw_role, raw_value in pairs:
        role, value = raw_role.strip().lower(), raw_value.strip()
        if role not in ROLES:
            raise ExtractionError(f"unknown semantic role for {option}: {role}; expected one of {', '.join(ROLES)}")
        if not value:
            raise ExtractionError(f"empty value for {option} {role}")
        if role in result:
            raise ExtractionError(f"duplicate {option} declaration for role: {role}")
        result[role] = value
    return result


def _resolve_selector(dossier: dict[str, Any], selector: str) -> tuple[dict[str, Any], dict[str, Any] | None]:
    layouts = dossier["layouts"]
    slides = dossier["slides"]
    if selector.startswith("slide:"):
        token = selector.split(":", 1)[1].strip()
        if not token.isdigit() or int(token) < 1:
            raise ExtractionError(f"invalid slide selector: {selector}; use slide:N with N >= 1")
        matches = [item for item in slides if item["index"] == int(token)]
        if not matches:
            raise ExtractionError(f"slide selector out of range: {selector}")
        slide = matches[0]
        layout = next((item for item in layouts if item["id"] == slide.get("layout_id")), None)
        if layout is None:
            raise ExtractionError(f"slide {token} has no readable source layout")
        return layout, slide
    if selector.startswith("layout-name:"):
        token = selector.split(":", 1)[1].strip()
        matches = [item for item in layouts if (item.get("name") or "") == token]
    elif selector.startswith("layout:"):
        token = selector.split(":", 1)[1].strip()
        if token.isdigit():
            matches = [item for item in layouts if item["index"] == int(token)]
        else:
            matches = [item for item in layouts if (item.get("name") or "") == token]
    else:
        raise ExtractionError(
            f"invalid role selector: {selector}; use slide:N, layout:N, or layout-name:NAME"
        )
    if not matches:
        raise ExtractionError(f"layout selector did not match: {selector}")
    if len(matches) > 1:
        raise ExtractionError(f"layout selector is ambiguous: {selector}; use the 1-based layout index")
    return matches[0], None


def _slot_from_sample(element: dict[str, Any]) -> dict[str, Any]:
    shape_type = element.get("shape_type")
    if (element.get("text") or "").strip():
        kind = "freeform_text"
    elif shape_type in {"picture", "media"}:
        kind = "freeform_image"
    else:
        kind = "freeform_content"
    slot = {
        "id": f"candidate-{element['id']}",
        "placeholder_index": None,
        "placeholder_type": "freeform",
        "slot_kind": kind,
        "semantic_name": kind,
        "system_field": False,
        "name": element.get("name"),
        "prompt_text": element.get("text") or "",
        "geometry_emu": copy.deepcopy(element["geometry_emu"]),
        "geometry": copy.deepcopy(element["geometry"]),
        "text_styles": copy.deepcopy(element.get("style", {}).get("text_styles", [])),
        "style": copy.deepcopy(element.get("style", {}).get("text", {})),
        "theme_context": copy.deepcopy(element.get("theme_context", {})),
        "source": {
            "layer": "slide",
            "part": element.get("source_part"),
            "shape_id": element.get("shape_id"),
        },
        "confidence": element.get("candidate_confidence", 0.65),
        "approval_required": True,
    }
    if element.get("asset_id"):
        slot["asset_id"] = element["asset_id"]
        slot["src"] = element.get("src")
    return slot


def _safe_area(slots: list[dict[str, Any]], canvas: dict[str, Any]) -> dict[str, Any]:
    content_slots = [slot for slot in slots if not slot.get("system_field")]
    if not content_slots:
        return {
            "status": "unknown",
            "geometry_emu": None,
            "geometry": None,
            "source": "no confirmed content slots",
            "confidence": 0.0,
        }
    xs = [slot["geometry_emu"]["x"] for slot in content_slots]
    ys = [slot["geometry_emu"]["y"] for slot in content_slots]
    rights = [slot["geometry_emu"]["x"] + slot["geometry_emu"]["w"] for slot in content_slots]
    bottoms = [slot["geometry_emu"]["y"] + slot["geometry_emu"]["h"] for slot in content_slots]
    x, y, right, bottom = min(xs), min(ys), max(rights), max(bottoms)
    geometry_emu = {"x": x, "y": y, "w": right - x, "h": bottom - y}
    geometry_for_projection = {
        **geometry_emu,
        "normalized": {
            "x": x / canvas["source_width_emu"],
            "y": y / canvas["source_height_emu"],
            "w": (right - x) / canvas["source_width_emu"],
            "h": (bottom - y) / canvas["source_height_emu"],
        },
    }
    return {
        "status": "candidate",
        "geometry_emu": geometry_emu,
        "geometry": _design_geometry(geometry_for_projection, canvas),
        "source": "union of non-system slot candidates",
        "confidence": min(slot.get("confidence", 0.5) for slot in content_slots),
    }


def _apply_role_semantics(slots: list[dict[str, Any]], role: str) -> list[str]:
    """Add draft semantic names without erasing factual placeholder kinds."""
    needs_confirmation: list[str] = []
    content_slots = [slot for slot in slots if not slot.get("system_field")]
    for slot in slots:
        label = f"{slot.get('name') or ''} {slot.get('prompt_text') or ''}".lower()
        semantic = slot.get("slot_kind", "content")
        confidence = 1.0
        keyword_semantics = (
            (("author", "presenter", "汇报人", "姓名"), "author"),
            (("date", "日期", "时间"), "date"),
            (("chapter", "section", "章节"), "chapter"),
            (("quote", "quotation", "金句", "引语"), "quote"),
            (("attribution", "source", "出处", "来源"), "attribution"),
        )
        for keywords, candidate in keyword_semantics:
            if any(keyword in label for keyword in keywords):
                semantic = candidate
                confidence = 0.9
                break
        slot["semantic_name"] = semantic
        slot["semantic_name_confidence"] = confidence

    # The caller explicitly mapped this source to the role, so these are
    # transparent draft hints rather than hidden automatic role activation.
    if role == "cover":
        freeform = [
            slot for slot in content_slots
            if slot.get("semantic_name") == "freeform_text"
        ]
        freeform.sort(key=lambda slot: (
            -float(slot.get("style", {}).get("font_size") or 0),
            float(slot.get("geometry", {}).get("y") or 0),
        ))
        if not any(slot.get("semantic_name") == "title" for slot in content_slots) and freeform:
            title = freeform.pop(0)
            title["semantic_name"] = "title"
            title["semantic_name_confidence"] = 0.72
            title["semantic_mapping_method"] = "cover-freeform-largest-text"
            title["approval_required"] = True
            needs_confirmation.append(
                f"confirm inferred cover title slot {title['id']} from a freeform text box"
            )
        if not any(slot.get("semantic_name") == "subtitle" for slot in content_slots) and freeform:
            subtitle = freeform.pop(0)
            subtitle["semantic_name"] = "subtitle"
            subtitle["semantic_name_confidence"] = 0.62
            subtitle["semantic_mapping_method"] = "cover-freeform-secondary-text"
            subtitle["approval_required"] = True
            needs_confirmation.append(
                f"confirm inferred cover subtitle slot {subtitle['id']} from a freeform text box"
            )
        if freeform:
            needs_confirmation.append(
                f"{len(freeform)} additional cover freeform text slot(s) remain unmapped"
            )
        if not any(slot.get("semantic_name") == "title" for slot in content_slots):
            needs_confirmation.append("cover has no confirmed title slot and cannot be approved")
    elif role == "section":
        title = next((slot for slot in content_slots if slot.get("slot_kind") == "title"), None)
        if title:
            # A PowerPoint section layout's title placeholder is the section
            # heading itself.  Do not silently reinterpret it as the optional
            # numeric chapter marker used by Feishu's default layout; a real
            # chapter-number placeholder is detected by the keyword pass above.
            title["semantic_name"] = "title"
            title["semantic_name_confidence"] = 0.95
            title["approval_required"] = True
    elif role == "quote" and content_slots:
        quote_slot = next(
            (slot for slot in content_slots if slot.get("semantic_name") == "quote"),
            content_slots[0],
        )
        quote_slot["semantic_name"] = "quote"
        quote_slot["semantic_name_confidence"] = max(
            float(quote_slot.get("semantic_name_confidence", 0)), 0.75
        )
        quote_slot["approval_required"] = True
        remainder = [slot for slot in content_slots if slot is not quote_slot]
        if remainder and not any(slot.get("semantic_name") == "attribution" for slot in remainder):
            remainder[-1]["semantic_name"] = "attribution"
            remainder[-1]["semantic_name_confidence"] = 0.6
            remainder[-1]["approval_required"] = True
            needs_confirmation.append(
                f"confirm inferred quote attribution slot {remainder[-1]['id']}"
            )
    return needs_confirmation


def _resolved_background(
    background: dict[str, Any] | None,
    *,
    theme: dict[str, Any] | None,
    color_map: dict[str, str],
    theme_context: dict[str, Any],
) -> dict[str, Any] | None:
    if not background:
        return None
    result = copy.deepcopy(background)
    result["theme_context"] = copy.deepcopy(theme_context)
    if result.get("fill") is not None:
        result["source_fill"] = copy.deepcopy(result["fill"])
        result["fill"] = _resolve_fill_fact(result["fill"], theme, color_map)
    reference = result.get("theme_reference")
    if isinstance(reference, dict) and reference.get("color"):
        result["source_theme_reference"] = copy.deepcopy(reference)
        resolved = _resolve_color_fact(reference.get("color"), theme, color_map)
        result["theme_reference"] = {**reference, "color": resolved}
        if result.get("fill") is None and resolved:
            result["fill"] = {"type": "solid", "color": resolved}
    result["color"] = _runtime_color(result.get("fill"))
    return result


def _native_layout(
    dossier: dict[str, Any],
    *,
    role: str,
    selector: str,
) -> dict[str, Any]:
    layout, slide = _resolve_selector(dossier, selector)
    master = next(item for item in dossier["masters"] if item["id"] == layout["master_id"])
    context_source = slide or layout
    theme = _theme_record(dossier["theme"]["themes"], context_source.get("theme_part"))
    color_map = context_source.get("color_map", master.get("color_map", _DEFAULT_CLR_MAP))
    theme_context = context_source.get("theme_context") or _theme_context(
        context_source["id"], theme, color_map,
    )
    fixed = copy.deepcopy(master["fixed_elements"] + layout["fixed_elements"])

    # Layout placeholders are the default.  A mapped sample slide may carry
    # explicit placeholder geometry/style overrides; replace by idx when found.
    slots_by_idx = {slot["placeholder_index"]: copy.deepcopy(slot) for slot in layout["placeholders"]}
    freeform: list[dict[str, Any]] = []
    if slide is not None:
        for slot in slide["placeholders"]:
            slots_by_idx[slot["placeholder_index"]] = copy.deepcopy(slot)
        freeform = [
            _slot_from_sample(element)
            for element in slide["sample_elements"]
            if element.get("replaceable_candidate")
        ]
    slots = list(slots_by_idx.values()) + freeform
    semantic_review = _apply_role_semantics(slots, role)
    for element in fixed:
        if element.get("source_stack_plane") == "interleaved":
            semantic_review.append(
                f"confirm interleaved fixed element {element['id']} relative to source placeholders"
            )
    return {
        "semantic_role": role,
        "fixed_elements": fixed,
        "slots": slots,
        "safe_area": _safe_area(slots, dossier["canvas"]),
        "background": _resolved_background(
            layout.get("background") or master.get("background"),
            theme=theme,
            color_map=color_map,
            theme_context=theme_context,
        ),
        "theme_context": copy.deepcopy(theme_context),
        "needs_confirmation": semantic_review,
        "source": {
            "kind": "slide" if slide is not None else "layout",
            "selector": selector,
            "layout_id": layout["id"],
            "layout_index": layout["index"],
            "layout_name": layout.get("name"),
            "master_id": layout["master_id"],
            "slide_id": slide["id"] if slide is not None else None,
            "slide_index": slide["index"] if slide is not None else None,
        },
        "status": "native",
        "confidence": 1.0,
    }


def _theme_tokens(dossier: dict[str, Any]) -> dict[str, Any]:
    colors = []
    fonts = []
    for theme in dossier["theme"]["themes"]:
        for name, value in theme.get("colors", {}).items():
            colors.append({"token": name, "value": value, "source_theme": theme.get("part")})
        for family, branch in theme.get("fonts", {}).items():
            for script, typeface in branch.items():
                if script == "scripts":
                    for script_name, script_typeface in typeface.items():
                        fonts.append({
                            "token": f"{family}.script.{script_name}",
                            "typeface": script_typeface,
                            "source_theme": theme.get("part"),
                        })
                else:
                    fonts.append({
                        "token": f"{family}.{script}",
                        "typeface": typeface,
                        "source_theme": theme.get("part"),
                    })
    contexts = []
    themes_by_part = {
        theme.get("part"): theme for theme in dossier["theme"]["themes"]
    }
    for master in dossier["masters"]:
        context_id = master["id"]
        theme = themes_by_part.get(master.get("theme_part"))
        color_map = master.get("color_map", {})
        contexts.append(copy.deepcopy(master.get("theme_context", {})))
        for alias, target in color_map.items():
            resolved = _resolve_color_fact(
                {"kind": "theme", "value": alias, "transforms": []},
                theme,
                color_map,
            )
            colors.append({
                "token": f"{context_id}.{alias}",
                "value": resolved,
                "source_theme": master.get("theme_part"),
                "source_alias": alias,
                "resolved_token": target,
            })
            if len(dossier["masters"]) == 1:
                colors.append({
                    "token": alias,
                    "value": copy.deepcopy(resolved),
                    "source_theme": master.get("theme_part"),
                    "source_alias": alias,
                    "resolved_token": target,
                })
    typography = []
    seen: set[str] = set()
    for layout in dossier["layouts"]:
        for slot in layout["placeholders"]:
            for style in slot.get("text_styles", []):
                key = json.dumps(style, ensure_ascii=False, sort_keys=True)
                if key not in seen:
                    seen.add(key)
                    typography.append(style)
    return {
        "colors": colors,
        "fonts": fonts,
        "typography": typography,
        "theme_contexts": contexts,
        "source_policy": "exact-ooxml-facts-no-normalization",
    }


def _walk_elements(elements: Iterable[dict[str, Any]]) -> Iterable[dict[str, Any]]:
    for element in elements:
        yield element
        yield from _walk_elements(element.get("children", []))


def build_template_pack(
    dossier: dict[str, Any],
    *,
    template_id: str,
    version: str = PACK_VERSION,
    role_mappings: dict[str, str] | Iterable[str] | None = None,
    aliases: dict[str, str] | Iterable[str] | None = None,
    derives: dict[str, str] | Iterable[str] | None = None,
) -> dict[str, Any]:
    mappings = _bindings(role_mappings, "--role")
    alias_map = _bindings(aliases, "--alias")
    derive_map = _bindings(derives, "--derive")
    for role in ROLES:
        modes = sum(role in item for item in (mappings, alias_map, derive_map))
        if modes > 1:
            raise ExtractionError(f"role {role} has conflicting --role/--alias/--derive declarations")
    for option, mapping in (("--alias", alias_map), ("--derive", derive_map)):
        for role, target in mapping.items():
            target = target.lower()
            if target not in ROLES:
                raise ExtractionError(f"{option} {role} targets unknown role: {target}")
            mapping[role] = target

    layouts: dict[str, dict[str, Any]] = {}
    coverage: dict[str, dict[str, Any]] = {}
    resolving: list[str] = []

    def resolve(role: str) -> str | None:
        if role in coverage:
            return coverage[role].get("layout_id") or coverage[role].get("source", {}).get("resolved_layout_id")
        if role in resolving:
            cycle = resolving[resolving.index(role):] + [role]
            raise ExtractionError(f"layout alias/derive cycle: {' -> '.join(cycle)}")
        resolving.append(role)
        try:
            if role in mappings:
                layout_id = f"role-{role}"
                layouts[layout_id] = _native_layout(dossier, role=role, selector=mappings[role])
                coverage[role] = {
                    "status": "native",
                    "layout_id": layout_id,
                    "source": copy.deepcopy(layouts[layout_id]["source"]),
                    "confidence": 1.0,
                }
                return layout_id
            if role in derive_map:
                target = derive_map[role]
                source_layout_id = resolve(target)
                if not source_layout_id:
                    raise ExtractionError(f"--derive {role}={target} cannot use unsupported role {target}")
                layout_id = f"role-{role}"
                derived = copy.deepcopy(layouts[source_layout_id])
                derived["semantic_role"] = role
                derived["status"] = "derived"
                derived["confidence"] = min(float(derived.get("confidence", 1.0)), 0.55)
                derived["source"] = {
                    "kind": "derived",
                    "derived_from_role": target,
                    "derived_from_layout_id": source_layout_id,
                    "approval_required": True,
                }
                for slot in derived["slots"]:
                    slot["approval_required"] = True
                    slot["confidence"] = min(float(slot.get("confidence", 1.0)), 0.55)
                semantic_review = _apply_role_semantics(derived["slots"], role)
                derived["needs_confirmation"] = list(dict.fromkeys(
                    list(derived.get("needs_confirmation") or []) + semantic_review
                ))
                derived["safe_area"] = _safe_area(derived["slots"], dossier["canvas"])
                layouts[layout_id] = derived
                coverage[role] = {
                    "status": "derived",
                    "layout_id": layout_id,
                    "derived_from": target,
                    "source": copy.deepcopy(derived["source"]),
                    "confidence": derived["confidence"],
                }
                return layout_id
            if role in alias_map:
                target = alias_map[role]
                source_layout_id = resolve(target)
                if not source_layout_id:
                    raise ExtractionError(f"--alias {role}={target} cannot target unsupported role {target}")
                coverage[role] = {
                    "status": "alias",
                    "alias_to": target,
                    "source": {
                        "kind": "alias",
                        "declared_target": target,
                        "resolved_layout_id": source_layout_id,
                    },
                    "confidence": coverage[target].get("confidence", 1.0),
                }
                return source_layout_id
            coverage[role] = {
                "status": "unsupported",
                "source": {"kind": "missing", "reason": "no explicit approved mapping"},
                "confidence": 1.0,
            }
            return None
        finally:
            resolving.pop()

    for semantic_role in ROLES:
        resolve(semantic_role)

    brand_refs = []
    asset_ids: set[str] = set()
    for layout_id, layout in layouts.items():
        for element in _walk_elements(layout["fixed_elements"]):
            if element.get("asset_id"):
                asset_ids.add(element["asset_id"])
                brand_refs.append({
                    "layout_id": layout_id,
                    "element_id": element["id"],
                    "asset_id": element["asset_id"],
                    "status": "candidate",
                    "confidence": 0.6,
                    "requires_confirmation": True,
                })

    missing_roles = [role for role in ROLES if coverage[role]["status"] == "unsupported"]
    derived_roles = [role for role in ROLES if coverage[role]["status"] == "derived"]
    layout_review_items = [
        f"{layout_id}: {item}"
        for layout_id, layout in layouts.items()
        for item in layout.get("needs_confirmation", [])
    ]
    return {
        "schema_version": SCHEMA_VERSION,
        "template_id": template_id,
        "version": version,
        "status": "draft",
        "source": copy.deepcopy(dossier["source"]),
        "canvas": copy.deepcopy(dossier["canvas"]),
        "tokens": _theme_tokens(dossier),
        "brand": {
            "lock_status": "pending_confirmation",
            "fixed_element_candidates": brand_refs,
            "asset_ids": sorted(asset_ids),
            "policy": "preserve source geometry and z-order after approval",
        },
        "layouts": layouts,
        "layout_coverage": {role: coverage[role] for role in ROLES},
        "policies": {
            "typography_mode": "strict",
            "allow_automatic_font_size_change": False,
            "allow_automatic_font_substitution": False,
            "overflow_resolution_order": [
                "shorten-copy",
                "switch-approved-layout",
                "split-slide",
                "request-confirmation",
            ],
            "preserve_fixed_elements": True,
            "preserve_logo_alignment": True,
            "missing_layout_behavior": "block",
            "whole_page_screenshots": "forbidden",
            "approval_required_before_activation": True,
        },
        "extraction_report": {
            "generated_at": _now(),
            "role_candidates": copy.deepcopy(dossier["role_candidates"]),
            "missing_roles": missing_roles,
            "derived_roles": derived_roles,
            "warnings": copy.deepcopy(dossier["warnings"]),
            "needs_confirmation": list(dict.fromkeys(
                copy.deepcopy(dossier["confidence"]["needs_confirmation"])
                + layout_review_items
            )),
        },
    }


def _preview_html(dossier: dict[str, Any], pack: dict[str, Any]) -> str:
    canvas = dossier["canvas"]
    rows = []
    for role in ROLES:
        item = pack["layout_coverage"][role]
        layout = pack["layouts"].get(item.get("layout_id", ""), {})
        source = item.get("source", {})
        source_label = (
            source.get("selector")
            or (f"alias → {item.get('alias_to')}" if item["status"] == "alias" else "—")
        )
        rows.append(
            "<tr>"
            f"<td><b>{html.escape(ROLE_LABELS[role])}</b><small>{role}</small></td>"
            f"<td><span class='status {item['status']}'>{html.escape(item['status'])}</span></td>"
            f"<td>{html.escape(str(source_label))}</td>"
            f"<td>{float(item.get('confidence', 0)):.2f}</td>"
            f"<td>{len(layout.get('fixed_elements', []))}</td>"
            f"<td>{len(layout.get('slots', []))}</td>"
            "</tr>"
        )

    candidate_rows = []
    for role in ROLES:
        candidates = dossier["role_candidates"].get(role, [])
        if not candidates:
            candidate_rows.append(
                f"<tr><td>{html.escape(ROLE_LABELS[role])}</td><td colspan='4'>没有可靠候选</td></tr>"
            )
            continue
        for index, candidate in enumerate(candidates[:3]):
            candidate_rows.append(
                "<tr>"
                f"<td>{html.escape(ROLE_LABELS[role]) if index == 0 else ''}</td>"
                f"<td>{html.escape(candidate['selector'])}</td>"
                f"<td>{html.escape(str(candidate.get('name') or '—'))}</td>"
                f"<td>{float(candidate['confidence']):.2f}</td>"
                f"<td>{html.escape(candidate['reason'])}</td>"
                "</tr>"
            )

    warnings = "".join(f"<li>{html.escape(item)}</li>" for item in pack["extraction_report"]["warnings"])
    return f"""<!doctype html>
<html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{html.escape(pack['template_id'])} · Template Review</title>
<style>
body{{margin:0;background:#f5f6f8;color:#1f2329;font:15px/1.55 -apple-system,BlinkMacSystemFont,"Segoe UI","PingFang SC",sans-serif}}
main{{max-width:1120px;margin:0 auto;padding:40px 24px 72px}}h1{{font-size:32px;margin:0 0 8px}}h2{{margin:34px 0 12px}}
.muted,small{{color:#8f959e}}.summary{{display:grid;grid-template-columns:repeat(4,1fr);gap:12px;margin:24px 0}}
.card{{background:#fff;border:1px solid #e5e6eb;border-radius:12px;padding:18px}}.card b{{display:block;font-size:22px}}
table{{width:100%;border-collapse:collapse;background:#fff;border-radius:12px;overflow:hidden}}th,td{{padding:12px 14px;border-bottom:1px solid #eff0f1;text-align:left;vertical-align:top}}th{{background:#f0f4ff}}td small{{display:block}}
.status{{padding:3px 9px;border-radius:999px;font-weight:600}}.native{{background:#dcfce7;color:#166534}}.derived{{background:#fef3c7;color:#92400e}}.alias{{background:#dbeafe;color:#1d4ed8}}.unsupported{{background:#f3f4f6;color:#6b7280}}
.notice{{background:#fff7e6;border:1px solid #ffd591;border-radius:12px;padding:14px 18px}}code{{background:#eef1f5;padding:2px 5px;border-radius:4px}}
</style></head><body><main>
<p class="muted">PPTX Template Design System · draft review</p><h1>{html.escape(pack['template_id'])}</h1>
<p>此报告只展示结构化事实和候选映射，没有整页截图，也没有自动激活缺失版式。</p>
<section class="summary">
<div class="card"><span>源尺寸</span><b>{canvas['width_inches']} × {canvas['height_inches']} in</b></div>
<div class="card"><span>页面比例</span><b>{html.escape(canvas['aspect_ratio']['label'])}</b></div>
<div class="card"><span>建议 H5 画布</span><b>{canvas['recommended_design_canvas']['width']} × {canvas['recommended_design_canvas']['height']}</b></div>
<div class="card"><span>源结构</span><b>{dossier['source']['master_count']} M / {dossier['source']['layout_count']} L / {dossier['source']['slide_count']} P</b></div>
</section>
<h2>六类 Layout 覆盖</h2><table><thead><tr><th>角色</th><th>状态</th><th>来源</th><th>置信度</th><th>固定元素</th><th>插槽</th></tr></thead><tbody>{''.join(rows)}</tbody></table>
<h2>候选页面 / Layout</h2><table><thead><tr><th>角色</th><th>选择器</th><th>名称</th><th>置信度</th><th>依据</th></tr></thead><tbody>{''.join(candidate_rows)}</tbody></table>
<h2>确认前注意</h2><div class="notice"><ul>{warnings}</ul><p>确认后再把 <code>status</code> 从 <code>draft</code> 提升为批准版本。</p></div>
</main></body></html>"""


def _atomic_json(path: Path, payload: dict[str, Any]) -> None:
    temporary = path.with_suffix(path.suffix + ".tmp")
    temporary.write_text(json.dumps(payload, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    os.replace(temporary, path)


def extract_template(
    pptx_path: Path,
    output_dir: Path,
    *,
    template_id: str | None = None,
    version: str = PACK_VERSION,
    role_mappings: dict[str, str] | Iterable[str] | None = None,
    aliases: dict[str, str] | Iterable[str] | None = None,
    derives: dict[str, str] | Iterable[str] | None = None,
    force: bool = False,
) -> tuple[dict[str, Any], dict[str, Any], Path]:
    output_dir = output_dir.resolve()
    dossier_path = output_dir / "template-dossier.json"
    pack_path = output_dir / "template-pack.json"
    preview_path = output_dir / "template-preview.html"
    conflicts = [path for path in (dossier_path, pack_path, preview_path) if path.exists()]
    if conflicts and not force:
        raise ExtractionError(
            "refusing to overwrite existing template artifacts without --force: "
            + ", ".join(str(path) for path in conflicts)
        )
    dossier = build_dossier(pptx_path.resolve(), output_dir)
    pack = build_template_pack(
        dossier,
        template_id=template_id or _slug(pptx_path.stem),
        version=version,
        role_mappings=role_mappings,
        aliases=aliases,
        derives=derives,
    )
    _atomic_json(dossier_path, dossier)
    _atomic_json(pack_path, pack)
    temporary = preview_path.with_suffix(".html.tmp")
    temporary.write_text(_preview_html(dossier, pack), encoding="utf-8")
    os.replace(temporary, preview_path)
    return dossier, pack, preview_path


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("output_dir", type=Path)
    parser.add_argument("--template-id", default=None)
    parser.add_argument("--version", default=PACK_VERSION)
    parser.add_argument("--role", action="append", default=[], metavar="ROLE=SELECTOR")
    parser.add_argument("--alias", action="append", default=[], metavar="ROLE=ROLE")
    parser.add_argument("--derive", action="append", default=[], metavar="ROLE=ROLE")
    parser.add_argument("--force", action="store_true")
    args = parser.parse_args(argv)
    try:
        dossier, pack, preview = extract_template(
            args.pptx,
            args.output_dir,
            template_id=args.template_id,
            version=args.version,
            role_mappings=args.role,
            aliases=args.alias,
            derives=args.derive,
            force=args.force,
        )
    except ExtractionError as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2
    output_dir = args.output_dir.resolve()
    statuses = ", ".join(
        f"{role}={pack['layout_coverage'][role]['status']}" for role in ROLES
    )
    print(f"wrote {output_dir / 'template-dossier.json'}")
    print(f"wrote {output_dir / 'template-pack.json'}")
    print(f"wrote {preview}")
    print(f"canvas {dossier['canvas']['aspect_ratio']['label']} -> "
          f"{dossier['canvas']['recommended_design_canvas']['width']}x"
          f"{dossier['canvas']['recommended_design_canvas']['height']}")
    print(f"coverage: {statuses}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
