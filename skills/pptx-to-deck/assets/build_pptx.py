#!/usr/bin/env python3
"""
build_pptx.py · PPTX → feishu-deck-h5 deck.json (layout:"canvas")

Native PowerPoint (.pptx) parser backend. Walks every slide with python-pptx
and emits each slide as a `layout:"canvas"` deck.json slide: a STRUCTURED list
of absolutely-positioned, typed, id'd elements (`data.elements[]`) — NOT an HTML
blob and NOT a screenshot. The user's own feishu-deck-h5 renderer
(deck-json/render-deck.py) turns elements[] into positioned HTML, and
sync-index-to-deck.py round-trips edits back into elements[] by id.

This is the PPTX side of the DECKJSON-UNIFIED-INTERMEDIATE-SPEC (§2/§3): one
intermediate layer = deck.json, structured JSON for PPTX (no HTML坨, no images),
one edit loop = edit→sync→deck.json.

.pptx is open OOXML, so python-pptx reads exact element geometry (EMU), text
runs (font / size / color / b / i), embedded images, and groups directly,
cross-platform — no AppleScript, no Keynote, no LibreOffice.

Element mapping (each → {id, type, x, y, w, h, ...} ; px on canvas 1920×1080):
  · TEXT_BOX / PLACEHOLDER-with-text → {type:"text", runs:[{text,bold,color,
        size}], anchor, insets}   — CLEAN STRUCTURED editable content
  · PICTURE                         → {type:"image", src:"input/<file>"}
        (embedded media extracted to input/ — that's original content, kept)
  · MEDIA (video)                   → {type:"image", src:<poster>} or a poster
        placeholder shape
  · AUTO_SHAPE / FREEFORM / LINE    → {type:"shape", kind, fill|gradient,
        border, radius, svg, style} — appearance fields so it renders right
  · GROUP                           → FLATTENED: children composed through the
        group transform and emitted as top-level elements (no group wrapper)

Un-reconstructable (live chart / SmartArt diagram / OLE object): the WHOLE slide
becomes a placeholder — {layout:"canvas", data:{placeholder:true, source_page:N,
elements:[]}} — and N is collected. At the end a report line prints:
  `unreconstructed slides: [N, ...]`  (empty list if none).

NO SCREENSHOTS: the old whole-slide / per-element raster fallback is RETIRED.
--raster / --full-raster are accepted as deprecated no-ops (they never produce
image-replica slides). Embedded PICTUREs stay as image elements (original
content, fine). This is "尽可能还原" (best-effort reconstruction), NOT pixel-
perfect: gradients are carried as CSS, freeform/line as inline SVG paths.

Usage:
  build_pptx.py <in.pptx> <out-dir>
       [--renderer DIR]   feishu-deck-h5 skill root
                          (default: ~/.claude/skills/feishu-deck-h5)
       [--limit N]        only first N slides
       [--raster]         DEPRECATED no-op (rasterization retired)
       [--full-raster]    DEPRECATED no-op (rasterization retired)
       [--inline]         single-file render output (base64-inline everything)
       [--title TEXT]     deck title
       [--no-render]      emit deck.json + assets only, skip the HTML render
"""
from __future__ import annotations

import argparse
import html as html_lib
import json
import re
import subprocess
import sys
from pathlib import Path
from typing import Optional

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn

# ── canvas ──────────────────────────────────────────────────────────────────
SLIDE_W, SLIDE_H = 1920, 1080
EMU_PER_PT = 12700

# NOTE: font names use SINGLE quotes — these strings are emitted inside a
# double-quoted style="" attribute, and nested double quotes would truncate
# the attribute (silently dropping every declaration after font-family,
# e.g. color/font-weight → text falls back to the inherited default).
DEFAULT_FONT_STACK = (
    "'PingFang SC','Microsoft YaHei','Source Han Sans SC',"
    "'Helvetica Neue',Arial,sans-serif"
)


# ── affine transform (EMU→EMU), used to flatten nested groups ────────────────
class Xf:
    """Maps a shape's container-local EMU coords to absolute slide EMU."""

    def __init__(self, ox=0.0, oy=0.0, sx=1.0, sy=1.0):
        self.ox, self.oy, self.sx, self.sy = ox, oy, sx, sy

    def x(self, v):  return self.ox + v * self.sx
    def y(self, v):  return self.oy + v * self.sy
    def w(self, v):  return v * self.sx
    def h(self, v):  return v * self.sy

    def enter_group(self, grp) -> "Xf":
        """Return the child transform for a GroupShape `grp`."""
        # group's own box, mapped to absolute slide EMU through *this* xf
        gx, gy = self.x(grp.left), self.y(grp.top)
        gw, gh = self.w(grp.width), self.h(grp.height)
        # child coordinate space (a:chOff / a:chExt)
        cox = coy = 0.0
        cex = grp.width or 1
        cey = grp.height or 1
        xfrm = grp._element.find(qn("p:grpSpPr"))
        if xfrm is not None:
            xfrm = xfrm.find(qn("a:xfrm"))
        if xfrm is not None:
            ch_off = xfrm.find(qn("a:chOff"))
            ch_ext = xfrm.find(qn("a:chExt"))
            if ch_off is not None:
                cox = float(ch_off.get("x", 0)); coy = float(ch_off.get("y", 0))
            if ch_ext is not None:
                cex = float(ch_ext.get("cx", 1)) or 1
                cey = float(ch_ext.get("cy", 1)) or 1
        sx = gw / cex
        sy = gh / cey
        return Xf(gx - cox * sx, gy - coy * sy, sx, sy)


# ── theme-color resolution ────────────────────────────────────────────────────
# Populated once per presentation in main(): MSO_THEME_COLOR member → "#RRGGBB".
# Without this, scheme-colored text (the deck author's "tx1/lt1/accent1") has no
# .rgb and falls back to a guessed default → light text on dark cards goes dark.
_THEME: dict = {}
# parallel map keyed by the raw schemeClr val string ("accent1","lt1","tx1"…),
# for resolving colors inside gradient stops / XML we read by hand.
_THEME_BY_NAME: dict = {}

# clrScheme child tag → theme-color enum members it backs (incl. the tx/bg
# aliases that the master's clrMap maps to dk/lt by default).
_SCHEME_TAGS = [
    ("dk1", [MSO_THEME_COLOR.DARK_1, MSO_THEME_COLOR.TEXT_1]),
    ("lt1", [MSO_THEME_COLOR.LIGHT_1, MSO_THEME_COLOR.BACKGROUND_1]),
    ("dk2", [MSO_THEME_COLOR.DARK_2, MSO_THEME_COLOR.TEXT_2]),
    ("lt2", [MSO_THEME_COLOR.LIGHT_2, MSO_THEME_COLOR.BACKGROUND_2]),
    ("accent1", [MSO_THEME_COLOR.ACCENT_1]),
    ("accent2", [MSO_THEME_COLOR.ACCENT_2]),
    ("accent3", [MSO_THEME_COLOR.ACCENT_3]),
    ("accent4", [MSO_THEME_COLOR.ACCENT_4]),
    ("accent5", [MSO_THEME_COLOR.ACCENT_5]),
    ("accent6", [MSO_THEME_COLOR.ACCENT_6]),
    ("hlink", [MSO_THEME_COLOR.HYPERLINK]),
    ("folHlink", [MSO_THEME_COLOR.FOLLOWED_HYPERLINK]),
]


def _clr_to_hex(node) -> Optional[str]:
    """Read a hex from a clrScheme child (<a:srgbClr> or <a:sysClr lastClr>)."""
    srgb = node.find(qn("a:srgbClr"))
    if srgb is not None and srgb.get("val"):
        return f"#{srgb.get('val')}"
    sysc = node.find(qn("a:sysClr"))
    if sysc is not None and sysc.get("lastClr"):
        return f"#{sysc.get('lastClr')}"
    return None


def build_theme_map(prs) -> dict:
    """Map theme-color members → hex from the first master's theme clrScheme."""
    from lxml import etree
    out: dict = {}
    try:
        master = prs.slide_masters[0]
        theme_part = None
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return out
        # the theme is a generic Part (no ._element) — parse its blob
        root = etree.fromstring(theme_part.blob)
        clr_scheme = root.find(f".//{qn('a:clrScheme')}")
        if clr_scheme is None:
            return out
        for tag, members in _SCHEME_TAGS:
            node = clr_scheme.find(qn(f"a:{tag}"))
            if node is None:
                continue
            hexc = _clr_to_hex(node)
            if hexc:
                for m in members:
                    out[m] = hexc
                _THEME_BY_NAME[tag] = hexc
        # tx/bg aliases → resolve via the MASTER's ACTUAL <p:clrMap>, not the
        # hardcoded default. A dark-themed master INVERTS these (clrMap
        # tx1="lt1" bg1="dk1"); hardcoding tx1→dk1 made every inheritance-only
        # title/body run resolve to dk1 (#000000 = invisible black on the dark
        # slide). Reading the real clrMap makes tx1→lt1 (light) on dark decks,
        # and is a no-op on normal light decks (clrMap == default).
        _DEFAULT_CLRMAP = {"tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2"}
        clrmap = {}
        try:
            cm = master.element.find(qn("p:clrMap"))
            if cm is not None:
                clrmap = dict(cm.attrib)
        except Exception:
            pass
        for alias in ("tx1", "bg1", "tx2", "bg2"):
            base = clrmap.get(alias, _DEFAULT_CLRMAP[alias])
            if base in _THEME_BY_NAME:
                _THEME_BY_NAME[alias] = _THEME_BY_NAME[base]
    except Exception:
        pass
    return out


def _scheme_name_hex(val: str) -> Optional[str]:
    """Resolve a raw schemeClr val ('accent1','lt1','tx1','bg1'…) to hex."""
    return _THEME_BY_NAME.get(val)


def _xml_color_hex(node) -> Optional[str]:
    """Resolve a DrawingML color child (<a:srgbClr>/<a:schemeClr>/<a:sysClr>)."""
    srgb = node.find(qn("a:srgbClr"))
    if srgb is not None and srgb.get("val"):
        return f"#{srgb.get('val')}"
    sysc = node.find(qn("a:sysClr"))
    if sysc is not None and sysc.get("lastClr"):
        return f"#{sysc.get('lastClr')}"
    sch = node.find(qn("a:schemeClr"))
    if sch is not None and sch.get("val"):
        return _scheme_name_hex(sch.get("val"))
    return None


def _alpha_of(node) -> float:
    """Opacity 0..1 from a DrawingML color node's resolved color child
    (<a:srgbClr>/<a:schemeClr>/<a:sysClr>) <a:alpha val=N> (N in 1/1000 of a
    percent). 1.0 (fully opaque) when absent. python-pptx never exposes this, so
    a half-transparent scrim/glass panel would otherwise flatten to opaque and
    cover the artwork beneath it."""
    for tag in ("a:srgbClr", "a:schemeClr", "a:sysClr"):
        c = node.find(qn(tag))
        if c is not None:
            a = c.find(qn("a:alpha"))
            if a is not None and a.get("val") is not None:
                try:
                    return max(0.0, min(1.0, int(a.get("val")) / 100000.0))
                except Exception:
                    return 1.0
            return 1.0
    return 1.0


def _hex_to_rgba(hexc: str, a: float) -> str:
    h = hexc.lstrip("#")
    if len(h) == 3:
        h = "".join(ch * 2 for ch in h)
    try:
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return f"rgba({r}, {g}, {b}, {a:.3f})"
    except Exception:
        return hexc


def _css_color(node) -> Optional[str]:
    """Resolve a DrawingML color child to a CSS color, carrying <a:alpha> as
    rgba() so semi-transparent fills (scrims / glass panels) keep their
    transparency. Plain hex when fully opaque."""
    hexc = _xml_color_hex(node)
    if not hexc:
        return None
    a = _alpha_of(node)
    return hexc if a >= 0.999 else _hex_to_rgba(hexc, a)


def gradient_css(shape) -> Optional[str]:
    """Parse a shape's <a:gradFill> into a CSS linear-gradient(), or None."""
    try:
        sp_pr = shape._element.find(qn("p:spPr"))
        if sp_pr is None:
            return None
        grad = sp_pr.find(qn("a:gradFill"))
        if grad is None:
            return None
        gs_lst = grad.find(qn("a:gsLst"))
        if gs_lst is None:
            return None
        stops = []
        for gs in gs_lst.findall(qn("a:gs")):
            col = _css_color(gs)   # carries <a:alpha> as rgba() (soft-glow fades)
            if col:
                stops.append((int(gs.get("pos", "0")) / 1000.0, col))
        if len(stops) < 2:
            return None
        # angle: PPT <a:lin ang> is 1/60000 deg clockwise from East (3 o'clock).
        # CSS linear-gradient 0deg = "to top"; convert via (ppt + 90) mod 360.
        ang_css = "to bottom"
        lin = grad.find(qn("a:lin"))
        if lin is not None and lin.get("ang") is not None:
            ang = (int(lin.get("ang")) / 60000.0 + 90) % 360
            ang_css = f"{ang:.0f}deg"
        stop_str = ", ".join(f"{c} {p:.0f}%" for p, c in stops)
        return f"linear-gradient({ang_css}, {stop_str})"
    except Exception:
        return None


_FONT_SCHEME: dict = {}
_WEIGHT_SUFFIX = re.compile(
    r"\s+(Light|Regular|Medium|DemiBold?|SemiBold?|Bold|Heavy|Black|Thin|"
    r"ExtraLight|UltraLight|Normal|Book)$", re.I)


def build_font_scheme(prs) -> None:
    """Theme fontScheme → {'mj-lt','mn-lt','mj-ea','mn-ea'} so a run typeface that
    is a theme reference (+mj-lt / +mn-ea) resolves to the real font name."""
    _FONT_SCHEME.clear()
    try:
        master = prs.slide_masters[0]
        theme_part = None
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return
        root = etree.fromstring(theme_part.blob)
        fs = root.find(f".//{qn('a:fontScheme')}")
        if fs is None:
            return
        for which, tag in (("mj", "a:majorFont"), ("mn", "a:minorFont")):
            fnode = fs.find(qn(tag))
            if fnode is None:
                continue
            for sc, t in (("lt", "a:latin"), ("ea", "a:ea")):
                node = fnode.find(qn(t))
                if node is not None and node.get("typeface"):
                    _FONT_SCHEME["%s-%s" % (which, sc)] = node.get("typeface")
            for f in fnode.findall(qn("a:font")):   # CJK via script="Hans" (ea 常空)
                if f.get("script") == "Hans" and f.get("typeface"):
                    _FONT_SCHEME.setdefault("%s-ea" % which, f.get("typeface"))
    except Exception:
        pass


def _resolve_theme_font(tf: str) -> str:
    return _FONT_SCHEME.get(tf[1:], "") if tf and tf.startswith("+") else tf


def _run_fonts(rPr) -> Optional[str]:
    """CSS font-family value from a run's <a:latin>/<a:ea> typefaces — latin first
    (ASCII/digits) then ea (CJK), mirroring PowerPoint's per-script font pick. The
    weight-stripped family is appended so 'FZLanTingHeiPro_GB18030 Medium' also
    matches the installed family 'FZLanTingHeiPro_GB18030'. None if unset."""
    if rPr is None:
        return None
    fams: list = []
    for tag in ("a:latin", "a:ea"):
        node = rPr.find(qn(tag))
        if node is None:
            continue
        tf = _resolve_theme_font(node.get("typeface") or "")
        for cand in (tf, _WEIGHT_SUFFIX.sub("", tf)):
            if cand and cand not in fams:
                fams.append(cand)
    return ", ".join('"%s"' % f for f in fams) if fams else None


def _run_gradient_css(rPr) -> Optional[str]:
    """Full CSS linear-gradient() from a run's <a:gradFill>, for gradient TEXT
    (rendered via background-clip:text). None when the run has no gradient. The
    flat first-stop color stays as the fallback in `color`."""
    if rPr is None:
        return None
    gf = rPr.find(qn("a:gradFill"))
    if gf is None:
        return None
    gs_lst = gf.find(qn("a:gsLst"))
    if gs_lst is None:
        return None
    stops = []
    for gs in gs_lst.findall(qn("a:gs")):
        col = _css_color(gs)
        if col:
            stops.append((int(gs.get("pos", "0")) / 1000.0, col))
    if len(stops) < 2:
        return None
    ang_css = "to bottom"
    lin = gf.find(qn("a:lin"))
    if lin is not None and lin.get("ang") is not None:
        ang = (int(lin.get("ang")) / 60000.0 + 90) % 360
        ang_css = f"{ang:.0f}deg"
    return (f"linear-gradient({ang_css}, "
            + ", ".join(f"{c} {p:.0f}%" for p, c in stops) + ")")


# ── color helpers ─────────────────────────────────────────────────────────────
def rgb_hex(color) -> Optional[str]:
    """Best-effort RGB hex from a python-pptx ColorFormat. None if unresolved.
    Resolves scheme/theme colors via the presentation theme map (_THEME)."""
    try:
        if color is None or color.type is None:
            return None
        # explicit RGB → direct (getattr swallows the AttributeError that
        # .rgb raises for non-RGB color types)
        rgb = getattr(color, "rgb", None)
        if rgb is not None:
            return f"#{str(rgb)}"
        # scheme/theme color → look up resolved hex
        tc = getattr(color, "theme_color", None)
        if tc is not None and tc in _THEME:
            return _THEME[tc]
    except Exception:
        return None
    return None


def fill_hex(fill) -> Optional[str]:
    """Safe solid-fill hex. Returns None for no-fill / gradient / picture /
    pattern / unresolved — never raises (python-pptx raises on .fore_color
    for non-solid fills)."""
    try:
        from pptx.enum.dml import MSO_FILL
        if fill is not None and fill.type == MSO_FILL.SOLID:
            return rgb_hex(fill.fore_color)
    except Exception:
        return None
    return None


def emu(v) -> float:
    return float(v) if v is not None else 0.0


# ── geometry → px ─────────────────────────────────────────────────────────────
class Canvas:
    """Uniform scale + center offset from slide EMU into 1920×1080 px."""

    def __init__(self, slide_w_emu: int, slide_h_emu: int):
        self.scale = min(SLIDE_W / slide_w_emu, SLIDE_H / slide_h_emu)
        self.off_x = (SLIDE_W - slide_w_emu * self.scale) / 2
        self.off_y = (SLIDE_H - slide_h_emu * self.scale) / 2

    def px_left(self, emu_x):  return emu_x * self.scale + self.off_x
    def px_top(self, emu_y):   return emu_y * self.scale + self.off_y
    def px_size(self, emu_d):  return emu_d * self.scale
    def pt_to_px(self, pt):    return pt * EMU_PER_PT * self.scale


# ── slide background ──────────────────────────────────────────────────────────
def slide_bg(slide, prs) -> str:
    """Resolve a background color from slide → layout → master. Default #FFF."""
    from pptx.enum.dml import MSO_FILL
    for src in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        try:
            bg = src.background
            if bg.fill.type == MSO_FILL.SOLID:
                hexc = rgb_hex(bg.fill.fore_color)
                if hexc:
                    return hexc
        except Exception:
            continue
    return "#FFFFFF"


# ── element id ─────────────────────────────────────────────────────────────────
class IdGen:
    """Stable per-slide element ids: e{slide}_{n} (1-based slide, 0-based n)."""

    def __init__(self, slide_no: int):
        self.slide_no = slide_no
        self.n = 0

    def next(self) -> str:
        eid = f"e{self.slide_no}_{self.n}"
        self.n += 1
        return eid


# ── geometry → element x/y/w/h (px on 1920×1080) ───────────────────────────────
def el_geom(cv: Canvas, xf: Xf, shape) -> dict:
    """Absolute px geometry for one shape, group transform composed in."""
    return {
        "x": round(cv.px_left(xf.x(emu(shape.left))), 1),
        "y": round(cv.px_top(xf.y(emu(shape.top))), 1),
        "w": round(cv.px_size(xf.w(emu(shape.width))), 1),
        "h": round(cv.px_size(xf.h(emu(shape.height))), 1),
    }


def _rotation_style(shape) -> Optional[str]:
    """Non-trivial rotation → a CSS `transform` string for the shape `style`
    escape-hatch (text/image stay clean; only shapes carry style)."""
    rot = getattr(shape, "rotation", 0) or 0
    if rot and abs(rot) > 0.01:
        return f"transform:rotate({rot:.2f}deg)"
    return None


# ── text element ───────────────────────────────────────────────────────────────
_ANCHOR = {
    MSO_ANCHOR.TOP: "top", MSO_ANCHOR.MIDDLE: "middle",
    MSO_ANCHOR.BOTTOM: "bottom",
}


def _bullet_char(p) -> Optional[str]:
    """Bullet marker for a paragraph, else None (prepended to the run text so it
    stays structured + editable rather than becoming layout markup)."""
    pPr = p._p.find(qn("a:pPr"))
    if pPr is None:
        return None
    if pPr.find(qn("a:buNone")) is not None:
        return None
    bc = pPr.find(qn("a:buChar"))
    if bc is not None:
        return bc.get("char", "•")
    if pPr.find(qn("a:buAutoNum")) is not None:
        return "•"  # numbered → approximate; v0 doesn't track the counter
    return None


def _emu_to_px_pt(cv: Canvas, sz) -> Optional[float]:
    """python-pptx font.size (EMU/Length) → px on the canvas."""
    if sz is None:
        return None
    try:
        return round(cv.px_size(int(sz)), 1)
    except Exception:
        return None


# ── text style inheritance (master txStyles + placeholder lstStyle) ─────────────
# python-pptx exposes only what is written on the RUN (a:rPr). Size/color that a
# slide inherits from the master's p:txStyles or a placeholder's lstStyle is
# invisible to it → on this real deck 40% of runs lost color and 13% lost size,
# rendering as 16px browser-black. This map + _resolve_run_style rebuild that
# inheritance chain: run → paragraph defRPr → shape lstStyle[level] → master
# txStyle[(placeholder-kind, level)] → theme tx1 default.
_TXSTYLE: dict = {}


def _defRPr_size_color(defRPr):
    """(size_pt, css_color) from an <a:defRPr> node. size_pt from @sz (1/100 pt);
    color from solidFill (alpha-aware) or first gradFill stop."""
    sz = None
    color = None
    if defRPr is not None:
        v = defRPr.get("sz")
        if v:
            try:
                sz = int(v) / 100.0
            except Exception:
                pass
        sf = defRPr.find(qn("a:solidFill"))
        if sf is not None:
            color = _css_color(sf)
        else:
            gf = defRPr.find(qn("a:gradFill"))
            if gf is not None:
                first = gf.find(f".//{qn('a:gs')}")
                if first is not None:
                    color = _css_color(first)
    return sz, color


def build_text_style_map(prs) -> None:
    """Populate _TXSTYLE[(style, level)] = {size_pt, color} from the master's
    p:txStyles (titleStyle / bodyStyle / otherStyle, levels 1..9)."""
    _TXSTYLE.clear()
    try:
        master = prs.slide_masters[0]
        txst = master.element.find(qn("p:txStyles"))
        if txst is None:
            return
        for style, tag in (("title", "p:titleStyle"),
                           ("body", "p:bodyStyle"),
                           ("other", "p:otherStyle")):
            st = txst.find(qn(tag))
            if st is None:
                continue
            for lvl in range(9):
                lvlpr = st.find(qn(f"a:lvl{lvl + 1}pPr"))
                if lvlpr is None:
                    continue
                sz, color = _defRPr_size_color(lvlpr.find(qn("a:defRPr")))
                _TXSTYLE[(style, lvl)] = {"size_pt": sz, "color": color}
    except Exception:
        pass


# Per-slide layout-placeholder style: idx -> {level: {size_pt, color}}. The
# slide-LAYOUT placeholder (matched by idx) overrides the master txStyle — this
# deck's title is master-44pt but layout-26pt, so without this level the title
# inherits 44pt and overflows its box. Set per slide in compose_slide.
_CUR_LAYOUT_PH: dict = {}


def set_layout_ph_map(slide) -> None:
    _CUR_LAYOUT_PH.clear()
    try:
        layout_phs = slide.slide_layout.placeholders
    except Exception:
        return
    for ph in layout_phs:
        # per-placeholder isolation: one malformed placeholder must NOT abort
        # override extraction for the rest of the slide's placeholders.
        try:
            idx = ph.placeholder_format.idx
            tb = ph.text_frame._txBody
            levels: dict = {}
            lst = tb.find(qn("a:lstStyle"))
            if lst is not None:
                for lvl in range(9):
                    lp = lst.find(qn(f"a:lvl{lvl + 1}pPr"))
                    if lp is not None:
                        sz, color = _defRPr_size_color(lp.find(qn("a:defRPr")))
                        if sz is not None or color:
                            levels[lvl] = {"size_pt": sz, "color": color}
            # fallback: a defRPr carried on the placeholder's first paragraph
            if 0 not in levels:
                d = tb.find(f".//{qn('a:defRPr')}")
                if d is not None:
                    sz, color = _defRPr_size_color(d)
                    if sz is not None or color:
                        levels[0] = {"size_pt": sz, "color": color}
            if levels:
                _CUR_LAYOUT_PH[idx] = levels
        except Exception:
            continue


def _ph_style(shape) -> str:
    """Which master txStyle bucket a shape inherits: title / body / other."""
    try:
        if not shape.is_placeholder:
            return "other"
        from pptx.enum.shapes import PP_PLACEHOLDER as PPP
        t = shape.placeholder_format.type
        if t in (PPP.TITLE, PPP.CENTER_TITLE):
            return "title"
        return "body"
    except Exception:
        return "other"


def _resolve_run_style(r, p, shape, cv: Canvas):
    """(size_px or None, css_color or None) for a run, resolving inheritance when
    the run itself omits size/color (the 承重墙 fix)."""
    rPr = r._r.find(qn("a:rPr"))
    # Resolve the run's OWN color from the RAW rPr XML instead of python-pptx's
    # r.font.color, for two reasons:
    #   1. Reading r.font.color MUTATES the element — it replaces <a:gradFill>
    #      with an empty <a:solidFill> — destroying gold-gradient text (titles /
    #      big stat numbers / card captions) and dropping it to a default.
    #   2. python-pptx's theme_color path is clrMap-BLIND: an explicit
    #      schemeClr "tx1" resolves to dk1 (#000000, invisible black) instead of
    #      the master's inverted lt1 (light) on dark-themed decks.
    # _css_color routes through _scheme_name_hex/_THEME_BY_NAME (clrMap-aware)
    # and carries <a:alpha> as rgba.
    color = None
    if rPr is not None:
        sf = rPr.find(qn("a:solidFill"))
        if sf is not None:
            color = _css_color(sf)
        if color is None:
            gf = rPr.find(qn("a:gradFill"))
            if gf is not None:
                first = gf.find(f".//{qn('a:gs')}")
                if first is not None:
                    color = _css_color(first)
    size = _emu_to_px_pt(cv, r.font.size)
    if color is None:
        color = rgb_hex(r.font.color)   # rare non-fill color paths
    if size is not None and color is not None:
        return size, color
    try:
        level = p.level or 0
    except Exception:
        level = 0
    # paragraph defRPr, then the shape's own lstStyle for this level
    cands = []
    pPr = p._p.find(qn("a:pPr"))
    if pPr is not None:
        cands.append(pPr.find(qn("a:defRPr")))
    lst = shape._element.find(f".//{qn('a:lstStyle')}")
    if lst is not None:
        lp = lst.find(qn(f"a:lvl{level + 1}pPr"))
        if lp is not None:
            cands.append(lp.find(qn("a:defRPr")))
    for d in cands:
        if d is None:
            continue
        s, c = _defRPr_size_color(d)
        if size is None and s is not None:
            size = cv.pt_to_px(s)
        if color is None and c:
            color = c
    # slide-LAYOUT placeholder (by idx) — overrides the master (this deck's
    # title is layout-26pt over master-44pt).
    if (size is None or color is None):
        try:
            if shape.is_placeholder:
                lvls = _CUR_LAYOUT_PH.get(shape.placeholder_format.idx, {})
                st = lvls.get(level) or lvls.get(0)
                if st:
                    if size is None and st.get("size_pt") is not None:
                        size = cv.pt_to_px(st["size_pt"])
                    if color is None and st.get("color"):
                        color = st["color"]
        except Exception:
            pass
    # master txStyles by placeholder kind + level
    if size is None or color is None:
        kind = _ph_style(shape)
        st = _TXSTYLE.get((kind, level)) or _TXSTYLE.get((kind, 0))
        if st:
            if size is None and st.get("size_pt") is not None:
                size = cv.pt_to_px(st["size_pt"])
            if color is None and st.get("color"):
                color = st["color"]
    # last resort: theme default text color (what PowerPoint paints by default)
    if color is None:
        color = _THEME_BY_NAME.get("tx1")
    return size, color


def text_runs(shape, cv: Canvas) -> list:
    """Flatten a text frame's paragraphs into a single runs[] list. Paragraph
    breaks and soft <a:br/> become a run whose text starts with '\\n' so the
    renderer's _esc_br turns it into a <br>. Bullets are prepended to the first
    run of the bulleted paragraph as text (stays structured + editable)."""
    from pptx.text.text import _Run
    tf = shape.text_frame
    runs: list = []
    first_para = True
    for p in tf.paragraphs:
        para_runs: list = []
        for child in p._p:
            tag = child.tag
            if tag == qn("a:r"):
                r = _Run(child, p)
                txt = r.text
                if not txt:
                    continue
                run = {"text": txt}
                if r.font.bold:
                    run["bold"] = True
                # gradient TEXT: capture the full gradient from the raw rPr
                # (before _resolve_run_style, which may touch font.color). The
                # flat first-stop color is still resolved as the fallback.
                _rPr = r._r.find(qn("a:rPr"))
                grad = _run_gradient_css(_rPr)
                font = _run_fonts(_rPr)
                size, col = _resolve_run_style(r, p, shape, cv)
                if col:
                    run["color"] = col
                if size is not None:
                    run["size"] = round(size, 1)
                if grad:
                    run["grad"] = grad
                if font:
                    run["font"] = font
                para_runs.append(run)
            elif tag == qn("a:br"):
                para_runs.append({"text": "\n"})
            elif tag == qn("a:fld"):  # slide number / date field → cached text
                t = child.find(qn("a:t"))
                if t is not None and t.text:
                    para_runs.append({"text": t.text})
        if not para_runs:
            continue
        # bullet → prepend to the paragraph's first real run
        bullet = _bullet_char(p)
        if bullet:
            para_runs[0]["text"] = f"{bullet} {para_runs[0]['text']}"
        # paragraph separator: a leading "\n" on the first run of every
        # paragraph after the first (renders as <br> between paragraphs).
        if not first_para:
            para_runs[0]["text"] = "\n" + para_runs[0]["text"]
        first_para = False
        runs.extend(para_runs)
    return runs


def _autofit_scale(shape) -> float:
    """PowerPoint's stored shrink-to-fit factor: <a:bodyPr><a:normAutofit
    fontScale=N> (N in 1/1000 %). A long title in a fixed box carries e.g.
    fontScale=62500 (62.5%); without applying it the inherited base size
    overflows the box and collides with the content below it. 1.0 if none."""
    try:
        bodyPr = shape.text_frame._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            na = bodyPr.find(qn("a:normAutofit"))
            if na is not None and na.get("fontScale"):
                return max(0.1, int(na.get("fontScale")) / 100000.0)
    except Exception:
        pass
    return 1.0


def text_element(shape, cv: Canvas, xf: Xf, eid: str) -> Optional[dict]:
    runs = text_runs(shape, cv)
    if not runs:
        return None
    # apply PowerPoint's shrink-to-fit factor to the (now inheritance-resolved)
    # run sizes so an autofit title doesn't overflow its box.
    scale = _autofit_scale(shape)
    if scale < 0.999:
        for run in runs:
            if run.get("size") is not None:
                run["size"] = round(run["size"] * scale, 1)
    el = {"id": eid, "type": "text"}
    el.update(el_geom(cv, xf, shape))
    try:
        anchor = _ANCHOR.get(shape.text_frame.vertical_anchor)
        if anchor:
            el["anchor"] = anchor
    except Exception:
        pass
    # text frame insets (left/right/top/bottom) in px
    try:
        tf = shape.text_frame
        insets = [tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom]
        if any(i is not None for i in insets):
            el["insets"] = [round(cv.px_size(int(i)), 1) if i is not None else 0
                            for i in insets]
    except Exception:
        pass
    # paragraph horizontal alignment (first paragraph) → renderer maps it to the
    # text wrapper's text-align. Only explicit algn here; inherited alignment is
    # resolved by the style-inheritance pass (R1).
    try:
        from pptx.enum.text import PP_ALIGN
        a = {PP_ALIGN.LEFT: "left", PP_ALIGN.CENTER: "center",
             PP_ALIGN.RIGHT: "right", PP_ALIGN.JUSTIFY: "justify"}.get(
            shape.text_frame.paragraphs[0].alignment)
        if a:
            el["align"] = a
    except Exception:
        pass
    el["runs"] = runs
    return el


# ── shape appearance (fill / gradient / border / radius / preset / svg) ─────────
def _preset_geom_name(shape) -> Optional[str]:
    """OOXML preset geometry name (<a:prstGeom prst="...">), e.g. rect /
    roundRect / ellipse. None for custom/freeform geometry."""
    try:
        sp_pr = shape._element.find(qn("p:spPr"))
        if sp_pr is None:
            return None
        prst = sp_pr.find(qn("a:prstGeom"))
        if prst is not None:
            return prst.get("prst")
    except Exception:
        pass
    return None


def _line_is_nofill(shape) -> bool:
    """True when the shape's <a:ln> explicitly carries <a:noFill/> — i.e.
    'width exists but the stroke is none', which PowerPoint renders as NO
    line. python-pptx only exposes ln.width (still > 0 here), so without
    reading the raw OOXML we'd treat it as a real border, fail to resolve a
    color, and fabricate a phantom #888888 box around every such element."""
    try:
        sp_pr = shape._element.find(qn("p:spPr"))
        ln_el = sp_pr.find(qn("a:ln")) if sp_pr is not None else None
        return ln_el is not None and ln_el.find(qn("a:noFill")) is not None
    except Exception:
        return False


def _border_obj(shape, cv: Canvas) -> Optional[dict]:
    try:
        if _line_is_nofill(shape):
            return None
        ln = shape.line
        if ln.width is not None and int(ln.width) > 0:
            color = rgb_hex(ln.color) or "#888888"
            width = max(1.0, round(int(ln.width) / EMU_PER_PT, 1))
            return {"color": color, "width": width}
    except Exception:
        pass
    return None


_ROUNDED_PRESETS = {"roundRect", "round1Rect", "round2SameRect",
                    "round2DiagRect", "roundRectCallout"}


def _prst_adj(shape, gd_name: str = "adj") -> Optional[float]:
    """The named preset-geometry adjust value as a fraction (0..1) from
    <a:prstGeom><a:avLst><a:gd name=... fmla="val N">. None if absent."""
    try:
        sp_pr = shape._element.find(qn("p:spPr"))
        prst = sp_pr.find(qn("a:prstGeom")) if sp_pr is not None else None
        av = prst.find(qn("a:avLst")) if prst is not None else None
        if av is None:
            return None
        for gd in av.findall(qn("a:gd")):
            if gd.get("name") == gd_name:
                fmla = gd.get("fmla", "")
                if fmla.startswith("val "):
                    return int(fmla.split()[1]) / 100000.0
    except Exception:
        pass
    return None


def _radius_px(shape, cv: Canvas) -> Optional[float]:
    """Corner radius (px) for rounded-rect presets, from the REAL adj value
    (<a:gd name="adj">) rather than a hardcoded 16% (which over-rounded cards
    ~3.5×). Radius = adj-fraction × short side; PPT's roundRect default is
    16.667% when no adj is authored."""
    name = _preset_geom_name(shape)
    if name in _ROUNDED_PRESETS:
        try:
            w = cv.px_size(emu(shape.width))
            h = cv.px_size(emu(shape.height))
            adj = _prst_adj(shape, "adj")
            frac = adj if adj is not None else 0.16667
            return round(min(w, h) * max(0.0, frac), 1)
        except Exception:
            return 14.0
    return None


def shape_appearance(shape, cv: Canvas) -> dict:
    """Fill / gradient / border / radius / kind for a non-freeform AUTO_SHAPE.
    Returns only the keys that resolve (clean elements carry no empty fields)."""
    out: dict = {}
    name = _preset_geom_name(shape)
    if name:
        out["kind"] = name
    # fill: solid hex OR gradient CSS
    try:
        from pptx.enum.dml import MSO_FILL
        fill = shape.fill
        if fill.type == MSO_FILL.SOLID:
            # read the solidFill node directly so <a:alpha> survives (a
            # half-transparent压暗蒙版 / 玻璃面板 would otherwise flatten to
            # opaque and盖死 the artwork below). schemeClr+alpha is invisible to
            # python-pptx, hence the XML path; rgb_hex is the opaque fallback.
            sp_pr = shape._element.find(qn("p:spPr"))
            sf = sp_pr.find(qn("a:solidFill")) if sp_pr is not None else None
            cssc = _css_color(sf) if sf is not None else None
            out["fill"] = cssc or rgb_hex(fill.fore_color) or out.get("fill")
            if not out.get("fill"):
                out.pop("fill", None)
        elif fill.type == MSO_FILL.GRADIENT:
            grad = gradient_css(shape)
            if grad:
                out["gradient"] = grad
    except Exception:
        pass
    border = _border_obj(shape, cv)
    if border:
        out["border"] = border
    # radius: explicit roundRect → computed px; ellipse → 50% via style
    if name == "ellipse":
        out["style"] = "border-radius:50%"
    else:
        radius = _radius_px(shape, cv)
        if radius is not None:
            out["radius"] = radius
    rot = _rotation_style(shape)
    if rot:
        out["style"] = (out.get("style") + ";" + rot) if out.get("style") else rot
    return out


# ── freeform / custGeom / line → inline SVG (normalized 0..100 box) ─────────────
def _custgeom_svg(shape) -> Optional[str]:
    """Parse <a:custGeom> path commands into an SVG <path>, normalized to a
    0..100 viewBox (renderer uses preserveAspectRatio:none, so coords map to %
    of the element box). Returns the inner SVG markup, or None."""
    try:
        sp_pr = shape._element.find(qn("p:spPr"))
        if sp_pr is None:
            return None
        cust = sp_pr.find(qn("a:custGeom"))
        if cust is None:
            return None
        path_lst = cust.find(qn("a:pathLst"))
        if path_lst is None:
            return None
        # fill / stroke from the shape
        fill = fill_hex(shape.fill) or "none"
        stroke, sw = "none", 0.0
        try:
            ln = shape.line
            # same noFill trap as _border_obj — width>0 + <a:noFill/> = no stroke
            if not _line_is_nofill(shape) and ln.width is not None and int(ln.width) > 0:
                stroke = rgb_hex(ln.color) or "#888888"
                sw = max(0.5, int(ln.width) / EMU_PER_PT)
        except Exception:
            pass
        if fill == "none" and stroke == "none":
            fill = "#888888"  # at least show the silhouette
        d_parts: list[str] = []
        for path in path_lst.findall(qn("a:path")):
            pw = float(path.get("w") or 0) or 1.0
            ph = float(path.get("h") or 0) or 1.0
            def nx(v): return round(float(v) / pw * 100, 2)
            def ny(v): return round(float(v) / ph * 100, 2)
            for cmd in path:
                t = cmd.tag
                if t == qn("a:moveTo"):
                    pt = cmd.find(qn("a:pt"))
                    d_parts.append(f"M{nx(pt.get('x'))},{ny(pt.get('y'))}")
                elif t == qn("a:lnTo"):
                    pt = cmd.find(qn("a:pt"))
                    d_parts.append(f"L{nx(pt.get('x'))},{ny(pt.get('y'))}")
                elif t == qn("a:cubicBezTo"):
                    pts = cmd.findall(qn("a:pt"))
                    if len(pts) == 3:
                        d_parts.append("C" + " ".join(
                            f"{nx(p.get('x'))},{ny(p.get('y'))}" for p in pts))
                elif t == qn("a:quadBezTo"):
                    pts = cmd.findall(qn("a:pt"))
                    if len(pts) == 2:
                        d_parts.append("Q" + " ".join(
                            f"{nx(p.get('x'))},{ny(p.get('y'))}" for p in pts))
                elif t == qn("a:close"):
                    d_parts.append("Z")
        if not d_parts:
            return None
        d = " ".join(d_parts)
        sw_attr = f' stroke="{stroke}" stroke-width="{sw:.1f}"' if stroke != "none" else ""
        return (f'<path d="{html_lib.escape(d, quote=True)}" '
                f'fill="{fill}"{sw_attr} vector-effect="non-scaling-stroke"/>')
    except Exception:
        return None


def line_svg(shape, cv: Canvas) -> str:
    """A LINE / connector as an SVG line in the normalized 0..100 box, honoring
    flipH / flipV so diagonals point the right way."""
    color, sw = None, 1.0
    try:
        color = rgb_hex(shape.line.color)
        if shape.line.width:
            sw = max(1.0, int(shape.line.width) / EMU_PER_PT)
    except Exception:
        pass
    color = color or "#8895AA"
    xfrm = shape._element.find(f".//{qn('a:xfrm')}")
    flip_h = xfrm is not None and xfrm.get("flipH") == "1"
    flip_v = xfrm is not None and xfrm.get("flipV") == "1"
    x1, y1, x2, y2 = (100 if flip_h else 0), (100 if flip_v else 0), \
                     (0 if flip_h else 100), (0 if flip_v else 100)
    return (f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" '
            f'stroke="{color}" stroke-width="{sw:.1f}" '
            f'vector-effect="non-scaling-stroke"/>')


def _pic_crop(shape):
    """[l, r, t, b] crop fractions from a picture's <a:srcRect> (vals are 1/1000
    of a percent; may be negative = padding). Visible region is x∈[l,1-r],
    y∈[t,1-b]. None when absent/zero. Without this the full image is stretched
    to the box (object-fit:fill) → wrong proportions for cropped pictures."""
    try:
        sr = shape._element.find(f".//{qn('a:srcRect')}")
        if sr is None:
            return None
        def frac(k):
            v = sr.get(k)
            return round(int(v) / 100000.0, 5) if v is not None else 0.0
        crop = [frac("l"), frac("r"), frac("t"), frac("b")]
        return crop if any(crop) else None
    except Exception:
        return None


def _media_poster_blob(shape):
    """(blob, ext) of a video/audio's poster (cover) frame, or None. python-pptx
    models a movie as a `Movie` with NO `.image`, so `shape.image` always raised
    → every video fell back to a full-bleed #0B0F18 rect that covered the slide's
    designed background. The poster IS an embedded image part: the shape's
    `<p:pic>` carries an `<a:blip r:embed>` pointing at it."""
    try:
        blip = shape._element.find(f".//{qn('a:blip')}")
        if blip is None:
            return None
        rId = blip.get(qn("r:embed"))
        if not rId:
            return None
        part = shape.part.related_part(rId)
        ct = (getattr(part, "content_type", "") or "").lower()
        ext = ("jpg" if "jpeg" in ct or "jpg" in ct
               else "gif" if "gif" in ct else "png")
        return part.blob, ext
    except Exception:
        return None


# ── per-shape → element dict(s) ────────────────────────────────────────────────
def emit_shape(shape, cv: Canvas, xf: Xf, idgen: IdGen, input_dir: Path,
               img_counter: list) -> list:
    """Return list of element dicts for one shape (GROUP → many, flattened)."""
    out: list = []
    st = shape.shape_type

    # GROUP → recurse with the composed child transform; emit children as
    # top-level elements (no wrapper).
    if st == MSO_SHAPE_TYPE.GROUP:
        child_xf = xf.enter_group(shape)
        for child in shape.shapes:
            out += emit_shape(child, cv, child_xf, idgen, input_dir, img_counter)
        return out

    # PICTURE → image element; blob extracted to input/ (original content).
    if st == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            img_counter[0] += 1
            ext = img.ext or "png"
            fname = f"img-{img_counter[0]:03d}.{ext}"
            (input_dir / fname).write_bytes(img.blob)
            el = {"id": idgen.next(), "type": "image", "src": f"input/{fname}"}
            el.update(el_geom(cv, xf, shape))
            crop = _pic_crop(shape)   # <a:srcRect> 裁剪,否则整图被 object-fit:fill 拉伸
            if crop:
                el["crop"] = crop
            out.append(el)
        except Exception:
            pass  # unreadable embedded image → silently drop (rare)
        return out

    # MEDIA (video/audio) → poster/cover frame if present, else a dark placeholder.
    if st == MSO_SHAPE_TYPE.MEDIA:
        poster = _media_poster_blob(shape)
        if poster:
            blob, ext = poster
            img_counter[0] += 1
            fname = f"media-{img_counter[0]:03d}.{ext}"
            (input_dir / fname).write_bytes(blob)
            el = {"id": idgen.next(), "type": "image", "src": f"input/{fname}"}
            el.update(el_geom(cv, xf, shape))
            out.append(el)
        else:
            el = {"id": idgen.next(), "type": "shape", "fill": "#0B0F18",
                  "kind": "rect"}
            el.update(el_geom(cv, xf, shape))
            out.append(el)
        return out

    # LINE / connector → shape with inline SVG.
    if st == MSO_SHAPE_TYPE.LINE:
        el = {"id": idgen.next(), "type": "shape", "svg": line_svg(shape, cv)}
        el.update(el_geom(cv, xf, shape))
        out.append(el)
        return out

    # FREEFORM (custGeom) → shape with inline SVG path.
    if st == MSO_SHAPE_TYPE.FREEFORM:
        svg = _custgeom_svg(shape)
        if svg:
            el = {"id": idgen.next(), "type": "shape", "svg": svg}
            el.update(el_geom(cv, xf, shape))
            out.append(el)
        # also emit any text the freeform carries (rare)
        if shape.has_text_frame and shape.text_frame.text.strip():
            tel = text_element(shape, cv, xf, idgen.next())
            if tel:
                out.append(tel)
        return out

    # TABLE → flatten cells to text + cell-fill shapes (no table layout type in
    # canvas; tables are rare and best reconstructed as positioned cells).
    if shape.has_table:
        out += _emit_table(shape, cv, xf, idgen)
        return out

    # AUTO_SHAPE / TEXT_BOX / PLACEHOLDER.
    # A non-textbox auto-shape with a custGeom (preset failed) → SVG fallback.
    is_textbox = st in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER)
    has_text = shape.has_text_frame and shape.text_frame.text.strip()

    # shape appearance (skip for pure text boxes/placeholders — those are clean
    # text only; a placeholder with a real fill still gets a backing shape).
    if not is_textbox:
        appearance = shape_appearance(shape, cv)
        # nothing reproducible (e.g. picture/pattern fill, no border) and no
        # text → try custGeom SVG, else drop (decorative, un-reconstructable).
        has_visible = bool(appearance.get("fill") or appearance.get("gradient")
                           or appearance.get("border") or appearance.get("svg"))
        if not has_visible:
            svg = _custgeom_svg(shape)
            if svg:
                appearance["svg"] = svg
                has_visible = True
        if has_visible:
            sel = {"id": idgen.next(), "type": "shape"}
            sel.update(el_geom(cv, xf, shape))
            sel.update(appearance)
            out.append(sel)

    # text → a separate clean text element on top (works for shapes-with-text
    # too: the backing shape + a text box at the same geometry).
    if has_text:
        tel = text_element(shape, cv, xf, idgen.next())
        if tel:
            out.append(tel)
    return out


def _emit_table(shape, cv: Canvas, xf: Xf, idgen: IdGen) -> list:
    """Flatten a PPT table into positioned cell shapes + cell text elements.
    Column/row sizes from the table's gridCol / row heights (EMU)."""
    out: list = []
    tbl = shape.table
    base_x = emu(shape.left)
    base_y = emu(shape.top)
    col_w = [emu(c.width) for c in tbl.columns]
    row_h = [emu(r.height) for r in tbl.rows]
    y = base_y
    for ri, row in enumerate(tbl.rows):
        x = base_x
        for ci, cell in enumerate(row.cells):
            cw = col_w[ci] if ci < len(col_w) else 0
            ch = row_h[ri] if ri < len(row_h) else 0
            gx = round(cv.px_left(xf.x(x)), 1)
            gy = round(cv.px_top(xf.y(y)), 1)
            gw = round(cv.px_size(xf.w(cw)), 1)
            gh = round(cv.px_size(xf.h(ch)), 1)
            # cell backing shape (fill + a light border so the grid reads)
            cfill = fill_hex(cell.fill)
            sel = {"id": idgen.next(), "type": "shape",
                   "x": gx, "y": gy, "w": gw, "h": gh, "kind": "rect",
                   "border": {"color": "#D0D5DD", "width": 1}}
            if cfill:
                sel["fill"] = cfill
            out.append(sel)
            # cell text
            txt = cell.text.strip()
            if txt:
                out.append({"id": idgen.next(), "type": "text",
                            "x": gx, "y": gy, "w": gw, "h": gh,
                            "anchor": "middle",
                            "insets": [8, 8, 6, 6],
                            "runs": [{"text": txt}]})
            x += cw
        y += row_h[ri] if ri < len(row_h) else 0
    return out


# ── un-reconstructable detection ───────────────────────────────────────────────
_DIAGRAM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


# An OLE object occupying less than this fraction of the slide is a degenerate
# remnant — e.g. a ~1px embedded-object stub left behind by an editor — NOT real
# content. Without this guard a single invisible remnant condemns an otherwise
# fully reconstructable page to an empty placeholder (real-deck regression: the
# 香格里拉 deck lost slides 24/25 — title + text + 2 full-bleed photos — to one
# 0.0000%-area OLE残骸 each).
_OLE_MIN_AREA_FRAC = 0.005


def _shape_is_hard(shape, cv: "Canvas | None" = None) -> bool:
    """True if a shape is structurally un-reconstructable: a live chart,
    a SmartArt diagram, or an OLE object. A sub-pixel OLE remnant (area below
    `_OLE_MIN_AREA_FRAC` of the slide) is treated as NOT hard so it cannot
    condemn a whole page."""
    st = shape.shape_type
    if st in (MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT):
        if cv is not None:
            try:
                px_w = cv.px_size(emu(shape.width))
                px_h = cv.px_size(emu(shape.height))
                if (px_w * px_h) / (SLIDE_W * SLIDE_H) < _OLE_MIN_AREA_FRAC:
                    return False
            except Exception:
                pass
        return True
    if st == MSO_SHAPE_TYPE.CHART:
        return True
    try:
        if getattr(shape, "has_chart", False):
            return True
    except Exception:
        pass
    if st == MSO_SHAPE_TYPE.DIAGRAM:
        return True
    # SmartArt = a graphicFrame whose graphicData references the diagram ns
    try:
        gd = shape._element.find(f".//{qn('a:graphicData')}")
        if gd is not None and _DIAGRAM_NS in (gd.get("uri") or ""):
            return True
    except Exception:
        pass
    return False


def _slide_is_hard(slide, cv: "Canvas | None" = None) -> bool:
    """A slide is un-reconstructable if ANY of its shapes is hard (chart/
    SmartArt/OLE). Groups are recursed. `cv` enables the sub-pixel OLE-remnant
    guard in `_shape_is_hard`."""
    def walk(shapes) -> bool:
        for sh in shapes:
            if _shape_is_hard(sh, cv):
                return True
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                if walk(sh.shapes):
                    return True
        return False
    return walk(slide.shapes)


# placeholder prompt text on layout/master that must NOT render as content
_PROMPT_RE = re.compile(
    r"单击此处|單擊此處|请输入|請輸入|点击.*编辑|按一下.*編輯|母版|母版文本样式|"
    r"Click to edit|XXXX|公司职称|人名$|^\s*第.*页\s*$"
)


def _emit_bg_picture(container, input_dir: Path, idgen: IdGen,
                     tagname: str, img_counter: list) -> Optional[dict]:
    """A slide/layout/master `<p:bg>` PICTURE fill → a full-bleed image element
    (the designed cover background lives in the layout/master, not the slide's
    own shapes)."""
    try:
        el = container._element
        bg = el.find(f".//{qn('p:bg')}")
        if bg is None:
            return None
        blip = bg.find(f".//{qn('a:blip')}")
        if blip is None:
            return None
        embed = blip.get(qn("r:embed"))
        if not embed:
            return None
        part = container.part.related_part(embed)
        ext = str(part.partname).rsplit(".", 1)[-1] or "png"
        img_counter[0] += 1
        fname = f"bg-{tagname}-{img_counter[0]:03d}.{ext}"
        (input_dir / fname).write_bytes(part.blob)
        return {"id": idgen.next(), "type": "image", "src": f"input/{fname}",
                "x": 0, "y": 0, "w": SLIDE_W, "h": SLIDE_H}
    except Exception:
        return None


def emit_template_elements(container, cv: Canvas, idgen: IdGen, input_dir: Path,
                           img_counter: list) -> list:
    """Layout/master decorative shapes (background pictures, logos, footers)
    UNDER the slide content. Skip placeholders (prompts) and prompt text."""
    out: list = []
    xf = Xf()
    for sh in container.shapes:
        try:
            if sh.is_placeholder:
                continue
            if sh.has_text_frame and _PROMPT_RE.search(sh.text_frame.text or ""):
                continue
        except Exception:
            pass
        out += emit_shape(sh, cv, xf, idgen, input_dir, img_counter)
    return out


# ── imported-provenance marker ──────────────────────────────────────────────────
def _lifted_marker(deck_stem: str, slide_no: int) -> str:
    """Provenance string for a PPTX-imported canvas slide, reusing the schema's
    existing `lifted` field (value = source ref). Format `pptx:<stem>#<N>` mirrors
    lift-slides' `<deck>#<index>` convention. The renderer emits this as
    `data-lifted=…`, and validate.py then downgrades this slide's
    CONTENT-AUTHORING violations (R05 ellipsis / R-KEY positional / R-LANG latin
    leaf / R10 palette / R06 / R-WHITE-TEXT / R-VIS-TIER / R-VIS-BODY-FLOOR) from
    error → WARNING: the content is faithfully CARRIED from the source PPTX, not
    hand-authored to the deck's content standards, so it is surfaced (a human
    chooses to fix) rather than blocking. STRUCTURAL / GEOMETRY rules (R-DOM /
    R-OVERFLOW / R-VIS-CARD-OVERFLOW / R-CSSVAR / duplicate-key R-KEY) stay
    full-severity error — a real overflow is a real failure regardless of import."""
    return f"pptx:{deck_stem}#{slide_no}"


# ── slide → canvas elements ────────────────────────────────────────────────────
def compose_slide(slide, idx: int, cv: Canvas, prs, input_dir: Path,
                  img_counter: list, deck_stem: str) -> dict:
    """Build one `layout:"canvas"` slide dict (data.elements[])."""
    slide_no = idx + 1
    idgen = IdGen(slide_no)
    set_layout_ph_map(slide)   # layout placeholder size/color overrides (per slide)

    # background as a backing shape so the slide isn't transparent.
    elements: list = []
    bg = slide_bg(slide, prs)
    if bg and bg.upper() != "#FFFFFF":
        elements.append({"id": idgen.next(), "type": "shape", "kind": "rect",
                         "x": 0, "y": 0, "w": SLIDE_W, "h": SLIDE_H,
                         "fill": bg})

    # TEMPLATE LAYER (under content): layout/master background image + decor.
    layout = slide.slide_layout
    master = layout.slide_master
    for cont, tg in ((master, "master"), (layout, "layout")):
        bgimg = _emit_bg_picture(cont, input_dir, idgen, tg, img_counter)
        if bgimg:
            elements.append(bgimg)
        elements += emit_template_elements(cont, cv, idgen, input_dir, img_counter)

    # the slide's own shapes
    xf = Xf()
    for shape in slide.shapes:
        elements += emit_shape(shape, cv, xf, idgen, input_dir, img_counter)

    return {
        "key": f"slide-{slide_no:03d}",
        "layout": "canvas",
        "screen_label": f"{slide_no:02d}",
        # IMPORTED PROVENANCE: this slide is verbatim-carried from a foreign
        # PPTX, not hand-authored — mark it `lifted` so the validator downgrades
        # CONTENT-AUTHORING rules to warnings (see _lifted_marker).
        "lifted": _lifted_marker(deck_stem, slide_no),
        "data": {
            "canvas_w": SLIDE_W,
            "canvas_h": SLIDE_H,
            "source_page": slide_no,
            "elements": elements,
        },
    }


def _default_renderer() -> Path:
    """Locate the feishu-deck-h5 renderer (deck-json/render-deck.py).
    pptx-to-deck is a TOP-LEVEL skill that uses feishu-deck-h5 as its render
    backend; feishu-deck-h5 normally sits as a SIBLING skill dir. Resolution:
      1. sibling  <skills>/feishu-deck-h5/   (this skill at <skills>/pptx-to-deck/)
      2. legacy nested grandparent (when still inside feishu-deck-h5/<sub>/)
      3. ~/.claude/skills/feishu-deck-h5 (registered symlink) fallback."""
    skills_dir = Path(__file__).resolve().parent.parent.parent  # <skills>/
    for cand in (skills_dir / "feishu-deck-h5",                 # sibling (new layout)
                 skills_dir,                                     # legacy nested
                 Path.home() / ".claude/skills/feishu-deck-h5"):  # registered symlink
        if (cand / "deck-json/render-deck.py").is_file():
            return cand
    return Path.home() / ".claude/skills/feishu-deck-h5"


# ── main ──────────────────────────────────────────────────────────────────────
def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description="PPTX → feishu-deck-h5 canvas deck.json")
    ap.add_argument("pptx", type=Path)
    ap.add_argument("out_dir", type=Path)
    ap.add_argument("--renderer", type=Path, default=_default_renderer())
    ap.add_argument("--limit", type=int, default=0)
    ap.add_argument("--raster", action="store_true",
                    help="DEPRECATED no-op (rasterization retired)")
    ap.add_argument("--full-raster", action="store_true",
                    help="DEPRECATED no-op (rasterization retired)")
    ap.add_argument("--inline", action="store_true")
    ap.add_argument("--no-render", action="store_true",
                    help="emit deck.json + assets only, skip the HTML render")
    ap.add_argument("--title", default=None)
    args = ap.parse_args(argv)

    if args.raster or args.full_raster:
        print("==> NOTE: --raster / --full-raster are retired no-ops "
              "(no screenshots; embedded pictures stay as image elements).",
              file=sys.stderr)

    if not args.pptx.is_file():
        print(f"ERROR: pptx not found: {args.pptx}", file=sys.stderr)
        return 1
    render_script = args.renderer / "deck-json/render-deck.py"
    if not args.no_render and not render_script.is_file():
        print(f"ERROR: render-deck.py not at {render_script}", file=sys.stderr)
        return 1

    args.out_dir.mkdir(parents=True, exist_ok=True)
    # image assets live in input/ (real, scannable paths in elements[].src).
    input_dir = args.out_dir / "input"
    input_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(args.pptx))
    global _THEME
    _THEME = build_theme_map(prs)
    print(f"==> theme colors resolved: {len(_THEME)} members")
    build_text_style_map(prs)
    print(f"==> master text styles resolved: {len(_TXSTYLE)} (style,level) entries")
    build_font_scheme(prs)
    cv = Canvas(prs.slide_width, prs.slide_height)

    print(f"==> {args.pptx.name}: {len(prs.slides)} slides, "
          f"canvas scale={cv.scale:.4f}")
    deck_slides = []
    unreconstructed: list[int] = []
    img_counter = [0]
    deck_stem = (args.title or args.pptx.stem)
    for idx, slide in enumerate(prs.slides):
        if args.limit and idx >= args.limit:
            break
        slide_no = idx + 1
        if _slide_is_hard(slide, cv):
            # live chart / SmartArt / OLE → placeholder + collect page number.
            unreconstructed.append(slide_no)
            deck_slides.append({
                "key": f"slide-{slide_no:03d}",
                "layout": "canvas",
                "screen_label": f"{slide_no:02d}",
                # imported provenance (see _lifted_marker): a placeholder is
                # still verbatim-carried foreign content, mark it lifted too.
                "lifted": _lifted_marker(deck_stem, slide_no),
                "data": {"placeholder": True, "source_page": slide_no,
                         "elements": []},
            })
            print(f"  · slide {slide_no:2d}  ←  UN-RECONSTRUCTABLE "
                  f"(chart/SmartArt/OLE) → placeholder")
            continue
        sd = compose_slide(slide, idx, cv, prs, input_dir, img_counter, deck_stem)
        deck_slides.append(sd)
        n = len(sd["data"]["elements"])
        print(f"  · slide {slide_no:2d}  ←  {len(slide.shapes)} shapes  "
              f"→ {n} elements")

    deck = {
        "version": "1.0",
        "deck": {
            "title": args.title or args.pptx.stem,
            "language": "zh-only",
            "mode": "rewrite",
        },
        "slides": deck_slides,
    }
    deck_path = args.out_dir / "deck.json"
    deck_path.write_text(json.dumps(deck, ensure_ascii=False, indent=2))
    print(f"==> wrote {deck_path}  ({len(deck_slides)} slides)")
    # the contract report line (empty list if none).
    print(f"unreconstructed slides: {unreconstructed}")

    if args.no_render:
        print(f"\n==> DONE (no render) → {deck_path}")
        return 0

    # Keep the DeckJSON schema validation gate ON (handoff contract: validate
    # before render). Skip the post-render HTML validator + the content/story
    # schema-fit refusal — neither applies to canvas slides. Skip render-deck's
    # own copy-assets: canvas_finish.make_portable does a path-independent,
    # CSS-url()-following self-contained pack (superset) right after.
    cmd = [sys.executable, str(render_script), str(deck_path), str(args.out_dir),
           "--skip-validate-html", "--skip-fit-check", "--skip-copy-assets"]
    if args.inline:
        cmd.append("--inline")
    print(f"==> rendering via {render_script}")
    res = subprocess.run(cmd, capture_output=True, text=True)
    if res.returncode != 0:
        print("RENDER FAILED:", file=sys.stderr)
        print(res.stdout, file=sys.stderr)
        print(res.stderr, file=sys.stderr)
        return 1
    for line in (res.stdout.strip().splitlines() or [""])[-6:]:
        print(f"    {line}")

    # render-layer finish (shared with rerender-deck.py): self-contained pack +
    # letterbox/fitText. fitText is the in-browser equivalent of PowerPoint's
    # autofit shrink-to-fit — single-line boxes that overflow get nowrap+scaleX
    # at runtime instead of clipping (measures real bbox, no estimation).
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from canvas_finish import make_portable, post_process  # noqa: E402
    make_portable(args.out_dir, args.renderer)
    post_process(args.out_dir, deck)

    # serve helper
    serve = args.out_dir / "serve.sh"
    serve.write_text(
        '#!/usr/bin/env bash\nset -e\ncd "$(dirname "$0")"\n'
        'PORT="${1:-8765}"\n'
        'echo "==> http://localhost:$PORT/index.html"\n'
        'python3 -m http.server "$PORT"\n')
    serve.chmod(0o755)
    print(f"\n==> DONE → {args.out_dir / 'index.html'}")
    print(f"    preview: bash {serve}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
