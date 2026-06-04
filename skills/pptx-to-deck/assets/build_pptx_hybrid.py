#!/usr/bin/env python3
"""build_pptx_hybrid.py · PPTX → 高保真「可编辑」HTML（混合管线）

与 build_pptx.py（纯代码重建）互补，这条管线解决"装饰细节对不齐"：

  · 背景层 = LibreOffice headless 渲染「剥光文字」的幻灯片 → 像素级保真装饰
            （渐变 / glow / 金焰 / 阴影 / 自由曲线 / 照片 / 图表，全部原样）
  · 文字层 = 从「带文字」的 LibreOffice 渲染里抽出每一行的真实 bbox + 字体 + 颜色，
            作为结构化、可编辑、纯色文字钉位叠加（用真实字体，nowrap + scaleX 贴合）

一条命令：
  build_pptx_hybrid.py <in.pptx> <out-dir> [--renderer SKILL] [--title T] [--soffice PATH]

依赖：LibreOffice（soffice headless，渲染改写过的 pptx 不像 PowerPoint 那样失败）、
      PyMuPDF（fitz）、python-pptx。产物 = out-dir/{index.html, deck.json, bg/}。

注意：这是 PPTX 复刻管线；背景装饰像素级但不可编辑，前景文字结构化可编辑。
"""
from __future__ import annotations
import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

CANVAS_W, CANVAS_H = 1920, 1080
_SUBSET_PREFIX = re.compile(r"^[A-Z]{6}\+")          # 子集字体名前缀 ABCDEF+
_FONT_FALLBACK = '"PingFang SC", "Microsoft YaHei", sans-serif'


# ── locate tools ────────────────────────────────────────────────────────────
def _find_soffice(explicit: str | None) -> str:
    if explicit:
        return explicit
    for c in ("/Applications/LibreOffice.app/Contents/MacOS/soffice",
              shutil.which("soffice") or "", shutil.which("libreoffice") or ""):
        if c and Path(c).exists():
            return c
    sys.exit("ERROR: LibreOffice (soffice) not found. Install: brew install --cask "
             "libreoffice. 这条混合管线必须用 LibreOffice 渲染无字背景。")


def _default_renderer() -> Path:
    """feishu-deck-h5 渲染后端定位:pptx-to-deck 是顶层 skill,feishu-deck-h5 通常
    是兄弟 skill 目录。优先兄弟 <skills>/feishu-deck-h5/,兼容旧的内嵌祖父布局,
    最后回退已注册的 ~/.claude/skills/feishu-deck-h5(symlink)。"""
    skills_dir = Path(__file__).resolve().parent.parent.parent  # <skills>/
    for cand in (skills_dir / "feishu-deck-h5",                 # 兄弟(新布局)
                 skills_dir,                                     # 旧内嵌
                 Path.home() / ".claude/skills/feishu-deck-h5"):  # 注册 symlink
        if (cand / "deck-json/render-deck.py").is_file():
            return cand
    return Path.home() / ".claude/skills/feishu-deck-h5"


# ── step 1 · strip text + content pictures (background = decoration only) ────
def _is_content_pic(sh, slide_w: int, slide_h: int) -> bool:
    """A PICTURE that is foreground content (photo / mockup / logo) — to be
    REMOVED from the background and re-overlaid from the original blob (lossless,
    sharp, no ghost). Full-bleed pictures (decoration backdrop) and vector
    WMF/EMF (browser can't render the raw blob) are kept in the LibreOffice bg."""
    if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return False
    try:
        if (sh.image.ext or "").lower() in ("wmf", "emf"):
            return False
    except Exception:
        return False
    full_bleed = sh.width >= slide_w * 0.95 and sh.height >= slide_h * 0.95
    return not full_bleed


def strip_for_bg(src: Path, dst: Path) -> None:
    """Clear all text AND remove content pictures → the LibreOffice render is a
    decoration-only background (gradients/glow/shapes/full-bleed images/WMF),
    leaving holes where photos go. Originals are overlaid later — so no ghost."""
    prs = Presentation(str(src))
    sw, sh_ = prs.slide_width, prs.slide_height

    def process(shapes):
        drop = []
        for sh in shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                process(sh.shapes)
            elif sh.has_text_frame:
                for r in sh.text_frame._txBody.iter(qn("a:r")):
                    t = r.find(qn("a:t"))
                    if t is not None:
                        t.text = ""
            if _is_content_pic(sh, sw, sh_):
                drop.append(sh)
        for sh in drop:
            sh._element.getparent().remove(sh._element)
    for slide in prs.slides:
        process(slide.shapes)
    prs.save(str(dst))


# ── original content images (overlaid lossless on the decoration bg) ─────────
def extract_original_images(pptx: Path, work: Path, out: Path, renderer: Path,
                            scale: float) -> dict[int, list]:
    """Run build_pptx (no render) → original image blobs + geometry + srcRect
    crop. Return {1-based page: [image elements]} for the SAME content pictures
    strip_for_bg removed (non-full-bleed, non-WMF). Blobs are copied into
    out/input, downsized to scale× their display box (retina-sharp, sane size)."""
    bp = Path(__file__).resolve().parent / "build_pptx.py"   # 同目录兄弟管线
    bp_out = work / "bp"
    try:
        subprocess.run([sys.executable, str(bp), str(pptx), str(bp_out),
                        "--renderer", str(renderer), "--no-render"],
                       check=True, capture_output=True, text=True, timeout=600)
        deck = json.loads((bp_out / "deck.json").read_text(encoding="utf-8"))
    except Exception as e:
        print(f"    (build_pptx 原图抽取失败,跳过原图叠加: {e})")
        return {}
    (out / "input").mkdir(parents=True, exist_ok=True)
    res: dict[int, list] = {}
    for idx, s in enumerate(deck.get("slides", []), 1):
        els = []
        for e in s.get("data", {}).get("elements", []):
            if e.get("type") != "image":
                continue
            src = e.get("src", "")
            if not src.startswith("input/"):
                continue
            ext = src.rsplit(".", 1)[-1].lower()
            full_bleed = e.get("w", 0) >= 1880 and e.get("h", 0) >= 1040
            if ext in ("wmf", "emf") or full_bleed:
                continue                    # 留给 LibreOffice 背景
            srcf = bp_out / src
            if not srcf.is_file():
                continue
            saved = _copy_optimized(srcf, out / "input" / srcf.name,
                                    int(e["w"] * scale), int(e["h"] * scale))
            ne = dict(e); ne["src"] = "input/" + saved.name
            els.append(ne)
        if els:
            res[idx] = els
    return res


def _copy_optimized(src: Path, dst: Path, max_w: int, max_h: int) -> Path:
    """Copy an image, downsized so it's at most max_w×max_h (keeps original
    detail at display×scale, trims oversized blobs). Opaque → JPEG q90.
    Returns the actually-written path (extension may change to .jpg)."""
    try:
        im = Image.open(src)
        if im.size[0] > max_w or im.size[1] > max_h:
            im.thumbnail((max(1, max_w), max(1, max_h)), Image.LANCZOS)
        alpha = "A" in im.mode or (im.mode == "P" and "transparency" in im.info)
        if alpha:
            im.save(dst)
            return dst
        dst = dst.with_suffix(".jpg")
        im.convert("RGB").save(dst, "JPEG", quality=90, optimize=True)
        return dst
    except Exception:
        shutil.copy2(src, dst)
        return dst


# ── step 2 · LibreOffice → PDF (hidden slides included → 1:1 page mapping) ────
def soffice_to_pdf(soffice: str, pptx: Path, outdir: Path) -> Path:
    outdir.mkdir(parents=True, exist_ok=True)
    flt = ('pdf:impress_pdf_Export:'
           '{"ExportHiddenSlides":{"type":"boolean","value":"true"}}')
    subprocess.run([soffice, "--headless", "--convert-to", flt,
                    "--outdir", str(outdir), str(pptx)],
                   check=True, capture_output=True, text=True, timeout=900)
    pdf = outdir / (pptx.stem + ".pdf")
    if not pdf.is_file():
        sys.exit(f"ERROR: LibreOffice did not produce {pdf}")
    return pdf


# ── step 3 · rasterize the no-text PDF → background JPEGs ─────────────────────
def rasterize(pdf: Path, bgdir: Path, scale: float = 2.0) -> int:
    """Render each page to scale×1920px wide (default 2× for retina sharpness;
    the 16:9 deck displays at 1920 CSS px so a 2× source stays crisp when the
    browser downscales it on a HiDPI screen). q92 JPEG."""
    bgdir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(pdf))
    for i, page in enumerate(doc, 1):
        m = (CANVAS_W * scale) / page.rect.width
        page.get_pixmap(matrix=fitz.Matrix(m, m)).pil_save(
            str(bgdir / f"page-{i:03d}.jpg"), format="JPEG", quality=92)
    n = len(doc)
    doc.close()
    return n


_WEIGHT = re.compile(r"[-,]?\s*(Light|Regular|Medium|DemiBold?|SemiBold?|Bold|"
                     r"Heavy|Black|Thin|ExtraLight|Normal|Book)$", re.I)


# ── step 4 · text layer from the WITH-text render (real bbox/color) ──────────
def deck_fonts(pptx: Path) -> str:
    """CSS font-family from the PPTX's OWN dominant typefaces (ea=CJK first, then
    latin). CRITICAL: do NOT trust the font names PyMuPDF reads off the
    LibreOffice render — LibreOffice substitutes missing fonts (Hiragino /
    LiberationSerif / Arial…) and those look wrong in the browser. The real
    typeface (e.g. FZLanTingHeiPro_GB18030, usually installed) renders correctly
    and matches widths. Weight-stripped base name appended so 'FZLanTing… Medium'
    also matches the installed family."""
    from collections import Counter
    ea, lat = Counter(), Counter()

    def walk(shapes):
        for sh in shapes:
            if sh.shape_type == 6:
                walk(sh.shapes)
            elif sh.has_text_frame:
                for rPr in sh.text_frame._txBody.iter(qn("a:rPr")):
                    for tag, cnt in (("a:ea", ea), ("a:latin", lat)):
                        node = rPr.find(qn(tag))
                        tf = node.get("typeface") if node is not None else None
                        if tf and not tf.startswith("+"):
                            cnt[tf] += 1
    try:
        for s in Presentation(str(pptx)).slides:
            walk(s.shapes)
    except Exception:
        pass
    fams: list = []
    for cnt in (ea, lat):                       # ea(CJK) first → consistent
        if cnt:
            tf = cnt.most_common(1)[0][0]
            for cand in (tf, _WEIGHT.sub("", tf)):
                if cand and cand not in fams:
                    fams.append(cand)
    css = ", ".join(f'"{f}"' for f in fams)
    return (css + ", " if css else "") + _FONT_FALLBACK


def line_text_elements(pdf: Path, font: str) -> dict[int, list]:
    """{1-based page: [text elements]} — one element per rendered LINE, runs are
    its spans (so mixed CJK/Latin/size flows on one line, not split-positioned)."""
    doc = fitz.open(str(pdf))
    out: dict[int, list] = {}
    for pno, page in enumerate(doc, 1):
        m = CANVAS_W / page.rect.width
        els, eid = [], 0
        for blk in page.get_text("dict")["blocks"]:
            for ln in blk.get("lines", []):
                spans = [s for s in ln.get("spans", []) if s.get("text")]
                if not any(s["text"].strip() for s in spans):
                    continue
                x0, y0, x1, y1 = ln["bbox"]
                runs = [{"text": s["text"],
                         "color": "#%06X" % (s.get("color", 0xFFFFFF) & 0xFFFFFF),
                         "size": round(s.get("size", 12) * m, 1),
                         "font": font}        # 用 PPTX 真实主字体,不用 LibreOffice 替换名
                        for s in spans]
                eid += 1
                els.append({"id": f"t{pno}_{eid}", "type": "text",
                            "x": round(x0 * m, 1), "y": round(y0 * m, 1),
                            "w": round((x1 - x0) * m, 1), "h": round((y1 - y0) * m, 1),
                            "anchor": "middle", "insets": [0, 0, 0, 0], "runs": runs})
        out[pno] = els
    doc.close()
    return out


# ── step 5 · assemble deck.json ──────────────────────────────────────────────
def build_deck(n_pages: int, texts: dict[int, list], images: dict[int, list],
               title: str) -> dict:
    slides = []
    for i in range(1, n_pages + 1):
        bg = {"id": f"bg{i}", "type": "image", "src": f"bg/page-{i:03d}.jpg",
              "x": 0, "y": 0, "w": CANVAS_W, "h": CANVAS_H}
        # z-order: 装饰背景 → 原始内容图(无损) → 可编辑文字
        elements = [bg] + images.get(i, []) + texts.get(i, [])
        slides.append({
            "key": f"slide-{i:03d}", "layout": "canvas",
            "screen_label": f"{i:02d}", "lifted": f"pptx-hybrid:{title}#{i}",
            "data": {"canvas_w": CANVAS_W, "canvas_h": CANVAS_H, "source_page": i,
                     "elements": elements},
        })
    return {"version": "1.0", "deck": {"title": title}, "slides": slides}


# ── step 7 · post-process: letterbox bg + nowrap + scaleX fit + hide progress ─
def post_process(out_dir: Path, deck: dict) -> None:
    rules = []
    for s in deck["slides"]:
        bg = next((e for e in s["data"]["elements"]
                   if e.get("type") == "image" and e["src"].startswith("bg/")), None)
        if bg:
            rules.append('.deck[data-mode="present"] .slide-frame:has(> '
                         '.slide[data-slide-key="%s"]){background:#000 url("%s") '
                         'center/cover no-repeat}' % (s["key"], bg["src"]))
    inject = (
        '<style>.deck-ui .deck-progress{display:none!important}\n'
        '.slide .el.tb .tb-inner{white-space:nowrap}\n' + "\n".join(rules) + '</style>\n'
        '<script>\n'
        'function fitText(){var c=document.querySelector(".slide-frame.is-current");'
        'if(!c)return;c.querySelectorAll(".el.tb").forEach(function(el){'
        'var i=el.querySelector(".tb-inner");if(!i||i.dataset.fit)return;'
        'i.style.whiteSpace="nowrap";i.style.display="inline-block";'
        'i.style.transformOrigin="left center";var n=i.getBoundingClientRect().width,'
        'b=el.getBoundingClientRect().width;if(n>1&&b>1)'
        'i.style.transform="scaleX("+(b/n)+")";i.dataset.fit="1";});}\n'
        'addEventListener("load",function(){setTimeout(fitText,300)});'
        'addEventListener("hashchange",function(){setTimeout(fitText,300)});\n'
        '</script>')
    p = out_dir / "index.html"
    h = p.read_text(encoding="utf-8").replace("</head>", inject + "</head>", 1)
    p.write_text(h, encoding="utf-8")


# ── step 8 · 资源自包含打包（产物可移动/打包/发送，不依赖技能目录） ──────────────
_REF_RX = re.compile(r'(?:href|src)=["\'](?P<a>[^"\']+)["\']|url\((?P<u>[^)]+)\)')
_CSS_URL_RX = re.compile(r'url\((?P<u>[^)]+)\)')
_EXT_PREFIX = ("data:", "http://", "https://", "//", "#", "mailto:")


def make_portable(out: Path, renderer: Path) -> None:
    """使产物自包含：把 index.html 里所有「解析后落在技能目录内」的框架引用
    （CSS / JS / 字体 / 图）拷进 out/assets/ 并改写为 deck 本地相对引用，再跟随每个
    已拷 CSS 内部的 url()（feishu-deck.css 里的 lark logo 等）一并拷齐。结果：整个
    out/ 夹拷走 / 打包 / 发给别人都不断（否则 ../../../../…/skills/… 路径一离开仓库
    就 404）。

    判据 = 「引用解析后是否落在 renderer 技能目录内」：是→框架资源，拷进来；否→
    deck 本地（bg/、input/）或外链，原样不动。**路径无关**——不依赖 copy-assets.py
    的 runs/<ts>/output 规范布局与 (\\.\\./)+skills/feishu-deck-h5/ 正则（那套只在
    repo 根 runs/ 下才匹配），无论输出到哪都能打包成自包含。"""
    renderer = Path(os.path.realpath(renderer))
    out_real = Path(os.path.realpath(out))
    assets = out / "assets"
    html_path = out / "index.html"
    html = html_path.read_text(encoding="utf-8")
    copied: dict[str, Path] = {}   # origin-realpath → 本地 target

    def _resolve(base_dir: Path, ref: str):
        """ref 相对 base_dir 解析；返回「应被打包进来的框架文件」Path，否则 None。
        判据 = 解析后是真实文件、**在技能 renderer 内、但不在 out 内**。
        · 在 out 内（bg/、input/、或二次运行时已打包的 assets/）→ 本就 deck 本地，不动；
          ⚠ 关键：产物常输出到 skills/feishu-deck-h5/runs/<name>，bg/input 也在 renderer
          之下，单看 renderer 会误判，必须再排除 out 内。
        · 技能外（外链 / data: / 系统字体等）→ 不碰。"""
        r = ref.strip().strip("\"'")
        if not r or r.startswith(_EXT_PREFIX):
            return None
        origin = Path(os.path.realpath(base_dir / r.split("?")[0].split("#")[0]))
        if not origin.is_file():
            return None
        if origin.is_relative_to(out_real):     # 已在 deck 内 → 本地，不动
            return None
        if not origin.is_relative_to(renderer):  # 技能外 → 不碰
            return None
        return origin                            # 技能内 + out 外 = 框架资源

    def _local_target(origin: Path) -> Path:
        """技能内文件 → out/assets/ 下的落点。assets/ 子树去掉首段（不要双层
        assets/assets/）；其它子树（deck-json/templates 等）保留相对路径。"""
        rel = origin.relative_to(renderer)
        parts = rel.parts
        if parts and parts[0] == "assets":
            rel = Path(*parts[1:])
        return assets / rel

    def _copy_in(origin: Path) -> Path:
        target = _local_target(origin)
        target.parent.mkdir(parents=True, exist_ok=True)
        if not target.exists() or target.stat().st_size != origin.stat().st_size:
            shutil.copy2(origin, target)
        copied[str(origin)] = target
        return target

    def repl_html(m):
        ref = m.group("a") or m.group("u")
        origin = _resolve(out, ref)
        if origin is None:
            return m.group(0)
        new = os.path.relpath(_copy_in(origin), out).replace(os.sep, "/")
        return m.group(0).replace(ref, new, 1)

    new_html = _REF_RX.sub(repl_html, html)
    if new_html != html:
        html_path.write_text(new_html, encoding="utf-8")

    # pass 2：已拷 CSS 内部 url() — 相对其 ORIGIN 解析，拷进同位并按 TARGET 改写
    n_fw = len(copied)
    for origin_str, target in list(copied.items()):
        if target.suffix.lower() != ".css":
            continue
        origin = Path(origin_str)
        css = target.read_text(encoding="utf-8")

        def repl_css(m):
            o = _resolve(origin.parent, m.group("u"))
            if o is None:
                return m.group(0)
            t = _copy_in(o)
            return 'url("%s")' % os.path.relpath(t, target.parent).replace(os.sep, "/")

        new_css = _CSS_URL_RX.sub(repl_css, css)
        if new_css != css:
            target.write_text(new_css, encoding="utf-8")
    print(f"    便携打包完成：{len(copied)} 个框架资源已拷入 assets/（产物自包含可移动）")


# ── main ─────────────────────────────────────────────────────────────────────
def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description="PPTX → 高保真可编辑 HTML (混合管线)")
    ap.add_argument("pptx", type=Path)
    ap.add_argument("out_dir", type=Path)
    ap.add_argument("--renderer", type=Path, default=_default_renderer())
    ap.add_argument("--title", default=None)
    ap.add_argument("--soffice", default=None)
    ap.add_argument("--scale", type=float, default=2.0,
                    help="背景渲染倍率(默认2×=Retina清晰;3=更锐但更大)")
    args = ap.parse_args(argv)

    soffice = _find_soffice(args.soffice)
    title = args.title or args.pptx.stem
    out = args.out_dir
    out.mkdir(parents=True, exist_ok=True)
    work = Path(tempfile.mkdtemp(prefix="pptx-hybrid-"))
    try:
        print("==> [1/8] 剥光文字 + 剥内容图 (python-pptx) → 装饰背景源")
        stripped = work / "notext.pptx"
        strip_for_bg(args.pptx, stripped)

        print("==> [2/8] LibreOffice 渲染装饰背景 → PDF")
        notext_pdf = soffice_to_pdf(soffice, stripped, work / "notext")
        print(f"==> [3/8] 栅格化装饰背景 → bg/*.jpg ({args.scale:g}× = {int(CANVAS_W*args.scale)}px宽)")
        n_pages = rasterize(notext_pdf, out / "bg", args.scale)
        print(f"    {n_pages} 页背景")

        print("==> [4/8] 抽取 PPT 原始图片(无损,原分辨率) → 叠加层")
        images = extract_original_images(args.pptx, work, out, args.renderer, args.scale)
        print(f"    {sum(len(v) for v in images.values())} 张原图")

        print("==> [5/8] LibreOffice 渲染带字版 → 真实文字位置/颜色（字体取自源PPTX）")
        # 字体取自 PPTX 源文件的真实 typeface, 绝不用 PyMuPDF 从 LibreOffice 渲染
        # 读出的字体名(那是 LibreOffice 的替换字体 Hiragino/LiberationSerif/Arial,
        # 喂给浏览器会渲得又杂又怪)。这个坑只能这样根治。
        font = deck_fonts(args.pptx)
        print(f"    文字层字体: {font}")
        withtext_pdf = soffice_to_pdf(soffice, args.pptx, work / "withtext")
        texts = line_text_elements(withtext_pdf, font)

        print("==> [6/8] 组装 deck.json")
        deck = build_deck(n_pages, texts, images, title)
        (out / "deck.json").write_text(
            json.dumps(deck, ensure_ascii=False, indent=1), encoding="utf-8")

        print("==> [7/8] 渲染 HTML (render-deck.py)")
        render = args.renderer / "deck-json/render-deck.py"
        subprocess.run([sys.executable, str(render), str(out / "deck.json"), str(out),
                        "--skip-copy-assets", "--skip-validate-html"],
                       check=True, capture_output=True, text=True, timeout=600)

        print("==> [8/8] 资源自包含打包 (框架→assets/) + 前端增强")
        make_portable(out, args.renderer)
        post_process(out, deck)

        n_text = sum(len(v) for v in texts.values())
        print(f"\n==> DONE → {out / 'index.html'}")
        print(f"    {n_pages} 页 · {n_text} 个可编辑文字元素 · 背景像素级保真 · 自包含可移动")
    finally:
        shutil.rmtree(work, ignore_errors=True)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
