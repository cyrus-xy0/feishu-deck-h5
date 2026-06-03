#!/usr/bin/env python3
"""
build_pptx.py · PPTX → feishu-deck-h5 deck.json (layout:"canvas")

Native PowerPoint (.pptx) parser backend. Walks every slide with python-pptx
and emits each slide as a `layout:"canvas"` deck.json slide: a STRUCTURED list
of absolutely-positioned, typed, id'd elements (`data.elements[]`) — NOT an HTML
blob and NOT a screenshot. The user's own feishu-deck-h5 renderer
(deck-json/render-deck.py) turns elements[] into positioned HTML, and
sync-index-to-deck.py round-trips edits back into elements[] by id.

This is the PPTX side of the DECKJSON-UNIFIED-INTERMEDIATE-SPEC (§2/§3): one
intermediate layer = deck.json, structured JSON for PPTX (no HTML坨, no images),
one edit loop = edit→sync→deck.json.

.pptx is open OOXML, so python-pptx reads exact element geometry (EMU), text
runs (font / size / color / b / i), embedded images, and groups directly,
cross-platform — no AppleScript, no Keynote, no LibreOffice.

Element mapping (each → {id, type, x, y, w, h, ...} ; px on canvas 1920×1080):
  · TEXT_BOX / PLACEHOLDER-with-text → {type:"text", runs:[{text,bold,color,
        size}], anchor, insets}   — CLEAN STRUCTURED editable content
  · PICTURE                         → {type:"image", src:"input/<file>"}
        (embedded media extracted to input/ — that's original content, kept)
  · MEDIA (video)                   → {type:"image", src:<poster>} or a poster
        placeholder shape
  · AUTO_SHAPE / FREEFORM / LINE    → {type:"shape", kind, fill|gradient,
        border, radius, svg, style} — appearance fields so it renders right
  · GROUP                           → FLATTENED: children composed through the
        group transform and emitted as top-level elements (no group wrapper)

Un-reconstructable (live chart / SmartArt diagram / OLE object): the WHOLE slide
becomes a placeholder — {layout:"canvas", data:{placeholder:true, source_page:N,
elements:[]}} — and N is collected. At the end a report line prints:
  `unreconstructed slides: [N, ...]`  (empty list if none).

NO SCREENSHOTS: the old whole-slide / per-element raster fallback is RETIRED.
--raster / --full-raster are accepted as deprecated no-ops (they never produce
image-replica slides). Embedded PICTUREs stay as image elements (original
content, fine). This is "尽可能还原" (best-effort reconstruction), NOT pixel-
perfect: gradients are carried as CSS, freeform/line as inline SVG paths.

Usage:
  build_pptx.py <in.pptx> <out-dir>
       [--renderer DIR]   feishu-deck-h5 skill root
                          (default: ~/.claude/skills/feishu-deck-h5)
       [--limit N]        only first N slides
       [--raster]         DEPRECATED no-op (rasterization retired)
       [--full-raster]    DEPRECATED no-op (rasterization retired)
       [--inline]         single-file render output (base64-inline everything)
       [--title TEXT]     deck title
       [--no-render]      emit deck.json + assets only, skip the HTML render
"""
from __future__ import annotations

import argparse
import html as html_lib
import json
import re
import subprocess
import sys
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn

# ── canvas ──────────────────────────────────────────────────────────────────
SLIDE_W, SLIDE_H = 1920, 1080
EMU_PER_PT = 12700

# NOTE: font names use SINGLE quotes — these strings are emitted inside a
# double-quoted style="" attribute, and nested double quotes would truncate
# the attribute (silently dropping every declaration after font-family,
# e.g. color/font-weight → text falls back to the inherited default).
DEFAULT_FONT_STACK = (
    "'PingFang SC','Microsoft YaHei','Source Han Sans SC',"
    "'Helvetica Neue',Arial,sans-serif"
)


# ── affine transform (EMU→EMU), used to flatten nested groups ────────────────
class Xf:
    """Maps a shape's container-local EMU coords to absolute slide EMU."""

    def __init__(self, ox=0.0, oy=0.0, sx=1.0, sy=1.0):
        self.ox, self.oy, self.sx, self.sy = ox, oy, sx, sy

    def x(self, v):  return self.ox + v * self.sx
    def y(self, v):  return self.oy + v * self.sy
    def w(self, v):  return v * self.sx
    def h(self, v):  return v * self.sy

    def enter_group(self, grp) -> "Xf":
        """Return the child transform for a GroupShape `grp`."""
        # group's own box, mapped to absolute slide EMU through *this* xf
        gx, gy = self.x(grp.left), self.y(grp.top)
        gw, gh = self.w(grp.width), self.h(grp.height)
        # child coordinate space (a:chOff / a:chExt)
        cox = coy = 0.0
        cex = grp.width or 1
        cey = grp.height or 1
        xfrm = grp._element.find(qn("p:grpSpPr"))
        if xfrm is not None:
            xfrm = xfrm.find(qn("a:xfrm"))
        if xfrm is not None:
            ch_off = xfrm.find(qn("a:chOff"))
            ch_ext = xfrm.find(qn("a:chExt"))
            if ch_off is not None:
                cox = float(ch_off.get("x", 0)); coy = float(ch_off.get("y", 0))
            if ch_ext is not None:
                cex = float(ch_ext.get("cx", 1)) or 1
                cey = float(ch_ext.get("cy", 1)) or 1
        sx = gw / cex
        sy = gh / cey
        return Xf(gx - cox * sx, gy - coy * sy, sx, sy)


# ── theme-color resolution ────────────────────────────────────────────────────
# Populated once per presentation in main(): MSO_THEME_COLOR member → "#RRGGBB".
# Without this, scheme-colored text (the deck author's "tx1/lt1/accent1") has no
# .rgb and falls back to a guessed default → light text on dark cards goes dark.
_THEME: dict = {}
# parallel map keyed by the raw schemeClr val string ("accent1","lt1","tx1"…),
# for resolving colors inside gradient stops / XML we read by hand.
_THEME_BY_NAME: dict = {}

# clrScheme child tag → theme-color enum members it backs (incl. the tx/bg
# aliases that the master's clrMap maps to dk/lt by default).
_SCHEME_TAGS = [
    ("dk1", [MSO_THEME_COLOR.DARK_1, MSO_THEME_COLOR.TEXT_1]),
    ("lt1", [MSO_THEME_COLOR.LIGHT_1, MSO_THEME_COLOR.BACKGROUND_1]),
    ("dk2", [MSO_THEME_COLOR.DARK_2, MSO_THEME_COLOR.TEXT_2]),
    ("lt2", [MSO_THEME_COLOR.LIGHT_2, MSO_THEME_COLOR.BACKGROUND_2]),
    ("accent1", [MSO_THEME_COLOR.ACCENT_1]),
    ("accent2", [MSO_THEME_COLOR.ACCENT_2]),
    ("accent3", [MSO_THEME_COLOR.ACCENT_3]),
    ("accent4", [MSO_THEME_COLOR.ACCENT_4]),
    ("accent5", [MSO_THEME_COLOR.ACCENT_5]),
    ("accent6", [MSO_THEME_COLOR.ACCENT_6]),
    ("hlink", [MSO_THEME_COLOR.HYPERLINK]),
    ("folHlink", [MSO_THEME_COLOR.FOLLOWED_HYPERLINK]),
]


def _clr_to_hex(node) -> Optional[str]:
    """Read a hex from a clrScheme child (<a:srgbClr> or <a:sysClr lastClr>)."""
    srgb = node.find(qn("a:srgbClr"))
    if srgb is not None and srgb.get("val"):
        return f"#{srgb.get('val')}"
    sysc = node.find(qn("a:sysClr"))
    if sysc is not None and sysc.get("lastClr"):
        return f"#{sysc.get('lastClr')}"
    return None


def build_theme_map(prs) -> dict:
    """Map theme-color members → hex from the first master's theme clrScheme."""
    from lxml import etree
    out: dict = {}
    try:
        master = prs.slide_masters[0]
        theme_part = None
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return out
        # the theme is a generic Part (no ._element) — parse its blob
        root = etree.fromstring(theme_part.blob)
        clr_scheme = root.find(f".//{qn('a:clrScheme')}")
        if clr_scheme is None:
            return out
        for tag, members in _SCHEME_TAGS:
            node = clr_scheme.find(qn(f"a:{tag}"))
            if node is None:
                continue
            hexc = _clr_to_hex(node)
            if hexc:
                for m in members:
                    out[m] = hexc
                _THEME_BY_NAME[tag] = hexc
        # tx/bg aliases default-mapped to dk/lt (master clrMap default)
        for alias, base in (("tx1", "dk1"), ("bg1", "lt1"),
                            ("tx2", "dk2"), ("bg2", "lt2")):
            if base in _THEME_BY_NAME:
                _THEME_BY_NAME[alias] = _THEME_BY_NAME[base]
    except Exception:
        pass
    return out


def _scheme_name_hex(val: str) -> Optional[str]:
    """Resolve a raw schemeClr val ('accent1','lt1','tx1','bg1'…) to hex."""
    return _THEME_BY_NAME.get(val)


def _xml_color_hex(node) -> Optional[str]:
    """Resolve a DrawingML color child (<a:srgbClr>/<a:schemeClr>/<a:sysClr>)."""
    srgb = node.find(qn("a:srgbClr"))
    if srgb is not None and srgb.get("val"):
        return f"#{srgb.get('val')}"
    sysc = node.find(qn("a:sysClr"))
    if sysc is not None and sysc.get("lastClr"):
        return f"#{sysc.get('lastClr')}"
    sch = node.find(qn("a:schemeClr"))
    if sch is not None and sch.get("val"):
        return _scheme_name_hex(sch.get("val"))
    return None


def gradient_css(shape) -> Optional[str]:
    """Parse a shape's <a:gradFill> into a CSS linear-gradient(), or None."""
    try:
        sp_pr = shape._element.find(qn("p:spPr"))
        if sp_pr is None:
            return None
        grad = sp_pr.find(qn("a:gradFill"))
        if grad is None:
            return None
        gs_lst = grad.find(qn("a:gsLst"))
        if gs_lst is None:
            return None
        stops = []
        for gs in gs_lst.findall(qn("a:gs")):
            col = _xml_color_hex(gs)
            if col:
                stops.append((int(gs.get("pos", "0")) / 1000.0, col))
        if len(stops) < 2:
            return None
        # angle: PPT <a:lin ang> is 1/60000 deg clockwise from East (3 o'clock).
        # CSS linear-gradient 0deg = "to top"; convert via (ppt + 90) mod 360.
        ang_css = "to bottom"
        lin = grad.find(qn("a:lin"))
        if lin is not None and lin.get("ang") is not None:
            ang = (int(lin.get("ang")) / 60000.0 + 90) % 360
            ang_css = f"{ang:.0f}deg"
        stop_str = ", ".join(f"{c} {p:.0f}%" for p, c in stops)
        return f"linear-gradient({ang_css}, {stop_str})"
    except Exception:
        return None


# ── color helpers ─────────────────────────────────────────────────────────────
def rgb_hex(color) -> Optional[str]:
    """Best-effort RGB hex from a python-pptx ColorFormat. None if unresolved.
    Resolves scheme/theme colors via the presentation theme map (_THEME)."""
    try:
        if color is None or color.type is None:
            return None
        # explicit RGB → direct (getattr swallows the AttributeError that
        # .rgb raises for non-RGB color types)
        rgb = getattr(color, "rgb", None)
        if rgb is not None:
            return f"#{str(rgb)}"
        # scheme/theme color → look up resolved hex
        tc = getattr(color, "theme_color", None)
        if tc is not None and tc in _THEME:
            return _THEME[tc]
    except Exception:
        return None
    return None


def fill_hex(fill) -> Optional[str]:
    """Safe solid-fill hex. Returns None for no-fill / gradient / picture /
    pattern / unresolved — never raises (python-pptx raises on .fore_color
    for non-solid fills)."""
    try:
        from pptx.enum.dml import MSO_FILL
        if fill is not None and fill.type == MSO_FILL.SOLID:
            return rgb_hex(fill.fore_color)
    except Exception:
        return None
    return None


def emu(v) -> float:
    return float(v) if v is not None else 0.0


# ── geometry → px ─────────────────────────────────────────────────────────────
class Canvas:
    """Uniform scale + center offset from slide EMU into 1920×1080 px."""

    def __init__(self, slide_w_emu: int, slide_h_emu: int):
        self.scale = min(SLIDE_W / slide_w_emu, SLIDE_H / slide_h_emu)
        self.off_x = (SLIDE_W - slide_w_emu * self.scale) / 2
        self.off_y = (SLIDE_H - slide_h_emu * self.scale) / 2

    def px_left(self, emu_x):  return emu_x * self.scale + self.off_x
    def px_top(self, emu_y):   return emu_y * self.scale + self.off_y
    def px_size(self, emu_d):  return emu_d * self.scale
    def pt_to_px(self, pt):    return pt * EMU_PER_PT * self.scale


# ── slide background ──────────────────────────────────────────────────────────
def slide_bg(slide, prs) -> str:
    """Resolve a background color from slide → layout → master. Default #FFF."""
    from pptx.enum.dml import MSO_FILL
    for src in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        try:
            bg = src.background
            if bg.fill.type == MSO_FILL.SOLID:
                hexc = rgb_hex(bg.fill.fore_color)
                if hexc:
                    return hexc
        except Exception:
            continue
    return "#FFFFFF"


# ── element id ─────────────────────────────────────────────────────────────────
class IdGen:
    """Stable per-slide element ids: e{slide}_{n} (1-based slide, 0-based n)."""

    def __init__(self, slide_no: int):
        self.slide_no = slide_no
        self.n = 0

    def next(self) -> str:
        eid = f"e{self.slide_no}_{self.n}"
        self.n += 1
        return eid


# ── geometry → element x/y/w/h (px on 1920×1080) ───────────────────────────────
def el_geom(cv: Canvas, xf: Xf, shape) -> dict:
    """Absolute px geometry for one shape, group transform composed in."""
    return {
        "x": round(cv.px_left(xf.x(emu(shape.left))), 1),
        "y": round(cv.px_top(xf.y(emu(shape.top))), 1),
        "w": round(cv.px_size(xf.w(emu(shape.width))), 1),
        "h": round(cv.px_size(xf.h(emu(shape.height))), 1),
    }


def _rotation_style(shape) -> Optional[str]:
    """Non-trivial rotation → a CSS `transform` string for the shape `style`
    escape-hatch (text/image stay clean; only shapes carry style)."""
    rot = getattr(shape, "rotation", 0) or 0
    if rot and abs(rot) > 0.01:
        return f"transform:rotate({rot:.2f}deg)"
    return None


# ── text element ───────────────────────────────────────────────────────────────
_ANCHOR = {
    MSO_ANCHOR.TOP: "top", MSO_ANCHOR.MIDDLE: "middle",
    MSO_ANCHOR.BOTTOM: "bottom",
}


def _bullet_char(p) -> Optional[str]:
    """Bullet marker for a paragraph, else None (prepended to the run text so it
    stays structured + editable rather than becoming layout markup)."""
    pPr = p._p.find(qn("a:pPr"))
    if pPr is None:
        return None
    if pPr.find(qn("a:buNone")) is not None:
        return None
    bc = pPr.find(qn("a:buChar"))
    if bc is not None:
        return bc.get("char", "•")
    if pPr.find(qn("a:buAutoNum")) is not None:
        return "•"  # numbered → approximate; v0 doesn't track the counter
    return None


def _emu_to_px_pt(cv: Canvas, sz) -> Optional[float]:
    """python-pptx font.size (EMU/Length) → px on the canvas."""
    if sz is None:
        return None
    try:
        return round(cv.px_size(int(sz)), 1)
    except Exception:
        return None


def text_runs(shape, cv: Canvas) -> list:
    """Flatten a text frame's paragraphs into a single runs[] list. Paragraph
    breaks and soft <a:br/> become a run whose text starts with '\\n' so the
    renderer's _esc_br turns it into a <br>. Bullets are prepended to the first
    run of the bulleted paragraph as text (stays structured + editable)."""
    from pptx.text.text import _Run
    tf = shape.text_frame
    runs: list = []
    first_para = True
    for p in tf.paragraphs:
        para_runs: list = []
        for child in p._p:
            tag = child.tag
            if tag == qn("a:r"):
                r = _Run(child, p)
                txt = r.text
                if not txt:
                    continue
                run = {"text": txt}
                if r.font.bold:
                    run["bold"] = True
                col = rgb_hex(r.font.color)
                if col:
                    run["color"] = col
                size = _emu_to_px_pt(cv, r.font.size)
                if size is not None:
                    run["size"] = size
                para_runs.append(run)
            elif tag == qn("a:br"):
                para_runs.append({"text": "\n"})
            elif tag == qn("a:fld"):  # slide number / date field → cached text
                t = child.find(qn("a:t"))
                if t is not None and t.text:
                    para_runs.append({"text": t.text})
        if not para_runs:
            continue
        # bullet → prepend to the paragraph's first real run
        bullet = _bullet_char(p)
        if bullet:
            para_runs[0]["text"] = f"{bullet} {para_runs[0]['text']}"
        # paragraph separator: a leading "\n" on the first run of every
        # paragraph after the first (renders as <br> between paragraphs).
        if not first_para:
            para_runs[0]["text"] = "\n" + para_runs[0]["text"]
        first_para = False
        runs.extend(para_runs)
    return runs


def text_element(shape, cv: Canvas, xf: Xf, eid: str) -> Optional[dict]:
    runs = text_runs(shape, cv)
    if not runs:
        return None
    el = {"id": eid, "type": "text"}
    el.update(el_geom(cv, xf, shape))
    try:
        anchor = _ANCHOR.get(shape.text_frame.vertical_anchor)
        if anchor:
            el["anchor"] = anchor
    except Exception:
        pass
    # text frame insets (left/right/top/bottom) in px
    try:
        tf = shape.text_frame
        insets = [tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom]
        if any(i is not None for i in insets):
            el["insets"] = [round(cv.px_size(int(i)), 1) if i is not None else 0
                            for i in insets]
    except Exception:
        pass
    el["runs"] = runs
    return el


# ── shape appearance (fill / gradient / border / radius / preset / svg) ─────────
def _preset_geom_name(shape) -> Optional[str]:
    """OOXML preset geometry name (<a:prstGeom prst="...">), e.g. rect /
    roundRect / ellipse. None for custom/freeform geometry."""
    try:
        sp_pr = shape._element.find(qn("p:spPr"))
        if sp_pr is None:
            return None
        prst = sp_pr.find(qn("a:prstGeom"))
        if prst is not None:
            return prst.get("prst")
    except Exception:
        pass
    return None


def _border_obj(shape, cv: Canvas) -> Optional[dict]:
    try:
        ln = shape.line
        if ln.width is not None and int(ln.width) > 0:
            color = rgb_hex(ln.color) or "#888888"
            width = max(1.0, round(int(ln.width) / EMU_PER_PT, 1))
            return {"color": color, "width": width}
    except Exception:
        pass
    return None


def _radius_px(shape, cv: Canvas) -> Optional[float]:
    """Corner radius (px) for rounded rects; 50% sentinel handled by emit via
    ellipse. Approximate: PPT roundRect adj defaults ~16.7% of the short side."""
    name = _preset_geom_name(shape)
    if name == "roundRect":
        try:
            w = cv.px_size(emu(shape.width))
            h = cv.px_size(emu(shape.height))
            return round(min(w, h) * 0.16, 1)
        except Exception:
            return 14.0
    return None


def shape_appearance(shape, cv: Canvas) -> dict:
    """Fill / gradient / border / radius / kind for a non-freeform AUTO_SHAPE.
    Returns only the keys that resolve (clean elements carry no empty fields)."""
    out: dict = {}
    name = _preset_geom_name(shape)
    if name:
        out["kind"] = name
    # fill: solid hex OR gradient CSS
    try:
        from pptx.enum.dml import MSO_FILL
        fill = shape.fill
        if fill.type == MSO_FILL.SOLID:
            hexc = rgb_hex(fill.fore_color)
            if hexc:
                out["fill"] = hexc
        elif fill.type == MSO_FILL.GRADIENT:
            grad = gradient_css(shape)
            if grad:
                out["gradient"] = grad
    except Exception:
        pass
    border = _border_obj(shape, cv)
    if border:
        out["border"] = border
    # radius: explicit roundRect → computed px; ellipse → 50% via style
    if name == "ellipse":
        out["style"] = "border-radius:50%"
    else:
        radius = _radius_px(shape, cv)
        if radius is not None:
            out["radius"] = radius
    rot = _rotation_style(shape)
    if rot:
        out["style"] = (out.get("style") + ";" + rot) if out.get("style") else rot
    return out


# ── freeform / custGeom / line → inline SVG (normalized 0..100 box) ─────────────
def _custgeom_svg(shape) -> Optional[str]:
    """Parse <a:custGeom> path commands into an SVG <path>, normalized to a
    0..100 viewBox (renderer uses preserveAspectRatio:none, so coords map to %
    of the element box). Returns the inner SVG markup, or None."""
    try:
        sp_pr = shape._element.find(qn("p:spPr"))
        if sp_pr is None:
            return None
        cust = sp_pr.find(qn("a:custGeom"))
        if cust is None:
            return None
        path_lst = cust.find(qn("a:pathLst"))
        if path_lst is None:
            return None
        # fill / stroke from the shape
        fill = fill_hex(shape.fill) or "none"
        stroke, sw = "none", 0.0
        try:
            ln = shape.line
            if ln.width is not None and int(ln.width) > 0:
                stroke = rgb_hex(ln.color) or "#888888"
                sw = max(0.5, int(ln.width) / EMU_PER_PT)
        except Exception:
            pass
        if fill == "none" and stroke == "none":
            fill = "#888888"  # at least show the silhouette
        d_parts: list[str] = []
        for path in path_lst.findall(qn("a:path")):
            pw = float(path.get("w") or 0) or 1.0
            ph = float(path.get("h") or 0) or 1.0
            def nx(v): return round(float(v) / pw * 100, 2)
            def ny(v): return round(float(v) / ph * 100, 2)
            for cmd in path:
                t = cmd.tag
                if t == qn("a:moveTo"):
                    pt = cmd.find(qn("a:pt"))
                    d_parts.append(f"M{nx(pt.get('x'))},{ny(pt.get('y'))}")
                elif t == qn("a:lnTo"):
                    pt = cmd.find(qn("a:pt"))
                    d_parts.append(f"L{nx(pt.get('x'))},{ny(pt.get('y'))}")
                elif t == qn("a:cubicBezTo"):
                    pts = cmd.findall(qn("a:pt"))
                    if len(pts) == 3:
                        d_parts.append("C" + " ".join(
                            f"{nx(p.get('x'))},{ny(p.get('y'))}" for p in pts))
                elif t == qn("a:quadBezTo"):
                    pts = cmd.findall(qn("a:pt"))
                    if len(pts) == 2:
                        d_parts.append("Q" + " ".join(
                            f"{nx(p.get('x'))},{ny(p.get('y'))}" for p in pts))
                elif t == qn("a:close"):
                    d_parts.append("Z")
        if not d_parts:
            return None
        d = " ".join(d_parts)
        sw_attr = f' stroke="{stroke}" stroke-width="{sw:.1f}"' if stroke != "none" else ""
        return (f'<path d="{html_lib.escape(d, quote=True)}" '
                f'fill="{fill}"{sw_attr} vector-effect="non-scaling-stroke"/>')
    except Exception:
        return None


def line_svg(shape, cv: Canvas) -> str:
    """A LINE / connector as an SVG line in the normalized 0..100 box, honoring
    flipH / flipV so diagonals point the right way."""
    color, sw = None, 1.0
    try:
        color = rgb_hex(shape.line.color)
        if shape.line.width:
            sw = max(1.0, int(shape.line.width) / EMU_PER_PT)
    except Exception:
        pass
    color = color or "#8895AA"
    xfrm = shape._element.find(f".//{qn('a:xfrm')}")
    flip_h = xfrm is not None and xfrm.get("flipH") == "1"
    flip_v = xfrm is not None and xfrm.get("flipV") == "1"
    x1, y1, x2, y2 = (100 if flip_h else 0), (100 if flip_v else 0), \
                     (0 if flip_h else 100), (0 if flip_v else 100)
    return (f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" '
            f'stroke="{color}" stroke-width="{sw:.1f}" '
            f'vector-effect="non-scaling-stroke"/>')


# ── per-shape → element dict(s) ────────────────────────────────────────────────
def emit_shape(shape, cv: Canvas, xf: Xf, idgen: IdGen, input_dir: Path,
               img_counter: list) -> list:
    """Return list of element dicts for one shape (GROUP → many, flattened)."""
    out: list = []
    st = shape.shape_type

    # GROUP → recurse with the composed child transform; emit children as
    # top-level elements (no wrapper).
    if st == MSO_SHAPE_TYPE.GROUP:
        child_xf = xf.enter_group(shape)
        for child in shape.shapes:
            out += emit_shape(child, cv, child_xf, idgen, input_dir, img_counter)
        return out

    # PICTURE → image element; blob extracted to input/ (original content).
    if st == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = shape.image
            img_counter[0] += 1
            ext = img.ext or "png"
            fname = f"img-{img_counter[0]:03d}.{ext}"
            (input_dir / fname).write_bytes(img.blob)
            el = {"id": idgen.next(), "type": "image", "src": f"input/{fname}"}
            el.update(el_geom(cv, xf, shape))
            out.append(el)
        except Exception:
            pass  # unreadable embedded image → silently drop (rare)
        return out

    # MEDIA (video/audio) → poster image if present, else a poster placeholder.
    if st == MSO_SHAPE_TYPE.MEDIA:
        try:
            img = shape.image  # poster/cover frame
            img_counter[0] += 1
            ext = img.ext or "png"
            fname = f"media-{img_counter[0]:03d}.{ext}"
            (input_dir / fname).write_bytes(img.blob)
            el = {"id": idgen.next(), "type": "image", "src": f"input/{fname}"}
            el.update(el_geom(cv, xf, shape))
            out.append(el)
        except Exception:
            el = {"id": idgen.next(), "type": "shape", "fill": "#0B0F18",
                  "kind": "rect"}
            el.update(el_geom(cv, xf, shape))
            out.append(el)
        return out

    # LINE / connector → shape with inline SVG.
    if st == MSO_SHAPE_TYPE.LINE:
        el = {"id": idgen.next(), "type": "shape", "svg": line_svg(shape, cv)}
        el.update(el_geom(cv, xf, shape))
        out.append(el)
        return out

    # FREEFORM (custGeom) → shape with inline SVG path.
    if st == MSO_SHAPE_TYPE.FREEFORM:
        svg = _custgeom_svg(shape)
        if svg:
            el = {"id": idgen.next(), "type": "shape", "svg": svg}
            el.update(el_geom(cv, xf, shape))
            out.append(el)
        # also emit any text the freeform carries (rare)
        if shape.has_text_frame and shape.text_frame.text.strip():
            tel = text_element(shape, cv, xf, idgen.next())
            if tel:
                out.append(tel)
        return out

    # TABLE → flatten cells to text + cell-fill shapes (no table layout type in
    # canvas; tables are rare and best reconstructed as positioned cells).
    if shape.has_table:
        out += _emit_table(shape, cv, xf, idgen)
        return out

    # AUTO_SHAPE / TEXT_BOX / PLACEHOLDER.
    # A non-textbox auto-shape with a custGeom (preset failed) → SVG fallback.
    is_textbox = st in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.PLACEHOLDER)
    has_text = shape.has_text_frame and shape.text_frame.text.strip()

    # shape appearance (skip for pure text boxes/placeholders — those are clean
    # text only; a placeholder with a real fill still gets a backing shape).
    if not is_textbox:
        appearance = shape_appearance(shape, cv)
        # nothing reproducible (e.g. picture/pattern fill, no border) and no
        # text → try custGeom SVG, else drop (decorative, un-reconstructable).
        has_visible = bool(appearance.get("fill") or appearance.get("gradient")
                           or appearance.get("border") or appearance.get("svg"))
        if not has_visible:
            svg = _custgeom_svg(shape)
            if svg:
                appearance["svg"] = svg
                has_visible = True
        if has_visible:
            sel = {"id": idgen.next(), "type": "shape"}
            sel.update(el_geom(cv, xf, shape))
            sel.update(appearance)
            out.append(sel)

    # text → a separate clean text element on top (works for shapes-with-text
    # too: the backing shape + a text box at the same geometry).
    if has_text:
        tel = text_element(shape, cv, xf, idgen.next())
        if tel:
            out.append(tel)
    return out


def _emit_table(shape, cv: Canvas, xf: Xf, idgen: IdGen) -> list:
    """Flatten a PPT table into positioned cell shapes + cell text elements.
    Column/row sizes from the table's gridCol / row heights (EMU)."""
    out: list = []
    tbl = shape.table
    base_x = emu(shape.left)
    base_y = emu(shape.top)
    col_w = [emu(c.width) for c in tbl.columns]
    row_h = [emu(r.height) for r in tbl.rows]
    y = base_y
    for ri, row in enumerate(tbl.rows):
        x = base_x
        for ci, cell in enumerate(row.cells):
            cw = col_w[ci] if ci < len(col_w) else 0
            ch = row_h[ri] if ri < len(row_h) else 0
            gx = round(cv.px_left(xf.x(x)), 1)
            gy = round(cv.px_top(xf.y(y)), 1)
            gw = round(cv.px_size(xf.w(cw)), 1)
            gh = round(cv.px_size(xf.h(ch)), 1)
            # cell backing shape (fill + a light border so the grid reads)
            cfill = fill_hex(cell.fill)
            sel = {"id": idgen.next(), "type": "shape",
                   "x": gx, "y": gy, "w": gw, "h": gh, "kind": "rect",
                   "border": {"color": "#D0D5DD", "width": 1}}
            if cfill:
                sel["fill"] = cfill
            out.append(sel)
            # cell text
            txt = cell.text.strip()
            if txt:
                out.append({"id": idgen.next(), "type": "text",
                            "x": gx, "y": gy, "w": gw, "h": gh,
                            "anchor": "middle",
                            "insets": [8, 8, 6, 6],
                            "runs": [{"text": txt}]})
            x += cw
        y += row_h[ri] if ri < len(row_h) else 0
    return out


# ── un-reconstructable detection ───────────────────────────────────────────────
_DIAGRAM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"


def _shape_is_hard(shape) -> bool:
    """True if a shape is structurally un-reconstructable: a live chart,
    a SmartArt diagram, or an OLE object."""
    st = shape.shape_type
    if st in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
              MSO_SHAPE_TYPE.LINKED_OLE_OBJECT):
        return True
    try:
        if getattr(shape, "has_chart", False):
            return True
    except Exception:
        pass
    if st == MSO_SHAPE_TYPE.DIAGRAM:
        return True
    # SmartArt = a graphicFrame whose graphicData references the diagram ns
    try:
        gd = shape._element.find(f".//{qn('a:graphicData')}")
        if gd is not None and _DIAGRAM_NS in (gd.get("uri") or ""):
            return True
    except Exception:
        pass
    return False


def _slide_is_hard(slide) -> bool:
    """A slide is un-reconstructable if ANY of its shapes is hard (chart/
    SmartArt/OLE). Groups are recursed."""
    def walk(shapes) -> bool:
        for sh in shapes:
            if _shape_is_hard(sh):
                return True
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                if walk(sh.shapes):
                    return True
        return False
    return walk(slide.shapes)


# placeholder prompt text on layout/master that must NOT render as content
_PROMPT_RE = re.compile(
    r"单击此处|單擊此處|请输入|請輸入|点击.*编辑|按一下.*編輯|母版|母版文本样式|"
    r"Click to edit|XXXX|公司职称|人名$|^\s*第.*页\s*$"
)


def _emit_bg_picture(container, input_dir: Path, idgen: IdGen,
                     tagname: str, img_counter: list) -> Optional[dict]:
    """A slide/layout/master `<p:bg>` PICTURE fill → a full-bleed image element
    (the designed cover background lives in the layout/master, not the slide's
    own shapes)."""
    try:
        el = container._element
        bg = el.find(f".//{qn('p:bg')}")
        if bg is None:
            return None
        blip = bg.find(f".//{qn('a:blip')}")
        if blip is None:
            return None
        embed = blip.get(qn("r:embed"))
        if not embed:
            return None
        part = container.part.related_part(embed)
        ext = str(part.partname).rsplit(".", 1)[-1] or "png"
        img_counter[0] += 1
        fname = f"bg-{tagname}-{img_counter[0]:03d}.{ext}"
        (input_dir / fname).write_bytes(part.blob)
        return {"id": idgen.next(), "type": "image", "src": f"input/{fname}",
                "x": 0, "y": 0, "w": SLIDE_W, "h": SLIDE_H}
    except Exception:
        return None


def emit_template_elements(container, cv: Canvas, idgen: IdGen, input_dir: Path,
                           img_counter: list) -> list:
    """Layout/master decorative shapes (background pictures, logos, footers)
    UNDER the slide content. Skip placeholders (prompts) and prompt text."""
    out: list = []
    xf = Xf()
    for sh in container.shapes:
        try:
            if sh.is_placeholder:
                continue
            if sh.has_text_frame and _PROMPT_RE.search(sh.text_frame.text or ""):
                continue
        except Exception:
            pass
        out += emit_shape(sh, cv, xf, idgen, input_dir, img_counter)
    return out


# ── imported-provenance marker ──────────────────────────────────────────────────
def _lifted_marker(deck_stem: str, slide_no: int) -> str:
    """Provenance string for a PPTX-imported canvas slide, reusing the schema's
    existing `lifted` field (value = source ref). Format `pptx:<stem>#<N>` mirrors
    lift-slides' `<deck>#<index>` convention. The renderer emits this as
    `data-lifted=…`, and validate.py then downgrades this slide's
    CONTENT-AUTHORING violations (R05 ellipsis / R-KEY positional / R-LANG latin
    leaf / R10 palette / R06 / R-WHITE-TEXT / R-VIS-TIER / R-VIS-BODY-FLOOR) from
    error → WARNING: the content is faithfully CARRIED from the source PPTX, not
    hand-authored to the deck's content standards, so it is surfaced (a human
    chooses to fix) rather than blocking. STRUCTURAL / GEOMETRY rules (R-DOM /
    R-OVERFLOW / R-VIS-CARD-OVERFLOW / R-CSSVAR / duplicate-key R-KEY) stay
    full-severity error — a real overflow is a real failure regardless of import."""
    return f"pptx:{deck_stem}#{slide_no}"


# ── slide → canvas elements ────────────────────────────────────────────────────
def compose_slide(slide, idx: int, cv: Canvas, prs, input_dir: Path,
                  img_counter: list, deck_stem: str) -> dict:
    """Build one `layout:"canvas"` slide dict (data.elements[])."""
    slide_no = idx + 1
    idgen = IdGen(slide_no)

    # background as a backing shape so the slide isn't transparent.
    elements: list = []
    bg = slide_bg(slide, prs)
    if bg and bg.upper() != "#FFFFFF":
        elements.append({"id": idgen.next(), "type": "shape", "kind": "rect",
                         "x": 0, "y": 0, "w": SLIDE_W, "h": SLIDE_H,
                         "fill": bg})

    # TEMPLATE LAYER (under content): layout/master background image + decor.
    layout = slide.slide_layout
    master = layout.slide_master
    for cont, tg in ((master, "master"), (layout, "layout")):
        bgimg = _emit_bg_picture(cont, input_dir, idgen, tg, img_counter)
        if bgimg:
            elements.append(bgimg)
        elements += emit_template_elements(cont, cv, idgen, input_dir, img_counter)

    # the slide's own shapes
    xf = Xf()
    for shape in slide.shapes:
        elements += emit_shape(shape, cv, xf, idgen, input_dir, img_counter)

    return {
        "key": f"slide-{slide_no:03d}",
        "layout": "canvas",
        "screen_label": f"{slide_no:02d}",
        # IMPORTED PROVENANCE: this slide is verbatim-carried from a foreign
        # PPTX, not hand-authored — mark it `lifted` so the validator downgrades
        # CONTENT-AUTHORING rules to warnings (see _lifted_marker).
        "lifted": _lifted_marker(deck_stem, slide_no),
        "data": {
            "canvas_w": SLIDE_W,
            "canvas_h": SLIDE_H,
            "source_page": slide_no,
            "elements": elements,
        },
    }


def _default_renderer() -> Path:
    """Locate the feishu-deck-h5 renderer. This is a SUB-skill nested inside
    the feishu-deck-h5 skill (.../feishu-deck-h5/pptx-to-html/assets/), so the
    renderer is the grandparent dir. Fall back to ~/.claude/skills/feishu-deck-h5."""
    parent_skill = Path(__file__).resolve().parent.parent.parent  # → feishu-deck-h5/
    if (parent_skill / "deck-json/render-deck.py").is_file():
        return parent_skill
    return Path.home() / ".claude/skills/feishu-deck-h5"


# ── main ──────────────────────────────────────────────────────────────────────
def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description="PPTX → feishu-deck-h5 canvas deck.json")
    ap.add_argument("pptx", type=Path)
    ap.add_argument("out_dir", type=Path)
    ap.add_argument("--renderer", type=Path, default=_default_renderer())
    ap.add_argument("--limit", type=int, default=0)
    ap.add_argument("--raster", action="store_true",
                    help="DEPRECATED no-op (rasterization retired)")
    ap.add_argument("--full-raster", action="store_true",
                    help="DEPRECATED no-op (rasterization retired)")
    ap.add_argument("--inline", action="store_true")
    ap.add_argument("--no-render", action="store_true",
                    help="emit deck.json + assets only, skip the HTML render")
    ap.add_argument("--title", default=None)
    args = ap.parse_args(argv)

    if args.raster or args.full_raster:
        print("==> NOTE: --raster / --full-raster are retired no-ops "
              "(no screenshots; embedded pictures stay as image elements).",
              file=sys.stderr)

    if not args.pptx.is_file():
        print(f"ERROR: pptx not found: {args.pptx}", file=sys.stderr)
        return 1
    render_script = args.renderer / "deck-json/render-deck.py"
    if not args.no_render and not render_script.is_file():
        print(f"ERROR: render-deck.py not at {render_script}", file=sys.stderr)
        return 1

    args.out_dir.mkdir(parents=True, exist_ok=True)
    # image assets live in input/ (real, scannable paths in elements[].src).
    input_dir = args.out_dir / "input"
    input_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(args.pptx))
    global _THEME
    _THEME = build_theme_map(prs)
    print(f"==> theme colors resolved: {len(_THEME)} members")
    cv = Canvas(prs.slide_width, prs.slide_height)

    print(f"==> {args.pptx.name}: {len(prs.slides)} slides, "
          f"canvas scale={cv.scale:.4f}")
    deck_slides = []
    unreconstructed: list[int] = []
    img_counter = [0]
    deck_stem = (args.title or args.pptx.stem)
    for idx, slide in enumerate(prs.slides):
        if args.limit and idx >= args.limit:
            break
        slide_no = idx + 1
        if _slide_is_hard(slide):
            # live chart / SmartArt / OLE → placeholder + collect page number.
            unreconstructed.append(slide_no)
            deck_slides.append({
                "key": f"slide-{slide_no:03d}",
                "layout": "canvas",
                "screen_label": f"{slide_no:02d}",
                # imported provenance (see _lifted_marker): a placeholder is
                # still verbatim-carried foreign content, mark it lifted too.
                "lifted": _lifted_marker(deck_stem, slide_no),
                "data": {"placeholder": True, "source_page": slide_no,
                         "elements": []},
            })
            print(f"  · slide {slide_no:2d}  ←  UN-RECONSTRUCTABLE "
                  f"(chart/SmartArt/OLE) → placeholder")
            continue
        sd = compose_slide(slide, idx, cv, prs, input_dir, img_counter, deck_stem)
        deck_slides.append(sd)
        n = len(sd["data"]["elements"])
        print(f"  · slide {slide_no:2d}  ←  {len(slide.shapes)} shapes  "
              f"→ {n} elements")

    deck = {
        "version": "1.0",
        "deck": {
            "title": args.title or args.pptx.stem,
            "language": "zh-only",
            "mode": "rewrite",
        },
        "slides": deck_slides,
    }
    deck_path = args.out_dir / "deck.json"
    deck_path.write_text(json.dumps(deck, ensure_ascii=False, indent=2))
    print(f"==> wrote {deck_path}  ({len(deck_slides)} slides)")
    # the contract report line (empty list if none).
    print(f"unreconstructed slides: {unreconstructed}")

    if args.no_render:
        print(f"\n==> DONE (no render) → {deck_path}")
        return 0

    # Keep the DeckJSON schema validation gate ON (handoff contract: validate
    # before render). Skip the post-render HTML validator + the content/story
    # schema-fit refusal — neither applies to canvas slides.
    cmd = [sys.executable, str(render_script), str(deck_path), str(args.out_dir),
           "--skip-validate-html", "--skip-fit-check"]
    if args.inline:
        cmd.append("--inline")
    print(f"==> rendering via {render_script}")
    res = subprocess.run(cmd, capture_output=True, text=True)
    if res.returncode != 0:
        print("RENDER FAILED:", file=sys.stderr)
        print(res.stdout, file=sys.stderr)
        print(res.stderr, file=sys.stderr)
        return 1
    for line in (res.stdout.strip().splitlines() or [""])[-6:]:
        print(f"    {line}")

    # serve helper
    serve = args.out_dir / "serve.sh"
    serve.write_text(
        '#!/usr/bin/env bash\nset -e\ncd "$(dirname "$0")"\n'
        'PORT="${1:-8765}"\n'
        'echo "==> http://localhost:$PORT/index.html"\n'
        'python3 -m http.server "$PORT"\n')
    serve.chmod(0o755)
    print(f"\n==> DONE → {args.out_dir / 'index.html'}")
    print(f"    preview: bash {serve}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
