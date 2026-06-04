#!/usr/bin/env python3
"""extract_media.py — pull video/GIF overlays out of a .pptx: their slide,
absolute position (group-transform-aware), and the media file itself. Then probe
each with ffprobe to decide GIF-style (no audio → autoplay loop) vs real video
(has audio → click-to-play).

  python3 extract_media.py deck.pptx --out ./work/media
Outputs:
  <out>/files/...                 extracted media (mp4/mov/gif) + gif→mp4, mov→mp4 remux
  <out>/media-raw.json            {"<slide>":[{file,left,top,width,height,gif}], ...}
After you upload <out>/files to a host, rewrite `file`→`url` to feed make_manifest --media.

Requires: python-pptx; ffmpeg/ffprobe on PATH (for audio probe + remux/convert).
"""
import argparse
import json
import os
import re
import shutil
import subprocess
from pptx import Presentation
from pptx.oxml.ns import qn

CW, CH = 1920, 1080
R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def group_tf(group):
    el = group._element
    xf = el.find('.//' + qn('a:xfrm'))
    chOff, chExt = xf.find(qn('a:chOff')), xf.find(qn('a:chExt'))
    cx0, cy0 = int(chOff.get('x')), int(chOff.get('y'))
    ccx, ccy = int(chExt.get('cx')), int(chExt.get('cy'))
    gx, gy, gcx, gcy = group.left, group.top, group.width, group.height
    sx = gcx / ccx if ccx else 1
    sy = gcy / ccy if ccy else 1
    return lambda x, y, w, h: (gx + (x - cx0) * sx, gy + (y - cy0) * sy, w * sx, h * sy)


def has_audio(path):
    try:
        r = subprocess.run(['ffprobe', '-v', 'error', '-select_streams', 'a',
                            '-show_entries', 'stream=codec_type', '-of', 'csv=p=0', path],
                           capture_output=True, text=True)
        return 'audio' in r.stdout
    except Exception:
        return True  # if unsure, treat as video (click-to-play)


def walk(shapes, sw, sh_, out, part_map, tfs=None):
    if tfs is None:
        tfs = []
    for sh in shapes:
        if sh.shape_type == 6:
            walk(sh.shapes, sw, sh_, out, part_map, tfs + [group_tf(sh)])
            continue
        el = sh._element
        # video: <p:nvPicPr><p:nvPr><a:videoFile r:link=.. or p14:media r:embed=..>
        is_video = el.find('.//' + qn('a:videoFile')) is not None or b'media' in (el.xml.encode() if hasattr(el, 'xml') else b'')
        blip = el.find('.//' + qn('a:blip'))
        is_gif = False
        rid = None
        if sh.shape_type == 13 and blip is not None:
            rid = blip.get(qn('r:embed'))
            tgt = part_map.get(rid, '')
            if tgt.lower().endswith('.gif'):
                is_gif = True
        vf = el.find('.//' + qn('a:videoFile'))
        if vf is not None:
            rid = vf.get(qn('r:link')) or rid
        if not (is_video and vf is not None) and not is_gif:
            continue
        x, y, w, h = sh.left, sh.top, sh.width, sh.height
        if x is None:
            continue
        for tf in reversed(tfs):
            x, y, w, h = tf(x, y, w, h)
        out.append({'rid': rid, 'is_gif': is_gif,
                    'left': round(x * CW / sw, 2), 'top': round(y * CH / sh_, 2),
                    'width': round(w * CW / sw, 2), 'height': round(h * CH / sh_, 2)})


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('pptx')
    ap.add_argument('--out', default='./work/media')
    a = ap.parse_args()
    files_dir = os.path.join(a.out, 'files')
    os.makedirs(files_dir, exist_ok=True)

    import zipfile
    zf = zipfile.ZipFile(a.pptx)
    p = Presentation(a.pptx)
    sw, sh_ = p.slide_width, p.slide_height
    result = {}

    for i, slide in enumerate(p.slides):
        part = slide.part
        rel_map = {rid: rel.target_ref for rid, rel in part.rels.items()}
        items = []
        walk(slide.shapes, sw, sh_, items, rel_map)
        if not items:
            continue
        outlist = []
        for it in items:
            tgt = rel_map.get(it['rid'], '')
            if not tgt:
                continue
            src = 'ppt/' + tgt.replace('../', '')
            base = os.path.basename(src)
            dst = os.path.join(files_dir, base)
            try:
                with zf.open(src) as f, open(dst, 'wb') as o:
                    o.write(f.read())
            except KeyError:
                continue
            # normalize: mov→mp4 remux, gif→mp4 (so all overlays are <video>)
            final = base
            stem, ext = os.path.splitext(base)
            if ext.lower() == '.mov':
                final = stem + '.mp4'
                subprocess.run(['ffmpeg', '-y', '-i', dst, '-c', 'copy', '-movflags', '+faststart',
                                os.path.join(files_dir, final)], capture_output=True)
            elif ext.lower() == '.gif':
                final = stem + '.mp4'
                subprocess.run(['ffmpeg', '-y', '-i', dst, '-movflags', '+faststart', '-pix_fmt', 'yuv420p',
                                '-vf', 'scale=trunc(iw/2)*2:trunc(ih/2)*2', os.path.join(files_dir, final)],
                               capture_output=True)
            gif = it['is_gif'] or not has_audio(os.path.join(files_dir, final))
            outlist.append({'file': final, 'left': it['left'], 'top': it['top'],
                            'width': it['width'], 'height': it['height'], 'gif': gif})
        if outlist:
            result[str(i + 1)] = outlist

    json.dump(result, open(os.path.join(a.out, 'media-raw.json'), 'w'), ensure_ascii=False, indent=2)
    n = sum(len(v) for v in result.values())
    print(f'media-raw.json: {n} overlays across {len(result)} slides; files in {files_dir}')
    print('next: upload files/ to a host, set each "file"->"url", feed make_manifest --media')


if __name__ == '__main__':
    main()
