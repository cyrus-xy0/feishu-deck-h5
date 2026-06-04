#!/usr/bin/env python3
"""extract.py — pull editable text structure from a .pptx and write a text-free copy.

Outputs:
  <out>/texts.json        per-slide list of text frames:
                            {left,top,width,height (px on 1920×1080),
                             insets:[l,r,t,b], anchor, paras:[{align,size,color,bold,text,spc_before}]}
  <out>/text-stripped.pptx every run's text emptied, all decoration kept (for the
                            no-text background render).

Key fidelity details handled:
  • Group transforms (chOff/chExt → off/ext) so grouped shapes get absolute coords.
  • Resolved paragraph alignment from lstStyle/lvlNpPr (python-pptx's
    para.alignment misses the lstStyle default where many decks set 居中).
  • Per-paragraph font size (mixed title/body sizes preserved).
  • bodyPr lIns/rIns/tIns/bIns → padding.
  • z-order: a text frame fully covered by a later picture is SKIPPED (it was
    hidden behind an image in the deck; don't float it on top).
  • Optional --fill-sizes <keynote-or-prior.json>: fill None font sizes by
    matching frame text (placeholder titles inherit size from the master, which
    python-pptx can't resolve).

Usage:
  python3 extract.py deck.pptx --out ./work [--fill-sizes prior.json]
"""
import argparse
import json
import os
import re
import sys
from pptx import Presentation
from pptx.oxml.ns import qn

CANVAS_W, CANVAS_H = 1920, 1080
ALGN = {'l': 'LEFT', 'ctr': 'CENTER', 'r': 'RIGHT', 'just': 'JUSTIFY', 'dist': 'JUSTIFY'}


def group_tf(group):
    el = group._element
    xf = el.find('.//' + qn('a:xfrm'))
    chOff, chExt = xf.find(qn('a:chOff')), xf.find(qn('a:chExt'))
    cx0, cy0 = int(chOff.get('x')), int(chOff.get('y'))
    ccx, ccy = int(chExt.get('cx')), int(chExt.get('cy'))
    gx, gy, gcx, gcy = group.left, group.top, group.width, group.height
    sx = gcx / ccx if ccx else 1
    sy = gcy / ccy if ccy else 1
    return lambda x, y, w, h: (gx + (x - cx0) * sx, gy + (y - cy0) * sy, w * sx, h * sy)


def resolve_align(para, lstStyle):
    pPr = para._p.find(qn('a:pPr'))
    algn = pPr.get('algn') if pPr is not None else None
    lvl = int(pPr.get('lvl', 0)) if (pPr is not None and pPr.get('lvl')) else 0
    if algn:
        return ALGN.get(algn)
    if lstStyle is not None:
        for tag in (f'a:lvl{lvl+1}pPr', 'a:lvl1pPr'):
            e = lstStyle.find(qn(tag))
            if e is not None and e.get('algn'):
                return ALGN.get(e.get('algn'))
    return None


def spc_before(para):
    pPr = para._p.find(qn('a:pPr'))
    if pPr is None:
        return None
    sb = pPr.find(qn('a:spcBef'))
    if sb is None:
        return None
    pts = sb.find(qn('a:spcPts'))
    return int(pts.get('val')) / 100.0 if (pts is not None and pts.get('val')) else None


def run_color(run):
    try:
        c = run.font.color
        if c and c.type is not None and c.rgb is not None:
            return f'#{c.rgb}'
    except Exception:
        pass
    return None


def collect(shapes, leaves, tfs=None, z=None):
    if tfs is None:
        tfs = []
    if z is None:
        z = [0]
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            collect(sh.shapes, leaves, tfs + [group_tf(sh)], z)
            continue
        x, y, w, h = sh.left, sh.top, sh.width, sh.height
        if x is None:
            z[0] += 1
            continue
        for tf in reversed(tfs):
            x, y, w, h = tf(x, y, w, h)
        zi = z[0]
        z[0] += 1
        leaves.append({'z': zi, 'x': x, 'y': y, 'w': w, 'h': h,
                       'is_pic': sh.shape_type == 13,
                       'has_text': bool(sh.has_text_frame and sh.text_frame.text.strip()),
                       'shape': sh})


def extract_frame(L, sw, sh_):
    sh = L['shape']
    tf = sh.text_frame
    lstStyle = tf._txBody.find(qn('a:lstStyle'))
    bodyPr = tf._txBody.find(qn('a:bodyPr'))

    def ins(attr, d):
        v = bodyPr.get(attr) if bodyPr is not None else None
        return int(v) if v is not None else d

    sx, sy = CANVAS_W / sw, CANVAS_H / sh_
    lIns, rIns = ins('lIns', 91440) * sx, ins('rIns', 91440) * sx
    tIns, bIns = ins('tIns', 45720) * sy, ins('bIns', 45720) * sy
    anchor = None
    try:
        va = tf.vertical_anchor
        anchor = str(va) if va else None
    except Exception:
        pass
    pxpt = CANVAS_W / (sw / 914400 * 72)
    paras = []
    for para in tf.paragraphs:
        txt = ''.join(r.text for r in para.runs)
        if not para.runs:
            paras.append({'align': None, 'size': None, 'color': None, 'bold': None, 'text': '', 'spc_before': None})
            continue
        size = color = bold = None
        for r in para.runs:
            if r.font.size is not None and size is None:
                size = r.font.size.pt
            if color is None:
                c = run_color(r)
                if c:
                    color = c
            if bold is None and r.font.bold is not None:
                bold = r.font.bold
        sb = spc_before(para)
        paras.append({'align': resolve_align(para, lstStyle),
                      'size': round(size * pxpt, 1) if size else None,
                      'color': color, 'bold': bold, 'text': txt,
                      'spc_before': round(sb * pxpt, 1) if sb else None})
    while paras and not paras[-1]['text']:
        paras.pop()
    if not paras:
        return None
    return {'left': round(L['x'] * sx, 2), 'top': round(L['y'] * sy, 2),
            'width': round(L['w'] * sx, 2), 'height': round(L['h'] * sy, 2),
            'insets': [round(lIns, 1), round(rIns, 1), round(tIns, 1), round(bIns, 1)],
            'anchor': anchor, 'paras': paras}


def covered(tL, leaves):
    tx0, ty0, tx1, ty1 = tL['x'], tL['y'], tL['x'] + tL['w'], tL['y'] + tL['h']
    for o in leaves:
        if not o['is_pic'] or o['z'] <= tL['z']:
            continue
        ox0, oy0, ox1, oy1 = o['x'], o['y'], o['x'] + o['w'], o['y'] + o['h']
        tol = min(tL['w'], tL['h']) * 0.12 + 1
        if ox0 <= tx0 + tol and oy0 <= ty0 + tol and ox1 >= tx1 - tol and oy1 >= ty1 - tol:
            return True
    return False


def clear_text(shapes):
    n = 0
    for sh in shapes:
        if sh.shape_type == 6:
            n += clear_text(sh.shapes)
        elif sh.has_text_frame and sh.text_frame.text.strip():
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ''
            n += 1
    return n


def norm(t):
    return re.sub(r'\s+', '', t or '')


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('pptx')
    ap.add_argument('--out', default='./work')
    ap.add_argument('--fill-sizes', default=None,
                    help='prior JSON ({slide:[{text,font_size}]}) to fill None sizes by text match')
    a = ap.parse_args()
    os.makedirs(a.out, exist_ok=True)

    p = Presentation(a.pptx)
    sw, sh_ = p.slide_width, p.slide_height
    result = {}
    for i, slide in enumerate(p.slides):
        leaves = []
        collect(slide.shapes, leaves)
        frames = []
        for L in leaves:
            if not L['has_text'] or covered(L, leaves):
                continue
            fr = extract_frame(L, sw, sh_)
            if fr:
                frames.append(fr)
        result[i + 1] = frames

    if a.fill_sizes and os.path.exists(a.fill_sizes):
        kn = json.load(open(a.fill_sizes))
        km = {}
        for sn, shps in kn.items():
            km[str(sn)] = [(norm(s['text']), s['font_size']) for s in shps if s.get('text') and s.get('font_size')]
        for sn, frames in result.items():
            kl = km.get(str(sn), [])
            for fr in frames:
                ft = norm(''.join(p2['text'] for p2 in fr['paras']))
                for pp in fr['paras']:
                    if pp['size'] is None and pp['text'].strip():
                        nt = norm(pp['text'])
                        cand = next((fs for knt, fs in kl if knt == nt or
                                     (len(nt) >= 4 and (knt.startswith(nt[:8]) or nt.startswith(knt[:8])))), None)
                        if cand is None:
                            cand = next((fs for knt, fs in kl if knt == ft or
                                         (len(ft) >= 6 and knt.startswith(ft[:10]))), None)
                        if cand:
                            pp['size'] = cand

    json.dump(result, open(os.path.join(a.out, 'texts.json'), 'w'), ensure_ascii=False, indent=2)

    p2 = Presentation(a.pptx)
    cleared = sum(clear_text(s.shapes) for s in p2.slides)
    p2.save(os.path.join(a.out, 'text-stripped.pptx'))

    nf = sum(len(v) for v in result.values())
    print(f'texts.json: {nf} frames across {len(result)} slides')
    print(f'text-stripped.pptx: cleared {cleared} text runs')


if __name__ == '__main__':
    main()
